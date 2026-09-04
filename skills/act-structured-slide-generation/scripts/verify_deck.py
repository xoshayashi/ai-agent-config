#!/usr/bin/env python3
"""Audit a BUILT .pptx: fonts, colors, overflow (real font metrics), chrome presence.

Usage: verify_deck.py <deck.pptx>
Exit 0 = all checks green, exit 1 = violations found (fix deck.json / builder, rebuild).
"""
from __future__ import annotations

import json
import math
import re
import sys
from pathlib import Path

from deck_text import (MEASURE_OK, _natural_lines, _words as words, drawn_line_h,
                       text_width_in as _measure, _safe_breaks, _unbreakable_spans, _inside_span)
from pptx import Presentation
from pptx.util import Emu

TOKENS = json.loads((Path(__file__).resolve().parent.parent / "references" / "tokens.json").read_text())
LINE_BREAK = TOKENS["line_break"]
ALLOWED = {v.upper() for v in TOKENS["colors"].values()}
FORBIDDEN = {v.upper() for v in TOKENS["color_policy"]["forbidden_colors"]}
OK_FONTS = {TOKENS["fonts"]["latin"], TOKENS["fonts"]["latin_semibold"],
            TOKENS["fonts"]["ea"], TOKENS["fonts"]["ea_semibold"], "+mn-lt", "+mj-lt"}
SLIDE_W_IN = TOKENS["slide"]["width_in"]
SLIDE_H_IN = TOKENS["slide"]["height_in"]
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
CELL_INSET_IN = 0.18   # 表セルの左右余白+折返しの安全余白(build_deck._cell_text_w と同じ)

# 計測は deck_text の単一実装を使う — ビルダーが「どこで切るか」を決めた物差しと、検証が
# 「はみ出すか/重なるか」を判定する物差しは、同じでなければ意味がない
_measure_ok = MEASURE_OK
text_width_in = _measure


def _para_weight(para) -> int:
    """600 (SemiBold) is expressed as a family name, not a bold flag — measure with its own metrics."""
    if any(r.font.bold for r in para.runs):
        return 700
    if any("SemiBold" in (r.font.name or "") for r in para.runs):
        return 600
    return 400


def _run_weight(run) -> int:
    if run.font.bold:
        return 700
    if "SemiBold" in (run.font.name or ""):
        return 600
    return 400


def _text_indent_in(para) -> float:
    """箇条書き段落の字下げ(marL)。本文が使える幅は箱の幅から これを引いた分しかない —
    引き忘れると1行に入る量を多く見積もり、実際は折り返す行を「折り返さない」と数える。"""
    pPr = para._p.find(f"{A}pPr")
    if pPr is None or pPr.find(f"{A}buChar") is None:
        return 0.0
    try:
        return max(0.0, int(pPr.get("marL", "0")) / 914400.0)
    except (TypeError, ValueError):
        return 0.0


def _display_lines(para):
    """段落を「実際に描かれる行」へ割る。<a:br/> はソフト改行なので、幅も行数も
    その手前で切れる — 段落全体を1行として測ると幅を過大に見積もり、折返し数を誤る。"""
    lines, cur = [], []
    for el in para._p:
        if el.tag == f"{A}r":
            cur.append(el)
        elif el.tag == f"{A}br":
            lines.append(cur)
            cur = []
    lines.append(cur)
    return lines


def _para_metrics(tf, w_in):
    """Yield (line_stack_h_in, space_after_in, widest_line_w_in, n_lines) per non-empty
    paragraph. overflow 判定と ink-box 判定が同じ行高モデルを共有するための単一実装。"""
    from pptx.text.text import _Run
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text.strip():
            continue
        size = max((r.font.size.pt if r.font.size else 11) for r in para.runs)
        text_w = max(0.05, w_in - _text_indent_in(para))
        widths = []
        for line in _display_lines(para):
            runs = [_Run(r_el, para) for r_el in line]
            widths.append(sum(
                text_width_in(r.text, r.font.size.pt if r.font.size else size, _run_weight(r))
                for r in runs if r.text
            ))
        lines = sum(max(1, -(-int(w * 100) // max(1, int(text_w * 100)))) for w in widths)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else None
        space_after = para.space_after.pt / 72.0 if para.space_after is not None else 0.0
        # 行の高さは「実際に描かれる高さ」で測る(レンダラはフォント本来の行高より低い行を作らない)
        yield lines * drawn_line_h(size, None, spacing), space_after, max(widths), lines


def check_natural_wrap(shape, warns, where, width_in: float | None = None):
    """その列に収まらない語を拾う。

    行の切れ目は、短いラベルなら文節へ寄せ、それ以外は自然に詰めながら語をまたぐときだけ
    その語を次行へ送る — どちらの経路でも語は割れない。割れるのは「1語が列幅より広い」ときで、
    それは組版ではなくコピーの問題(語を短くするか、列を広げる)。ここで見えるようにする。"""
    tf = shape.text_frame
    # 折返しの判定は、ビルダーが行を決めたときと同じ幅で行う — ここで数 mm でも差をつけると、
    # ぎりぎり1行に収まった文を「2行になる」と読み違え、直すところのない警告が出る
    w_in = width_in if width_in is not None else Emu(shape.width).inches
    if w_in <= 0.05:
        return
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text.strip():
            continue
        size = max((r.font.size.pt if r.font.size else 11) for r in para.runs)
        avail = max(0.05, w_in - _text_indent_in(para))
        weight = _para_weight(para)
        if len({r.font.size.pt if r.font.size else 11 for r in para.runs}) > 1:
            # 値と単位のように大きさの違う走りが並ぶ行。1つの大きさで行を組み直しても
            # 実際の描かれ方にならない — この行は「収まるか」だけを、走りごとに測って見る
            drawn_w = sum(text_width_in(r.text, r.font.size.pt if r.font.size else size,
                                        _run_weight(r)) for r in para.runs)
            if drawn_w > avail:
                warns.append(f"{where}: 値と単位が列に収まらない — '{text[:20]}'"
                             f"(数字の大きさを下げるか、列を広げる。単位が割れて描かれる)")
            continue
        # 「語」は、割れると読みづらい範囲(deck_text._unbreakable_spans): カタカナ語、英単語、
        # 数量と単位。漢字・ひらがなはどこでも折り返せるので語とは数えない(利用者の指示、2026-09-04)
        _spans = _unbreakable_spans(text)
        segs = [text[a:b] for a, b in _spans]
        # 列幅より広い語の判定には、4字以上の漢字の連なり(電子帳簿保存法対応)も数える — 折返しの
        # 対象にはしないが、列に入らず途中で割れて描かれるならコピーか列幅の問題として知らせる
        segs += [m.group() for m in re.finditer(r"[\u4e00-\u9fff]{4,}", text)]
        too_wide = [w for w in segs if len(w) > 1 and text_width_in(w, size, weight) > avail]
        if too_wide:
            warns.append(f"{where}: 列幅に収まらない語 — '{too_wide[0]}'"
                         f"(語を短くするか列を広げる。語の途中で割れて描かれる)")
            continue
        # 語を割らずに組むと行が増える = その列にはコピーが密すぎる。行が短く階段状に並ぶ
        drawn = len(para._p.findall(f"{A}br")) + 1
        cap = max(0.05, avail - 0.3 * size / 72.0)
        natural = math.ceil(text_width_in(text, size, weight) / cap - 1e-9)
        if drawn > natural >= 2:      # 1行→2行は、契約で行数を決めた見出し(表紙の副題)もある
            warns.append(f"{where}: 語を割らずに組むと行が増える — '{text[:20]}'"
                         f"({natural}行ぶんのコピーが{drawn}行になる。短く言い切る)")
            continue
        # 折返しを任せた段落。レンダラの行が語をまたぐなら、その列にはこの文が長すぎる
        if drawn == 1 and text_width_in(text, size, weight) > cap:
            # ビルダーが手を入れなかった段落。レンダラの自然折返し(箱幅、0.3em 狭い幅)の切れ目が
            # 安全な境界の外に落ちるなら、その語が途切れて描かれる
            for width in (avail, cap):
                pos, start = [], 0
                for ln in _natural_lines(text, width, size, weight):
                    start += len(ln)
                    pos.append(start)
                bad = [q for q in pos[:-1] if _inside_span(q, _spans)]
                if bad:
                    q = bad[0]
                    warns.append(f"{where}: 語を割らずには組めない文 — '{text[max(0, q - 4):q + 4]}' が行をまたぐ"
                                 f"(コピーを短くするか列を広げる)")
                    break


def check_table_wrap(shape, warns, where):
    """表のセルも列である。セル内で使える幅は、列幅から左右の内側余白を引いた分しかない —
    ビルダーが折返しを決めた幅と同じ幅で見る(scripts/build_deck._cell_text_w)。"""
    table = shape.table
    widths = [Emu(c.width).inches for c in table.columns]
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            if not cell.text.strip():
                continue
            avail = max(0.4, widths[ci] - CELL_INSET_IN)
            check_natural_wrap(cell, warns, f"{where}: 表 r{ri + 1}c{ci + 1}", width_in=avail)


def check_overflow(shape, issues, where):
    w_in = Emu(shape.width).inches - 0.02
    h_in = Emu(shape.height).inches
    if w_in <= 0.05:
        return
    used_h = sum(h + sa for h, sa, _, _ in _para_metrics(shape.text_frame, w_in))
    # textboxes are allowed to visually overrun their nominal box a bit (top-anchored,
    # autosize off) as long as they don't collide; flag only meaningful overruns
    if used_h > h_in * 1.35 + 0.22:
        issues.append(f"{where}: text likely overflows box ({used_h:.2f}in used vs {h_in:.2f}in high)")


def _bbox(shape):
    l = Emu(shape.left).inches
    t = Emu(shape.top).inches
    return (l, t, l + Emu(shape.width).inches, t + Emu(shape.height).inches)


def _overlap_frac(a, b) -> float:
    """Intersection area as a fraction of the smaller box."""
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    if w <= 0 or h <= 0:
        return 0.0
    inter = w * h
    area = min((a[2] - a[0]) * (a[3] - a[1]), (b[2] - b[0]) * (b[3] - b[1]))
    return inter / area if area > 0 else 0.0


def _ink_bbox(shape):
    """描かれる文字の範囲(インク箱)。段落を積んだ箱では、段落後スペースも文字を押し下げる —
    行の高さだけで測ると、下の段落のぶんが見えず、重なりを見逃す。"""
    box = _bbox(shape)
    if not _measure_ok:
        return box
    w_in = box[2] - box[0]
    max_w, used_h = 0.0, 0.0
    for line_h, space_after, width, _lines in _para_metrics(shape.text_frame, w_in):
        used_h += line_h + space_after
        max_w = max(max_w, min(width, w_in))
    if max_w <= 0:
        return box
    from pptx.enum.text import MSO_ANCHOR as _MA
    y0 = ((box[1] + box[3]) / 2 - used_h / 2
          if shape.text_frame.vertical_anchor == _MA.MIDDLE else box[1])
    x0 = box[0]
    try:
        from pptx.enum.text import PP_ALIGN as _PA
        al = shape.text_frame.paragraphs[0].alignment
        if al == _PA.CENTER:
            x0 = box[0] + (w_in - max_w) / 2
        elif al == _PA.RIGHT:
            x0 = box[2] - max_w
    except Exception:
        pass
    return (x0, y0, x0 + max_w, y0 + used_h)


def _solid_shapes(slide):
    """塗りのある自動図形(バー・カード・矢羽・帯)の枠。線(コネクタ)・文字だけの箱・塗りなしの
    枠線図形は除く — 文字が「縁をまたぐ」ことが問題になるのは、面のある図形だけ。"""
    from pptx.enum.dml import MSO_FILL
    out = []
    for sh in slide.shapes:
        if sh.shape_type != 1 or sh.left is None or sh.width is None:      # 1 = AUTO_SHAPE
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        try:
            if sh.fill.type != MSO_FILL.SOLID:
                continue
            rgb = str(sh.fill.fore_color.rgb)
        except Exception:
            continue
        out.append((_bbox(sh), sh, rgb))
    # 同色の図形に完全に含まれる図形(角を四角くする詰め物など)は、見た目には1つの面 —
    # その内側の縁は存在しないので、またぎ判定の対象から外す
    def _inside(a, b):
        return a[0] >= b[0] - 0.005 and a[1] >= b[1] - 0.005 and a[2] <= b[2] + 0.005 and a[3] <= b[3] + 0.005
    keep = []
    for i, (bb, sh, rgb) in enumerate(out):
        if any(j != i and rgb == rgb2 and _inside(bb, bb2) and bb != bb2
               for j, (bb2, _s2, rgb2) in enumerate(out)):
            continue
        keep.append((bb, sh))
    return keep


def check_straddle(slide, idx, issues) -> None:
    """文字が塗り図形の「縁をまたぐ」= バーの上に載った値ラベル、カードからはみ出した本文。
    文字が図形の中に収まる(カードの中の文)のも、外に離れている(バーの上のラベル)のも正しい。
    中途半端に重なる状態だけが欠陥 — 面積の 12〜88% が図形にかかっていれば、字面は縁の上にある。
    機械ゲートが text↔text と text↔chart しか見ていなかったため、ウォーターフォールの値ラベルが
    バーに食い込んだまま「0 failures」で通っていた(監査 2026-09-03)。"""
    solids = _solid_shapes(slide)
    if not solids:
        return
    for shape in slide.shapes:
        if shape.shape_type != 17 or not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        if shape.left is None:
            continue
        tb = _ink_bbox(shape)
        t_area = max(1e-6, (tb[2] - tb[0]) * (tb[3] - tb[1]))
        for sb, _sh in solids:
            ow = min(tb[2], sb[2]) - max(tb[0], sb[0])
            oh = min(tb[3], sb[3]) - max(tb[1], sb[1])
            if ow <= 0.08 or oh <= 0.04:
                continue
            frac = ow * oh / t_area
            if 0.12 < frac < 0.88:
                issues.append(f"slide {idx}: テキストが図形の縁をまたぐ({frac:.0%} が面にかかる) — "
                              f"'{shape.text_frame.text[:18]}'(ラベルを図形の外へ出すか、中へ収める)")
                break


CARD_FILL_FLOOR = TOKENS["layout"].get("fill", {}).get("card_text_floor", 0.45)


def check_card_fill(slide, idx, warns) -> None:
    """カード(surface_tint の塗り)の中身が薄くないか。占有契約(fit_band)はカードの高さを本文帯に
    合わせて育てるので、コピーが短いとカードの下半分が空く — 見出し・帯・他のカードと釣り合わない
    (2026-09-04、利用者の指摘)。カードの内側に収まる文字のインク範囲(上端〜下端)が、カード高の
    card_text_floor(既定 45%)に届かなければ警告する。1.2in 未満の低いカード(帯・結論ストリップ)は対象外。
    直し方はコピーを足すこと。型を大きくして埋めるのは、ここでは選ばない。"""
    # 対象は「縦に読むカード」: 高さ 1.2in 以上で、幅が高さの3倍未満のもの。要約ページの横長の
    # 行(幅12in × 高さ1.6in)は中身を上下中央に置く帯なので、下半分が空く問題は起きない
    # 対象はビルダーが "card" と名付けた図形(process_flow / column_framework / two_column の本文カード)。
    # 章扉の側面パネルやロードマップのセルは同じ塗りでもカードではない
    cards = [(bb, sh) for bb, sh in _solid_shapes(slide)
             if getattr(sh, "name", "") == "card" and (bb[3] - bb[1]) >= 1.2]
    if not cards:
        return
    texts = [_ink_bbox(sh) for sh in slide.shapes
             if sh.shape_type == 17 and sh.has_text_frame and sh.text_frame.text.strip()
             and sh.left is not None]
    for cb, _sh in cards:
        # 面積の半分以上がカードに載る文字を「中身」と数える。あふれた本文も数えないと、
        # 溢れているカードを「薄い」と誤判定する。範囲はカードの内側で切る
        inside = [tb for tb in texts if _overlap_frac(tb, cb) >= 0.5
                  and (tb[2] - tb[0]) * (tb[3] - tb[1]) <= (cb[2] - cb[0]) * (cb[3] - cb[1])]
        if not inside:
            continue
        span = min(max(tb[3] for tb in inside), cb[3]) - max(min(tb[1] for tb in inside), cb[1])
        card_h = cb[3] - cb[1]
        ratio = span / card_h if card_h > 0 else 1.0
        if ratio < CARD_FILL_FLOOR:
            warns.append(f"slide {idx}: カードの中身が薄い(文字がカード高の {ratio:.0%}、下限 {CARD_FILL_FLOOR:.0%}) — "
                         "項目や説明を足して、見出し・帯と釣り合う量にする(型は大きくしない)")


def check_frame_overlaps(slide, idx, warns) -> None:
    """テキストボックスの「枠」どうしの重なり。描かれる文字が正しくても、枠が重なった pptx は
    編集で掴み違える(下の枠を選べない)。枠は文字を囲むだけの大きさに保つこと。"""
    boxes = [(_bbox(sh), sh.text_frame.text[:16]) for sh in slide.shapes
             if sh.shape_type == 17 and sh.has_text_frame and sh.text_frame.text.strip()
             and sh.left is not None]
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i][0], boxes[j][0]
            ov_w = min(a[2], b[2]) - max(a[0], b[0])
            ov_h = min(a[3], b[3]) - max(a[1], b[1])
            if ov_w > 0.01 and ov_h > 0.01:
                warns.append(f"slide {idx}: テキストボックスの枠が重なる {ov_h * 72:.1f}pt — "
                             f"'{boxes[i][1]}' × '{boxes[j][1]}'(枠は文字ぶんの大きさに)")


def check_collisions(slide, idx, issues) -> None:
    """Rendered-ink collisions: text↔text and large-text↔chart/table overlaps.
    Small overlays on charts (YoY badges etc.) are intentional and skipped."""
    texts, frames = [], []
    for shape in slide.shapes:
        if shape.left is None or shape.width is None:
            continue
        if shape.shape_type == 17 and shape.has_text_frame and shape.text_frame.text.strip():
            texts.append((_ink_bbox(shape), shape.text_frame.text[:18]))
        elif getattr(shape, "has_chart", False) or shape.shape_type == 19:
            frames.append((_bbox(shape), "chart/table"))
    for i in range(len(texts)):
        for j in range(i + 1, len(texts)):
            frac = _overlap_frac(texts[i][0], texts[j][0])
            if frac > 0.25:
                issues.append(
                    f"slide {idx}: テキストの重なり {frac:.0%} — '{texts[i][1]}' × '{texts[j][1]}'")
    for tb, label in texts:
        t_area = max(0.0, (tb[2] - tb[0]) * (tb[3] - tb[1]))
        for fb, _ in frames:
            f_area = (fb[2] - fb[0]) * (fb[3] - fb[1])
            if f_area <= 0 or t_area / f_area < 0.10:
                continue  # small annotation overlay on a chart is by design
            if _overlap_frac(tb, fb) > 0.3:
                issues.append(
                    f"slide {idx}: テキストが図表に重なる — '{label}'")


def audit(path: Path) -> tuple[list[str], list[str], int]:
    """Run every check on a built .pptx and return (failures, warnings, slide_count).
    The CLI prints them; stress_deck / tests call this directly."""
    prs = Presentation(path)
    issues: list[str] = []
    warns: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        xml = slide._element.xml
        # color audit: every literal srgbClr must be in the Act palette
        for m in set(re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml)):
            mu = m.upper()
            if mu in FORBIDDEN:
                issues.append(f"slide {idx}: forbidden color #{mu}")
            elif mu not in ALLOWED:
                issues.append(f"slide {idx}: off-palette color #{mu}")
        # font audit
        for m in set(re.findall(r'typeface="([^"]+)"', xml)):
            if m not in OK_FONTS:
                issues.append(f"slide {idx}: off-system font '{m}'")
        # East-Asian coverage: any run holding JP text must have <a:ea>
        for r_el in slide._element.iter(f"{A}r"):
            t = r_el.find(f"{A}t")
            if t is None or not t.text or not any(ord(c) > 0x2E7F for c in t.text):
                continue
            rPr = r_el.find(f"{A}rPr")
            if rPr is None or rPr.find(f"{A}ea") is None:
                issues.append(f"slide {idx}: Japanese run without <a:ea> font: '{t.text[:20]}'")
        # placeholder rot
        low = xml.lower()
        for bad in ("lorem", "ipsum", "xxxx", "placeholder"):
            if bad in low:
                issues.append(f"slide {idx}: placeholder text '{bad}' present")
        # autofit shrink is the canonical header-consistency breaker: nominally
        # identical titles silently render at different sizes per slide
        if "normAutofit" in xml and "fontScale" in xml:
            issues.append(f"slide {idx}: autofit font shrink (<a:normAutofit fontScale>) — サイズ縮小ではなく短文化/分割で収める")
        # geometry + chrome
        n_text = 0
        for shape in slide.shapes:
            if shape.left is None:
                continue
            l, t = Emu(shape.left).inches, Emu(shape.top).inches
            r = l + Emu(shape.width).inches
            b = t + Emu(shape.height).inches
            if l < -0.01 or t < -0.01 or r > SLIDE_W_IN + 0.01 or b > SLIDE_H_IN + 0.01:
                issues.append(f"slide {idx}: shape '{shape.shape_type}' out of slide bounds ({l:.2f},{t:.2f})-({r:.2f},{b:.2f})")
            if shape.has_text_frame and shape.text_frame.text.strip():
                n_text += 1
                if _measure_ok:
                    check_overflow(shape, issues, f"slide {idx}")
                    check_natural_wrap(shape, warns, f"slide {idx}")
            elif getattr(shape, "has_table", False):
                # 表もテキストである。セルを見ないと、列に収まらない語がそのまま割れて描かれる
                n_text += 1
                if _measure_ok:
                    check_table_wrap(shape, warns, f"slide {idx}")
        if n_text == 0:
            issues.append(f"slide {idx}: no text at all")
        check_collisions(slide, idx, issues)
        if _measure_ok:
            check_straddle(slide, idx, issues)
            check_card_fill(slide, idx, warns)
        check_frame_overlaps(slide, idx, warns)

    if not _measure_ok:
        warns.append("NotoSansJP-{400,600,700}.ttf not found — overflow measurement skipped (install fonts)")
    return issues, warns, len(list(prs.slides))


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__)
        return 1
    issues, warns, n = audit(Path(sys.argv[1]))
    for w in warns:
        print(f"WARN: {w}")
    for i in issues:
        print(f"FAIL: {i}")
    print(f"\n{n} slides / {len(issues)} failures / {len(warns)} warnings")
    return 1 if issues else 0


if __name__ == "__main__":
    sys.exit(main())
