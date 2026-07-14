#!/usr/bin/env python3
"""Build an Act-styled 16:9 banker-grade .pptx from a deck.json spec.

Usage: build_deck.py <deck.json> [-o out.pptx]

The LLM owns content (deck.json); this script owns geometry, typography, and color. The
engine keeps the Act system consistent, but pattern renderers must branch by evidence shape
and available data rather than forcing one fixed layout.
"""
from __future__ import annotations

import argparse
import json
import math
import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

TOKENS = json.loads((Path(__file__).resolve().parent.parent / "references" / "tokens.json").read_text())

C = {k: RGBColor.from_string(v) for k, v in TOKENS["colors"].items()}
TS = TOKENS["type_scale_pt"]
LAY = TOKENS["layout"]
LINE_BREAK = TOKENS["line_break"]
CH = TOKENS["chart_style"]
FONTS = TOKENS["fonts"]
VALUE_DELTA_GAP_IN = 0.18
SLIDE_W = Inches(TOKENS["slide"]["width_in"])
SLIDE_H = Inches(TOKENS["slide"]["height_in"])
MX = LAY["margin_x_in"]
CONTENT_W = TOKENS["slide"]["width_in"] - 2 * MX
FOOT_Y = LAY["footer"]["y_in"]
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ---------------------------------------------------------------- utilities

# 表示長(_ja_len)と全角→半角(hw)は deck_text.py の単一実装を使う — validate_spec の
# 字数バジェット・lint_render の readback 照合と同一実装であることが契約
from deck_text import (MEASURE_OK, drawn_line_h, footer_text, header_slots, hw, ink_slacks,
                       is_prose, ja_len as _ja_len, leading, text_width_in, wrap_display,
                       wrap_prose)


def _label_w(text: str, size_pt: float, weight: int = 400) -> float:
    """1行ラベルの枠幅(in)。枠は文字を囲むだけの大きさに保つ — 余った幅を含めると隣の枠へ
    食い込み、切り詰めると折り返して2行になる。実測幅ならどちらも起きない。

    余裕は折返しの安全マージン(wrap_display が箱幅から引く 0.3em)より広く取る。
    ぴったりの幅で渡すと「収まらない」と判断され、1行のはずのラベルが割れる。"""
    return text_width_in(text, size_pt, weight) + size_pt / 72.0 * 0.4 + 0.03


def _text_lines(text: str, width_in: float, size_pt: float, weight: int = 400) -> int:
    """描かれる行数。箱の高さを決める側と、改行を打つ側は同じ答えを見なければならない —
    字数近似で数えると、文節で1行増えたぶんがカードの下側余白から黙って差し引かれる。

    幅が 0 以下で呼ばれるのは呼び出し側の算術の誤り。黙って丸めると「142行」のような
    答えを返し、それが高さへ流れ込んでスライドの外へ出る図形を生む(実際に起きた)。"""
    if width_in <= 0:
        raise ValueError(f"_text_lines: 描画幅が正でない ({width_in:.3f}in) — 呼び出し側の算術を疑う")
    if not text:
        return 0
    total = 0
    for line in display_wrap_text(str(text), width_in, size_pt, weight).split("\n"):
        w = text_width_in(line, size_pt, weight)
        total += max(1, math.ceil(w / max(0.05, width_in) - 1e-9))
    return max(1, total)


def _statement_lines(stmt: str, width_in: float, size_pt: float) -> list[str]:
    """Break a centered closing statement into display lines on 、 clause boundaries so
    the tail never renders as an orphaned one/two-character line. Clauses are greedily
    packed to the box width; short statements (no 、) keep their natural single-run wrap."""
    if "\n" in stmt:
        return stmt.split("\n")
    if "、" not in stmt:
        return [stmt]
    raw = stmt.split("、")
    clauses = [c + "、" for c in raw[:-1]] + ([raw[-1]] if raw[-1] else [])
    if len(clauses) < 2:
        return [stmt]
    # 床値 8.0 は意図的に _text_lines の 4.0 と異なる: ステートメント箱は常に広く、
    # 節パッキングの下限を狭くすると1節1行の退化分割になるため(極端入力のみ差が出る)
    cap = max(8.0, width_in / (size_pt / 72.0))
    lines: list[str] = []
    cur = ""
    for c in clauses:
        if cur and _ja_len(cur) + _ja_len(c) > cap:
            lines.append(cur)
            cur = c
        else:
            cur += c
    if cur:
        lines.append(cur)
    return lines


# 象徴的なメッセージ(結び・ステートメント)は、文字を流し込むのではなく「形」を作る。
# 行長(measure)を選び、節の切れ目で割り、行の長さを揃え、周りの余白との占有バランスを見る。
# 端から端まで文字で埋めた行は、読み手に「詰め込まれた」と感じさせる。
def _units_of(text: str, avail_w: float, size_pt: float, weight: int) -> list[str]:
    """長い節を、行の切れ目になれる単位へ割る。まず文節、収まらないものだけ語へ。
    文節で割れば、行頭に助詞が立たない(「…存在感 / をプログラム…」にならない)。"""
    from deck_text import _segments, _words
    out: list[str] = []
    for chunk in _segments(text)[0]:
        if text_width_in(chunk, size_pt, weight) > avail_w * 0.9:
            out.extend(_words(chunk))
        else:
            out.append(chunk)
    return out


MESSAGE_SUPPORT_MEASURE = 0.66     # 支え文の行長(使える幅に対する比)。右に余白を残す
DASHES = "—―─－"


def _clauses(stmt: str) -> list[str]:
    """節に割る。読点とダッシュは節の切れ目 — 意味の変わり目で行が変わると、文体が形に出る。
    ダッシュは連なって1つの記号(——)なので、途中では割らない。"""
    out, cur = [], ""
    i = 0
    while i < len(stmt):
        ch = stmt[i]
        if ch in DASHES:
            j = i
            while j < len(stmt) and stmt[j] in DASHES:
                j += 1
            out.append(cur + stmt[i:j])
            cur = ""
            i = j
            continue
        cur += ch
        if ch in "、，":
            out.append(cur)
            cur = ""
        i += 1
    if cur:
        out.append(cur)
    return [c for c in out if c]


def split_message(spec: dict) -> tuple[str, str | None]:
    """メッセージを「言い切る一行(lead)」と「それを支える一文(support)」に分ける。

    象徴的なスライドは、まず一行で言い切り、その根拠や含意を小さな一文で支える形が読みやすい。
    ダッシュでつないだ長い一文(A——B)は、B が言い切りで A が支えになっていることが多いので、
    その形に開く。lead を明示したいときは spec.lead に書く。"""
    stmt = spec.get("statement", "")
    lead = spec.get("lead")
    if lead:
        return lead, stmt or None
    parts = [p.strip() for p in re.split(f"[{DASHES}]+", stmt) if p.strip()]
    if len(parts) == 2:
        return parts[1], parts[0]                 # ダッシュの後ろが言い切り、前が支え
    return stmt, None


def shape_message(stmt: str, avail_w: float, size_pt: float, weight: int = 700,
                  max_lines: int = 4) -> tuple[list[str], float]:
    """メッセージの行(lines)と、その行を組む幅(measure)を返す。

    象徴的な一文は、流し込みではなく形として組む:
      - 行は節(読点・ダッシュ)の切れ目で変わる。意味の変わり目が行の変わり目になる
      - 行の長さを揃える(行長の分散を最小化する分割を選ぶ)
      - 行長は使える幅より短く取り、右に余白を残す。端まで埋めた行は詰め込んで見える
      - 語は割らない(節が長すぎる行は語の単位で詰める)
    """
    clauses = _clauses(stmt)
    units: list[tuple[str, bool]] = []           # (文字列, 節の終わりか)
    for c in clauses:
        if text_width_in(c, size_pt, weight) > avail_w * 0.9:
            parts = _units_of(c, avail_w, size_pt, weight)
            for i, part in enumerate(parts):
                units.append((part, i == len(parts) - 1))
        else:
            units.append((c, True))
    widths = [text_width_in(u, size_pt, weight) for u, _ in units]
    total = sum(widths)
    n = len(units)
    if n < 2:
        return [stmt], min(avail_w, total + 0.06)

    best, best_score = None, float("inf")
    for k in range(1, max_lines + 1):
        if k > n:
            break
        target = total / k
        INF = float("inf")
        dp = [[INF] * (k + 1) for _ in range(n + 1)]
        cut = [[0] * (k + 1) for _ in range(n + 1)]
        dp[0][0] = 0.0
        for i in range(1, n + 1):
            for j in range(1, k + 1):
                w = 0.0
                for start_i in range(i - 1, -1, -1):
                    w += widths[start_i]
                    if w > avail_w * 0.94 and i - start_i > 1:
                        break
                    if dp[start_i][j - 1] == INF:
                        continue
                    # 行長のばらつき + 節の途中で行が変わることの罰
                    cost = (target - w) ** 2
                    if i < n and not units[i - 1][1]:
                        cost += 0.35
                    if dp[start_i][j - 1] + cost < dp[i][j]:
                        dp[i][j] = dp[start_i][j - 1] + cost
                        cut[i][j] = start_i
        if dp[n][k] == INF:
            continue
        lines, i, j = [], n, k
        while j > 0:
            st = cut[i][j]
            lines.append("".join(u for u, _ in units[st:i]))
            i, j = st, j - 1
        lines.reverse()
        line_w = [text_width_in(ln, size_pt, weight) for ln in lines]
        measure = max(line_w)
        if measure > avail_w:
            continue
        even = sum((measure - lw) ** 2 for lw in line_w)
        fill = measure / avail_w
        score = even + max(0.0, fill - 0.88) * 12 + max(0, k - 3) * 0.6
        if score < best_score:
            best, best_score = (lines, measure), score
    if not best:
        # 収まる形が見つからない = コピーが長すぎる。自然に流し、行数は実測で数える
        # (1行と誤って数えると、下の要素がその上に載る)
        lines = wrap_prose(stmt, avail_w, size_pt, weight).split("\n")
        drawn: list[str] = []
        for ln in lines:
            n = max(1, math.ceil(text_width_in(ln, size_pt, weight) / avail_w - 1e-9))
            drawn.extend([ln] if n == 1 else [ln] * n)   # 折返しぶんも行として数える
        return drawn, avail_w
    lines, measure = best
    # 箱は行の実測幅より少し広く取る(折返しの安全マージン 0.3em ぶん)。ぴったりの幅で渡すと
    # レンダラが先に折り返し、行末のダッシュや約物だけが次の行へこぼれる
    return lines, min(avail_w, measure + 0.3 * size_pt / 72.0 + 0.06)


def grid(col_start: int, col_span: int) -> tuple[float, float]:
    """Return (x_in, w_in) for a 12-column grid region (0-indexed start)."""
    ncols = LAY["grid_columns"]
    gut = LAY["gutter_in"]
    colw = (CONTENT_W - gut * (ncols - 1)) / ncols
    x = MX + col_start * (colw + gut)
    w = col_span * colw + (col_span - 1) * gut
    return x, w


def _set_run_fonts(run, size_pt: float, weight: int = 400, color: RGBColor | None = None,
                   ea_latin: bool = False, latin_ea: bool = False,
                   lang: str | None = None) -> None:
    """Apply Act typography to a run: Geist + Noto Sans JP, weight via family/bold."""
    f = run.font
    f.size = Pt(size_pt)
    f.color.rgb = color if color is not None else C["ink"]
    if weight >= 700:
        f.name, ea = FONTS["latin"], FONTS["ea"]
        f.bold = True
    elif weight >= 600:
        f.name, ea = FONTS["latin_semibold"], FONTS["ea_semibold"]
        f.bold = False
    else:
        f.name, ea = FONTS["latin"], FONTS["ea"]
        f.bold = False
    if latin_ea:
        f.name = ea  # 数字ラン: Latin 側も EA フォント(半角数字 0.55em)で組む(下記参照)
    rPr = run._r.get_or_add_rPr()
    if lang:
        # ラン言語を明示する。未指定だと PowerPoint 系ビューアが編集ロケール(ja-JP)で
        # スクリプト整形し、ASCII の数字・欧文を EA フォント側で描画することがある
        # (半角コードポイントでも全角風の見た目になる)。ASCII ランは en-US、
        # CJK ランは ja-JP に固定する。
        rPr.set("lang", lang)
    if ea_latin:
        ea = f.name  # ASCII 専用ラン: EA 側も Latin フォントに固定(下記参照)
    for tag in ("ea", "cs"):
        el = rPr.find(f"{A_NS}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{A_NS}{tag}")
        el.set("typeface", ea)


_ASCII_SEG = re.compile(r"[\x20-\x7E]+")
_DIGIT_SEG = re.compile(r"^[0-9 /:.,()%+\-]*[0-9][0-9 /:.,()%+\-]*$")


def _add_script_runs(p, text, size_pt, weight, color):
    """1段落ぶんのテキストを追加する。改行(\n)は段落を割らず、段落内のソフト改行
    (DrawingML の <a:br/>)として打つ。

    段落を割ってしまうと、(1) 箇条書きでは折返し行の頭にもう1つ記号が付き、
    (2) 段落後スペース(spcAft)が「同じ一続きの文」の行間に入り込む。ソフト改行なら
    行だけが増え、ぶら下げインデント・記号・行間はその段落のものが保たれる。"""
    for i, line in enumerate(str(text).split("\n")):
        if i:
            br = etree.SubElement(p._p, f"{A_NS}br")
            rPr = etree.SubElement(br, f"{A_NS}rPr")
            rPr.set("sz", str(int(size_pt * 100)))   # 改行の行高はこの rPr で決まる
            rPr.set("lang", "ja-JP")
        _add_line_runs(p, line, size_pt, weight, color)


def _add_line_runs(p, text, size_pt, weight, color):
    """1行ぶんを ASCII 連続区間と CJK 区間の別ランに分割して追加する。
    CJK 文中の数字・欧文は、ビューアのスクリプト判定(LibreOffice の Asian/Western
    itemization 等)次第で Asian フォント側に割り当てられ、半角コードポイントでも
    全角風の幅で描画されることがある。ASCII 区間を独立ランにし、そのランの
    <a:ea> も Latin フォントへ固定すれば、どのビューアでも数字は Latin フォントで
    半角の見た目に確定する(カバーの日付に限らず全描画テキスト共通)。"""
    s = hw(text)
    has_cjk = any(ord(ch) > 0x2E7F for ch in s)
    pos = 0
    for m in _ASCII_SEG.finditer(s):
        if m.start() > pos:
            r = p.add_run(); r.text = s[pos:m.start()]
            _set_run_fonts(r, size_pt, weight, color, lang="ja-JP")
        r = p.add_run(); r.text = m.group()
        # CJK 混在文中の数字だけの区間(日付・時刻・件数)は EA フォントの半角数字で組む。
        # Geist の数字は 0.625em と広く(Noto Sans JP の半角数字は 0.55em)、和文中では
        # 全角数字のように見える。欧文単語(文字を含む区間)はブランドフォントの Geist を維持し、
        # CJK を含まない純欧文文字列(KPI 値等)には適用しない。
        if has_cjk and _DIGIT_SEG.match(m.group()):
            _set_run_fonts(r, size_pt, weight, color, latin_ea=True, lang="en-US")
        else:
            _set_run_fonts(r, size_pt, weight, color, ea_latin=True, lang="en-US")
        pos = m.end()
    if pos < len(s) or not s:
        r = p.add_run(); r.text = s[pos:]
        _set_run_fonts(r, size_pt, weight, color, lang="ja-JP")


DISPLAY_WRAP_MAX_CHARS = LINE_BREAK["max_display_chars_ja"]   # これを超えたら本文 — 自然に行を埋める
DISPLAY_WRAP_MAX_LINES = LINE_BREAK["max_lines"]              # 文節で割るときに許す行数の上限
LABEL_MAX_CHARS = LINE_BREAK["label_max_chars_ja"]            # ここまでが「ラベル」 — 文節で割る


def display_wrap_text(text: str, w_in: float, size_pt: float, weight: int = 400) -> str:
    """表示テキスト(ラベル・見出し・セル・結論句)の改行位置を文節の切れ目へ寄せる。
    手で改行してあるもの、長すぎる本文、そして文章(読点を持つ文)は素通し — 文の折返しは
    レンダラに任せる(wrap_display 側で判定)。

    行数は「自然折返しの行数 + 1」までを許す。語の途中で割るくらいなら1行増やす、という
    優先順位 — 増えた1行が箱に入らないなら、それは組版ではなくコピーが長すぎるという意味で、
    verify_deck の overflow/重なり検査が拾う(直すのはコピーの側)。"""
    if not isinstance(text, str) or "\n" in text:
        return text
    if not is_prose(text) and _ja_len(text) <= LABEL_MAX_CHARS:
        # 短いラベル・見出し・結論句は、意味の切れ目(文節)で割る — 行の切れ目が文体に表れる
        natural = math.ceil(text_width_in(text, size_pt, weight) / max(0.05, w_in))
        broken = wrap_display(text, w_in, size_pt,
                              min(DISPLAY_WRAP_MAX_LINES, natural + 1), weight)
        if "\n" in broken:
            return broken
    # 文章・長い表示テキスト、および文節で割れなかったもの: 自然に行を埋め、語が割れる行だけ
    # その語を次行へ送る(行の埋まりを保ったまま、語は割らない)
    return wrap_prose(text, w_in, size_pt, weight)


def add_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing: float | None = None, space_after_pt: float = 0.0, wrap=True,
             autosize_off=True, display_wrap=True, role: str | None = None):
    """Add a textbox. `runs` is a list of paragraphs; each paragraph is a list of
    (text, size_pt, weight, color) tuples.

    display_wrap: 折返しの起きる短い表示テキスト(ラベル・見出し・結論句)は、改行位置を
    文節の切れ目へ寄せてから組む。レンダラ任せの折返しは箱の幅だけを見て語の途中で割るので、
    「導入費＋固定利用/料」のように意味の切れ目と行の切れ目がずれる。長い本文には手を出さない
    (DISPLAY_WRAP_MAX_CHARS 超、または DISPLAY_WRAP_MAX_LINES に収まらないものは素通り)。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # word_wrap=False の箱でも折返しは起きる(レンダラは箱の幅で割る)。だから折返し位置の
    # 判断はここで行い、wrap フラグに関係なく文節へ寄せる
    if display_wrap:
        runs = [
            [(display_wrap_text(para[0][0], w, para[0][1], para[0][2]), *para[0][1:])]
            if len(para) == 1
            else para
            for para in runs
        ]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # 行間は級数が決める(tokens.leading)。段落ごとに、その段落の級数で引く —
        # 箱ごとに行間を決めると、同じ役割の文章が箱によって違う密度で見える
        size = max((r[1] for r in para), default=TS["body"])
        p.line_spacing = line_spacing if line_spacing else leading(size, role)
        if space_after_pt:
            p.space_after = Pt(space_after_pt)
        for text, size_pt, weight, color in para:
            _add_script_runs(p, text, size_pt, weight, color)
    return tb


# 箇条書きは「本物の箇条書き」で組む — 段落の bullet(DrawingML buChar)+ぶら下げインデント。
# 図形の円を本文の横に置く方式はやめた: 記号が本文と連動せず、編集で行が増減しても点が
# 取り残される(PowerPoint で開くと、点だけ浮いた妙な箇条書きになる)。
#
# ただし記号はベースライン揃えで打たれるので、縦位置は「グリフのインク中心 × buSzPct」で
# 決まる。ここを外すと点が下にずれる:
#   和文字面(全角の四角)の中心 : ベースラインから 0.3805em ← 記号を乗せたい高さ
#   Geist の「●」の中心        : 0.357em → 60% に縮めると 0.214em、本文16ptで約2.7pt 下
#   和文の中黒(・/･)の中心      : 0.380em → 字面中央と一致(等倍のときだけ)
# 縮小率は中心も比例して下げるので、中央に乗る倍率は 100% だけ。したがって
# 「和文フォントの中黒を等倍で」以外の組み合わせにしないこと。LibreOffice は独自に
# 中央寄せして欠陥を隠すが、PowerPoint / Keynote はベースライン揃えのまま描く。
#
# 記号は半角中黒(U+FF65)を使う。全角中黒(U+30FB)は送り幅が 1em あり、点はその中央に
# 打たれるので、点の右に 0.5em の空白がぶら下がる — 本文との間が間延びして、点だけが
# 左へ取り残されて見える。半角中黒は縦中心が同じ 0.380em のまま送り幅が 0.5em なので、
# 縦位置を保ったまま点を本文へ寄せられる。
#   記号の水平位置 = 行頭(marL + indent) + グリフのインク中心(半角中黒 0.249em)
#   本文の左端     = marL   ← 折返し行もここに揃う(ぶら下げインデント)
# ぶら下げは indent = -marL(行頭に記号、折返しは marL)。これを崩すと LibreOffice が
# 折返し行を記号の下へ回してしまう。点と本文の間隔は marL そのもので決める —
# 半角中黒の送り幅(0.5em = 本文16ptで 0.111in)に必要な余白を足した値にする。
BULLET_CHAR = "･"            # U+FF65 半角中黒(縦中心 0.380em / 送り幅 0.5em)
BULLET_SIZE_PCT = "100000"   # 等倍。縮小すると記号の中心が字面中央から下へずれる
BULLET_INDENT_IN = 0.16      # marL: 本文の左端(折返し行もここ)。半角中黒 0.111in + 余白


# カード内の縦積みは「箱」ではなく「インク(字面)」の間隔で組む。
#
# 箱の高さは級数ごとに違い、インクは箱の中で上下対称に座らない(和文の行ボックスは下に
# ディセンダ分の余白を持ち、欧文数字は字面が上寄り)。したがって箱の隙間を等しくしても、
# 見える隙間は等しくならない — 値の上と下で余白が食い違って見えるのはこれが理由である。
# stack_optical はインク矩形を並べ、その隙間を役割ごとの gap で決める。上下のインセットも
# 対称になる(残りを均等に配る)ので、「上は空いているのに下は詰まっている」も同時に消える。
# 較正値(ink_ratio / ink_center_offset_em)は 300dpi の実レンダーから採った(tokens 参照)。
OPT = LAY["optical_stack"]


def _parts_lines(parts, width_in: float) -> int:
    """そのブロックが実際に何行で描かれるか。行数はブロックの持ち物であって、呼び出し側が
    渡す前提にはしない(渡し忘れると1行として高さを積み、下の要素へ食い込む)。"""
    if len(parts) == 1:
        text, size, weight, _c = parts[0]
        return _text_lines(str(text), width_in, size, weight)
    width = sum(text_width_in(str(t), sz, wt) for t, sz, wt, _c in parts)
    return max(1, math.ceil(width / max(0.05, width_in) - 1e-9))


def _ink_h(b, width_in: float) -> float:
    """1ブロックのインク高(in)。折返し行は行送りぶん積み、最後の行だけ字面高で数える。"""
    n = max(1, _parts_lines(b["parts"], width_in))
    line_h = drawn_line_h(b["size"], b.get("role"))
    return (n - 1) * line_h + b["size"] / 72.0 * OPT["ink_ratio"][b.get("kind", "text")]


def _stack_gaps(blocks) -> tuple[list[float], list[float]]:
    """段落の継ぎ目ごとに (見えるインクの隙間, 段落後スペース) を返す。

    行ボックスは、すでに上下に余白を持っている(和文は下のディセンダ、大きな数字は上が広い)。
    見える隙間はその余白どうしの和(=下限)から始まるので、段落後スペースはそこから差し引いて置く。
    同じ役割の継ぎ目(同じ gap_name)は、見える隙間を揃える — 下限の大きい側に合わせるので、
    数値の上だけが広い、といった食い違いが起きない。"""
    floors, names, wants = [], [], []
    for i in range(len(blocks) - 1):
        cur, nxt = blocks[i], blocks[i + 1]
        _, below = ink_slacks(cur["size"], cur.get("kind", "text"),
                              leading(cur["size"], cur.get("role")))
        above, _ = ink_slacks(nxt["size"], nxt.get("kind", "text"),
                              leading(nxt["size"], nxt.get("role")))
        floors.append(below + above)
        names.append(nxt.get("gap_name") or f"#{i}")
        wants.append(nxt.get("gap_before", 0.0))
    want_by_name: dict[str, float] = {}
    for nm, floor, want in zip(names, floors, wants):
        want_by_name[nm] = max(want_by_name.get(nm, 0.0), want, floor)
    gaps = [want_by_name[nm] for nm in names]
    spc = [g - f for g, f in zip(gaps, floors)]
    return gaps, spc


def stack_optical_height(blocks, width_in: float) -> float:
    """縦積みが必要とするインク総高(in)。器の高さを決める側と描く側が同じ式を見る。"""
    if not blocks:
        return 0.0
    gaps, _spc = _stack_gaps(blocks)
    return sum(_ink_h(b, width_in) for b in blocks) + sum(gaps)


def _stack_drawn_h(blocks, w) -> float:
    """塊が実際に占める高さ(実描画の行高 + 段落後スペース)。器の高さと、中央に置くときの
    基準は、この「描かれる高さ」で測る。"""
    if not blocks:
        return 0.0
    _gaps, spc = _stack_gaps(blocks)
    total = sum(max(1, _parts_lines(b["parts"], w)) * drawn_line_h(b["size"], b.get("role"))
                for b in blocks)
    return total + sum(spc) + 0.02


def stack_block(slide, x, y, w, h, blocks, *, align=PP_ALIGN.LEFT):
    """1つのテキストボックスに、blocks を段落として積む。

    分ける必要のないもの(ラベル→数値→注記、見出し→本文)は同じ箱の中で仕上げる。箱を分けると
    位置を人手で合わせることになり、行間・余白・枠の重なりのズレがそこから生まれる。1つの箱なら、
    段落の行間(級数が決める)と段落後スペースだけで組み上がる。

    blocks: [{parts, size, kind, gap_before, gap_name, align, role}]
      gap_before : 直前のブロックとの「インク(字面)の隙間」(in)
      gap_name   : 同じ役割の継ぎ目につける名前。同名の隙間は見た目が揃う
    返り値: (shape, インク総高(in))
    """
    if not blocks:
        return None, 0.0
    gaps, spc = _stack_gaps(blocks)
    ink_total = sum(_ink_h(b, w) for b in blocks) + sum(gaps)
    first = blocks[0]
    above0, _ = ink_slacks(first["size"], first.get("kind", "text"),
                           leading(first["size"], first.get("role")))
    top = y + max(0.0, (h - ink_total) / 2) - above0     # インクが領域の中央に来るよう箱を置く

    tb = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = b.get("align", align)
        p.line_spacing = leading(b["size"], b.get("role"))
        if i < len(spc) and spc[i] > 0:
            p.space_after = Pt(spc[i] * 72)
        parts = b["parts"]
        if len(parts) == 1:
            parts = [(display_wrap_text(parts[0][0], w, parts[0][1], parts[0][2]), *parts[0][1:])]
        for part in parts:
            _add_script_runs(p, part[0], part[1], part[2], part[3])
    tb.height = Inches(_stack_drawn_h(blocks, w))
    return tb, ink_total


def add_bullets(slide, x, y, w, h, items, size, color, *, line_spacing=1.3,
                space_after_pt=8, anchor=MSO_ANCHOR.TOP, weight=400):
    """箇条書き: 段落の bullet(buChar)+ぶら下げインデント。折返し行は記号の右に揃う。

    記号は本文と同じ段落に属するので、行が増減しても、編集されても、どのビューアでも
    本文と一緒に動く。縦位置は和文中黒の字形設計で字面中央に乗る(上の注記)。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    indent = Inches(BULLET_INDENT_IN)
    # 折返しは記号の右(marL)から始まるので、本文が使える幅はぶら下げぶん狭い
    text_w = max(0.2, w - BULLET_INDENT_IN)
    for i, txt in enumerate(items):
        txt = display_wrap_text(txt, text_w, size, weight)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # 段落後スペースは項目の「間」にだけ置く。最終段落にも付けると、中身の下に
        # 見えない余白が1つ増え、縦中央寄せ(MIDDLE)がその分だけ上へずれる
        # (LibreOffice は末尾を除外するが、PowerPoint は数えるビューアがある)
        if i < len(items) - 1:
            p.space_after = Pt(space_after_pt)
        _add_script_runs(p, txt, size, weight, color)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(indent)))
        pPr.set("indent", str(-int(indent)))   # ぶら下げ: 行頭に記号、折返しは marL へ
        buClr = etree.SubElement(pPr, f"{A_NS}buClr")
        etree.SubElement(buClr, f"{A_NS}srgbClr").set("val", str(color))
        etree.SubElement(pPr, f"{A_NS}buSzPct").set("val", BULLET_SIZE_PCT)
        etree.SubElement(pPr, f"{A_NS}buFont").set("typeface", FONTS["ea"])
        etree.SubElement(pPr, f"{A_NS}buChar").set("char", BULLET_CHAR)
    return tb


P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def _strip_style(sh):
    """Remove <p:style> so LibreOffice/PowerPoint don't apply theme effects (shadows)."""
    el = sh._element
    style = el.find(f"{P_NS}style")
    if style is not None:
        el.remove(style)


def add_rect(slide, x, y, w, h, fill: RGBColor | None, *, line: RGBColor | None = None,
             line_w_pt: float = 0.75, radius_pt: float = 0.0, dash: str | None = None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_pt > 0 else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius_pt > 0:
        frac = min(0.5, (radius_pt / 72.0) / max(0.01, min(w, h)))
        sh.adjustments[0] = frac
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w_pt)
        if dash:
            ln = sh.line._get_or_add_ln()
            d = ln.find(f"{A_NS}prstDash")
            if d is None:
                d = etree.SubElement(ln, f"{A_NS}prstDash")
            d.set("val", dash)
    sh.shadow.inherit = False
    _strip_style(sh)
    return sh


def add_line(slide, x1, y1, x2, y2, color: RGBColor, w_pt: float = 0.75, dash: str | None = None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w_pt)
    ln.shadow.inherit = False
    _strip_style(ln)
    if dash:
        el = ln.line._get_or_add_ln()
        d = etree.SubElement(el, f"{A_NS}prstDash")
        d.set("val", dash)
    return ln


def add_chevron(slide, x, y, w, h, fill: RGBColor, tip_in: float = 0.30):
    # 矢羽は常に左辺フラットの PENTAGON(ホームベース形)。左に窪みのある CHEVRON は
    # 「前工程から食い込まれる」形で、Act ではどの位置のステップにも使わない(全パターン共通)
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    _strip_style(sh)
    # 矢先の水平深さを高さに依らず一定(~tip_in)に保つ。PENTAGON の既定調整値は
    # 矢先深さ=min(w,h)×adj のため、背の高い矢羽(roadmap の 2 行ラベル)だと矢先が
    # 深く尖りすぎる。process_flow の浅い矢先感に揃えるため adj を実寸から算出する。
    try:
        ss = min(w, h)
        if ss > 0:
            sh.adjustments[0] = max(0.12, min(0.5, tip_in / ss))
    except Exception:
        pass
    return sh


def _delta_color(direction: str | None, positive_is_good: bool = True) -> RGBColor:
    if direction == "up":
        return C["success"] if positive_is_good else C["danger"]
    if direction == "down":
        return C["danger"] if positive_is_good else C["success"]
    return C["ink_subtle"]


def _kpi_delta_text(k: dict) -> str:
    """The KPI delta string exactly as rendered. Signed (+/-/△/▲/▼) and comparison
    ('前年同期', '業界' …) deltas keep their own marker; only a bare magnitude gets a
    direction arrow. Shared by the height budget and the placement so a long delta
    reserves the right number of lines and never collides with the note below it."""
    delta_text = str(k.get("delta", "")).lstrip()
    is_comparison_text = delta_text.startswith(("前年同期", "前期", "業界", "目標", "計画", "中央値"))
    arrow = "" if delta_text.startswith(("+", "-", "△", "▲", "▼")) or is_comparison_text else \
        {"up": "▲ ", "down": "▼ "}.get(k.get("delta_dir"), "")
    return arrow + delta_text


# ---------------------------------------------------------------- chrome

def slot_lines(pattern: str, slot: str, text: str) -> int:
    """見出しスロットが実際に占める行数。契約(header_contract)が宣言する行数を下限に、
    実測の折返し行数と突き合わせる。validate_spec を通ったスペックでは常に契約どおりの
    行数になる — 実測側は、検証を経ずに build へ回されたスペックで文字が箱から
    はみ出さないための保険であって、字数超過を許容する意味ではない。"""
    cfg = next(c for c in header_slots(pattern) if c["slot"] == slot)
    measured = _text_lines(text, cfg["width_in"], cfg["size_pt"],
                           700 if slot == "title" else 400) if text else 0
    declared = cfg["lines"] if text else 0
    return min(3, max(declared, measured, text.count("\n") + 1 if text else 0))


def title_lines(title: str, pattern: str = "default") -> int:
    return slot_lines(pattern, "title", title)


def header_offset(spec: dict) -> float:
    return 0.26 if spec.get("kicker") else 0.0


def header_metrics(spec: dict) -> dict:
    """Deterministic header geometry: title cap-top to the bottom of the last
    header line (title or subtitle) fixes where the body region starts.

    行数はヘッダー契約(slot_lines)から取る — 幅・級数の出所を validate と共有し、
    箱の高さだけが別式で決まる状態を作らない。"""
    hdr = LAY["header"]
    off = header_offset(spec)
    pattern = spec.get("pattern", "default")
    title = spec.get("title", "")
    lines = title_lines(title, pattern)
    # タイトルサイズは行数によらず一定。縮小して収める操作は AutoFit と同種の一貫性破壊
    # (スライド間でタイトルの見た目が揺れる)。契約上タイトルは常に1行だが、検証を経ずに
    # build へ回されたスペックでも文字が箱からはみ出さないよう、実測行数で箱を高くする。
    t_size = TS["action_title"]
    title_y = hdr["title_y_in"] + off
    title_h = lines * (t_size / 72.0) * 1.27
    sub_y = title_y + title_h + hdr["title_subtitle_gap_in"]
    sub_lines = slot_lines(pattern, "subtitle", spec.get("subtitle", ""))
    sub_h = max(1, sub_lines) * (TS["subtitle"] / 72.0) * 1.35
    block_bottom = sub_y + sub_h if spec.get("subtitle") else title_y + title_h
    return {
        "off": off, "lines": lines, "t_size": t_size,
        "title_y": title_y, "title_h": title_h,
        "sub_y": sub_y, "sub_h": sub_h,
        "body_y": block_bottom + hdr["body_gap_in"],
    }


def add_chrome(slide, spec: dict, page_no: int, total: int, deck: dict) -> None:
    """Header (optional kicker + action title + subtitle) and footer.
    ヘッダーに縦アクセントバーは置かない — タイトル・サブタイトルの級差だけで階層を作る。
    背景は敷かない — 全面を覆う矩形は、編集時に本文の下で毎回つかんでしまう邪魔な
    オブジェクトになるだけで、地の色はスライドの背景に任せればよい。"""
    hdr = LAY["header"]
    m = header_metrics(spec)
    if spec.get("kicker"):
        add_text(slide, hdr["title_x_in"], 0.18, hdr["title_w_in"], 0.26,
                 [[(spec["kicker"], 13, 600, C["ink_faint"])]])
    add_text(slide, hdr["title_x_in"], m["title_y"], hdr["title_w_in"], m["title_h"],
             [[(spec.get("title", ""), m["t_size"], 700, C["ink"])]], line_spacing=1.16)
    if spec.get("subtitle"):
        add_text(slide, hdr["title_x_in"], m["sub_y"], hdr["title_w_in"], m["sub_h"],
                 [[(spec["subtitle"], TS["subtitle"], 400, C["ink_subtle"])]], line_spacing=1.18)
    add_footer(slide, spec, page_no, total, deck)


def add_footer(slide, spec: dict, page_no: int, total: int, deck: dict) -> None:
    foot = LAY["footer"]
    text = footer_text(spec)      # 描画も検証も deck_text の単一実装を見る
    if text:
        frags = [(text, TS["footnote"], 400, C["ink_faint"])]
        # フッターは max_lines 行まで。1行に収まらない出典・前提・注記は折り返して2行目へ
        # 落ちる — 行送りと箱の高さをトークンから取り、下端の外周パディングを侵さない。
        # 2行を超える長さは仕様側の欠陥(validate_spec が footnote_max_chars_ja で弾く)。
        # 帯の中では縦中央に置く: 上寄せだと1行のときだけ帯の上端に張り付き、2行のときと
        # 光学的な位置が変わってしまう。中央寄せなら1行は2行ぶんの帯の中間に座る
        add_text(slide, foot["source_x_in"], foot["y_in"], CONTENT_W, foot["h_in"], [frags],
                 line_spacing=foot["line_spacing"], anchor=MSO_ANCHOR.MIDDLE)



def body_region(spec: dict) -> tuple[float, float]:
    """(y, h) available for body content, honoring kicker/insight if present.

    The content area is vertically centered between the header block and the bottom edge
    (footer top, or the insight band when present): the same breathing gap `body_gap_in`
    is applied BELOW the header and ABOVE the bottom edge. Symmetric margins are what make
    any pattern that centers its content (tables, 2x2 rails, hero clusters) read as
    centered instead of drifting downward. Only the floor moves versus the old asymmetric
    margin — centered content rises to the true midpoint and top-anchored content simply
    gains clearance above the footer, so nothing can overflow."""
    gap = LAY["header"]["body_gap_in"]                       # symmetric top/bottom breathing
    y = header_metrics(spec)["body_y"]                       # already = header block bottom + gap
    floor = FOOT_Y - (0.74 if spec.get("insight") else 0.0)  # insight band reserves its strip
    # 帯は存在するときだけ予約する: フッター(source/assumption/note)も insight も無い
    # スライドでは、フッター文の上に確保していた呼吸 0.34 は何も守っていない。
    # 下端はスライド余白そのものなので小さなガードだけ残して本文領域を解放する
    # (カード束や表がこの高さを呼吸に変換できる)。
    has_footer = bool(footer_text(spec))
    bottom_gap = gap if (spec.get("insight") or has_footer) else 0.12
    h = (floor - bottom_gap) - y
    return y, h


def add_insight_bar(slide, spec: dict) -> None:
    """Optional bottom band. style 'accent' = selective judgment (pale yellow),
    style 'primary' = conclusion restating the so-what (pale petrol, centered)."""
    if not spec.get("insight"):
        return
    y = FOOT_Y - 0.74
    h = 0.58
    if spec.get("insight_style", "accent") == "primary":
        add_rect(slide, MX, y, CONTENT_W, h, C["primary_pale"], radius_pt=4)
        add_text(slide, MX + 0.18, y, CONTENT_W - 0.36, h,
                 [[(spec["insight"], TS["insight"], 600, C["primary_deep"])]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    else:
        add_rect(slide, MX, y, CONTENT_W, h, C["accent_pale"], line=C["accent_line"], line_w_pt=1.0, radius_pt=4)
        add_text(slide, MX + 0.18, y, CONTENT_W - 0.36, h,
                 [[(spec["insight"], TS["insight"], 600, C["ink"])]], anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- charts

CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
}


def _chart_ea_fonts(chart) -> None:
    """Ensure every defRPr in the chart XML carries latin+ea Act fonts."""
    root = chart._chartSpace
    ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
          "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for defRPr in root.findall(".//a:defRPr", ns):
        for tag in ("latin", "ea", "cs"):
            el = defRPr.find(f"a:{tag}", ns)
            if el is None:
                el = etree.SubElement(defRPr, f"{A_NS}{tag}")
            el.set("typeface", FONTS["latin"] if tag == "latin" else FONTS["ea"])


def add_act_chart(slide, x, y, w, h, cspec: dict):
    """Native chart with Act styling. cspec:
    {type, categories, series:[{name, values, focal|color}], unit, value_labels,
     focal_category (index highlighted for single-series), number_format}
    """
    ctype = cspec.get("type", "column")
    data = CategoryChartData()
    data.categories = [hw(str(c)) for c in cspec["categories"]]
    for s in cspec["series"]:
        data.add_series(hw(s["name"]), tuple(float(v) for v in s["values"]))
    gf = slide.shapes.add_chart(CHART_TYPES[ctype], Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = gf.chart
    chart.has_title = False
    nseries = len(cspec["series"])
    chart.has_legend = bool(cspec.get("legend", nseries > 1 and ctype != "donut"))
    if chart.has_legend:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(TS["chart_axis"])
        chart.legend.font.color.rgb = C["ink_subtle"]
    chart.font.size = Pt(TS["chart_axis"])
    chart.font.color.rgb = C["ink_subtle"]
    chart.font.name = FONTS["latin"]

    palette = [RGBColor.from_string(c) for c in CH["comparison_series_colors"]]
    focal_cat = cspec.get("focal_category")
    for si, (s, ser) in enumerate(zip(cspec["series"], chart.plots[0].series)):
        color = RGBColor.from_string(s["color"]) if s.get("color") else (
            C["chart_gray"] if s.get("focal") is False else palette[si % len(palette)])
        if ctype == "line":
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.25 if s.get("focal", si == 0) else 1.5)
            ser.smooth = False
        elif ctype == "donut":
            for pi, pt_ in enumerate(ser.points):
                pt_.format.fill.solid()
                pt_.format.fill.fore_color.rgb = palette[pi % len(palette)] if focal_cat is None else (
                    C["primary"] if pi == focal_cat else C["chart_gray"])
        else:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
            forecast_from = cspec.get("forecast_from") if nseries == 1 else None
            if nseries == 1 and (focal_cat is not None or forecast_from is not None):
                for pi, pt_ in enumerate(ser.points):
                    if forecast_from is not None and pi >= int(forecast_from):
                        # IBCS-style forecast: pale fill + brand outline (vs solid actuals).
                        # A focal category in the forecast zone stays pale (it is still a
                        # forecast) but gets a heavier deep outline for emphasis.
                        pt_.format.fill.solid()
                        pt_.format.fill.fore_color.rgb = C["primary_pale"]
                        emphasized = focal_cat is not None and pi == focal_cat
                        pt_.format.line.color.rgb = C["primary_deep"] if emphasized else C["primary"]
                        pt_.format.line.width = Pt(1.75 if emphasized else 1.0)
                        continue
                    pt_.format.fill.solid()
                    pt_.format.fill.fore_color.rgb = (
                        color if focal_cat is None or pi == focal_cat else C["chart_gray"])

    if ctype != "donut":
        plot = chart.plots[0]
        if ctype in ("column", "bar", "stacked_column"):
            plot.gap_width = CH["gap_width_pct"]
            if nseries > 1 and ctype != "stacked_column":
                plot.overlap = CH["overlap_pct"]
        cat_ax, val_ax = chart.category_axis, chart.value_axis
        show_labels = cspec.get("value_labels", nseries == 1)
        # 積み上げ棒のセグメント内ラベル(各構成値を帯の中に直接表示。IR頻出の
        # 「セグメント値+合計」型の証拠を型に依存せず出せるようにする opt-in)。
        seg_labels = bool(cspec.get("segment_labels")) and ctype == "stacked_column"
        axisless = cspec.get("axis_less", show_labels and nseries == 1 and ctype in ("column", "bar", "line"))
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = C["chart_gray_light"]
        val_ax.major_gridlines.format.line.width = Pt(CH["gridline_width_pt"])
        val_ax.format.line.fill.background()
        val_ax.major_tick_mark = XL_TICK_MARK.NONE
        val_ax.tick_labels.font.size = Pt(TS["chart_axis"])
        val_ax.tick_labels.font.color.rgb = C["ink_faint"]
        axis_fmt = cspec.get("axis_number_format") or cspec.get("number_format")
        if axis_fmt:
            val_ax.tick_labels.number_format = axis_fmt
            val_ax.tick_labels.number_format_is_linked = False
        if cspec.get("y_max") is not None:
            val_ax.maximum_scale = float(cspec["y_max"])
        elif cspec.get("annotation"):
            top = max(v for s in cspec["series"] for v in s["values"])
            val_ax.maximum_scale = top * 1.28
        if cspec.get("y_min") is None and all(
                v >= 0 for s in cspec["series"] for v in s["values"]):
            val_ax.minimum_scale = 0.0
        if cspec.get("y_min") is not None:
            val_ax.minimum_scale = float(cspec["y_min"])
        if axisless or cspec.get("hide_value_axis"):
            val_ax.visible = False
            val_ax.has_major_gridlines = False
        # 基線は遠目でも構造が読めるよう少し太く・濃く(Qwen: 軸が細い)
        cat_ax.format.line.color.rgb = C["ink_faint"]
        cat_ax.format.line.width = Pt(1.25)
        cat_ax.major_tick_mark = XL_TICK_MARK.NONE
        cat_ax.tick_labels.font.size = Pt(TS["chart_axis"])
        cat_ax.tick_labels.font.color.rgb = C["ink"]
        if show_labels or seg_labels:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(TS["chart_label"])
            dl.font.bold = True  # 主要値を太字にして視認性を上げる(Qwen)
            dl.font.color.rgb = C["ink"]
            dl.number_format = cspec.get("number_format", "#,##0.0")
            dl.number_format_is_linked = False
            if seg_labels:
                dl.position = XL_LABEL_POSITION.CENTER   # 帯の中央にセグメント値
            elif ctype in ("column", "bar"):
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            elif ctype == "line":
                dl.position = XL_LABEL_POSITION.ABOVE
    _chart_ea_fonts(chart)
    return gf


def add_unit_note(slide, x, y, unit: str | None):
    """縦軸の単位を、グラフの内側・軸側の上端に置く。

    単位はプロットを読むための添え物なので、グラフの外に別行として積むのではなく、軸の内側に
    小さく寄り添わせる。位置(インセット)はトークンが1か所で決めるので、どのグラフでも同じ位置に付く。
    枠は文字ぶんの大きさで、プロットの上端に揃う。(x, y) はグラフ領域の左上。"""
    if not unit:
        return
    cfg = LAY["chart"]["unit_note"]
    text = "(" + unit + ")"
    pt = TS["chart_axis"]
    w = _label_w(text, pt, 400)
    h = pt / 72.0 * leading(pt) + 0.02
    return add_text(slide, x + cfg["inset_x_in"], y + cfg["inset_y_in"], w, h,
                    [[(text, pt, 400, C["ink_faint"])]])


# ---------------------------------------------------------------- patterns

def p_cover(slide, spec, deck):
    m = deck.get("meta", {})
    add_rect(slide, MX, 2.28, 0.5, 0.055, C["primary"])
    # 箱の高さはタイトルの行数から取る。固定高のままだと、1行のタイトルでも箱が副題の
    # 位置まで伸び、描画は正しいのに枠だけが重なる(PowerPoint で編集すると掴み違える)
    title = spec.get("title", m.get("title", ""))
    title_h = max(0.6, slot_lines("cover", "title", title) * (TS["cover_title"] / 72.0) * 1.18 + 0.06)
    add_text(slide, MX, 2.62, CONTENT_W, title_h,
             [[(title, TS["cover_title"], 700, C["ink"])]], line_spacing=1.18)
    if spec.get("subtitle"):
        # 副題の行数は契約(header_contract: cover.subtitle.lines)が決める。ボックス高は
        # その行数から算出する — 1行分の固定高のままだと2行目が下端で切れる
        sub_lines = slot_lines("cover", "subtitle", spec["subtitle"])
        sub_h = max(0.6, sub_lines * (TS["cover_subtitle"] / 72.0) * 1.45 + 0.06)
        add_text(slide, MX, 3.95, CONTENT_W, sub_h,
                 [[(spec["subtitle"], TS["cover_subtitle"], 400, C["ink_subtle"])]],
                 line_spacing=1.3)
    meta_bits = [b for b in [spec.get("date", m.get("date")), spec.get("author", m.get("author"))] if b]
    if meta_bits:
        # メタ行(日付・作成者)は出所情報であり強調対象ではない。weight 400 で組む —
        # SemiBold の数字グリフは幅が広く、和文中では全角数字のように見えてしまう
        add_text(slide, MX, 6.48, CONTENT_W, 0.4,
                 [[("  |  ".join(meta_bits), TS["cover_meta"], 400, C["ink_subtle"])]])
    if m.get("confidential"):
        add_text(slide, TOKENS["slide"]["width_in"] - MX - 3.0, 0.42, 3.0, 0.3,
                 [[(m["confidential"], 11, 600, C["ink_faint"])]], align=PP_ALIGN.RIGHT)

def p_agenda(slide, spec, deck):
    items = spec.get("items", [])
    y0, h = body_region(spec)
    x, w = grid(1, 10)
    row_h = min(0.95, h / max(1, len(items)))
    y0 += max(0.0, (h - len(items) * row_h) * 0.42)
    for i, it in enumerate(items):
        label = it if isinstance(it, str) else it.get("label", "")
        desc = None if isinstance(it, str) else it.get("desc")
        y = y0 + 0.2 + i * row_h
        add_text(slide, x, y, 0.7, 0.4,
                 [[(f"{i + 1:02d}", 20, 600, C["primary"])]])
        add_text(slide, x + 0.85, y - 0.02, w - 3.0, 0.34,
                 [[(label, 17, 600, C["ink"])]])
        if desc:
            add_text(slide, x + 0.85, y + 0.36, w - 1.0, 0.3,
                     [[(desc, TS["body"], 400, C["ink_subtle"])]])
        if i < len(items) - 1:
            add_line(slide, x, y + row_h - 0.12, x + w, y + row_h - 0.12, C["rule"], 0.5)


def p_section_divider(slide, spec, deck):
    """章扉: 巨大数字(96pt、タイトル比3:1)が主役のスケールコントラスト。右は
    surface_tint の面パネル+章ナビ(現在章スポットライト)で、扉をサムネイルでも
    識別できるマクロ形状にする。左ブロックは光学中心(50%-0.25in リフト)に置く。"""
    W, H = TOKENS["slide"]["width_in"], TOKENS["slide"]["height_in"]
    chapters = deck.get("_chapters") or []
    # 章扉のパネル位置とテキスト折返し幅は tokens(layout.divider)が単一ソース —
    # header_contract の 'divider' ボックス幅もここから導出される(検証と描画が同じ幅を見る)
    dv = LAY["divider"]
    panel_x, text_pad = dv["panel_x_in"], dv["text_x_pad_in"]
    text_w = panel_x - MX - dv["text_gap_in"]
    if len(chapters) >= 2:
        add_rect(slide, panel_x, 0, W - panel_x, H, C["surface_tint"])
    num_pt = 96
    has_num = bool(spec.get("number"))
    # number 省略時は数字の段をレイアウトから除く(96pt 分の空洞をバーがブラケットしない)
    num_h = num_pt / 72.0 * 1.12 if has_num else 0.0
    num_gap = 0.08 if has_num else 0.0
    title_h = 0.55
    desc_h = 0.55 if spec.get("desc") else 0.0
    block_h = num_h + num_gap + title_h + desc_h
    by = H / 2 - block_h / 2 - 0.25
    block_cy = by + block_h / 2
    # 縦バーは先頭行のキャップトップから最終行の下端までをブラケットする(固定長にしない)
    bar_top = by + num_h * 0.16 if has_num else by + 0.06
    bar_bottom = by + num_h + num_gap + title_h + desc_h - 0.10
    add_rect(slide, MX - 0.1, bar_top, 0.075, bar_bottom - bar_top, C["primary"])
    if has_num:
        add_text(slide, MX + text_pad, by, 3.4, num_h, [[(str(spec["number"]).zfill(2), num_pt, 600, C["primary"])]])
    add_text(slide, MX + text_pad, by + num_h + num_gap, text_w, title_h,
             [[(spec.get("title", ""), TS["divider_title"], 700, C["ink"])]])
    if spec.get("desc"):
        add_text(slide, MX + text_pad, by + num_h + num_gap + title_h, text_w, desc_h,
                 [[(spec["desc"], TS["divider_desc"], 400, C["ink_subtle"])]], line_spacing=1.3)
    if len(chapters) >= 2:
        row_h = 0.5
        nav_h = len(chapters) * row_h
        nx = panel_x + 0.5
        ny = block_cy - nav_h / 2
        for num, title in chapters:
            cur = title == spec.get("title", "")
            add_text(slide, nx, ny, 0.6, 0.32,
                     [[(str(num).zfill(2), 15, 600, C["primary"] if cur else C["ink_faint"])]])
            add_text(slide, nx + 0.6, ny + 0.01, W - nx - 1.1, 0.32,
                     [[(title, 15, 700 if cur else 400, C["ink"] if cur else C["ink_faint"])]])
            ny += row_h


def p_executive_summary(slide, spec, deck):
    points = spec.get("points", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(points))
    # 見出しと本文はひとつの意味ペア: ペア内ギャップ(space_after)はカード間ギャップより
    # 必ず小さく保つ(近接の原則)。カード高は等分割ではなく実測内容+上下対称の内側余白
    # から導出する(和文行ボックス補正 1.13 込み)。等分割はカード内の上下余白が内容量に
    # 依存して痩せ、文字がカード縁に近接する。
    GAP = 0.08           # カード間ギャップ(サブタイトル込みヘッダーと2行本文を両立する床値)
    PAD = 0.12           # カード内の上下余白の最小値(対称)
    PAIR_GAP_PT = 3      # 見出し↔本文のペア内ギャップ(pt)
    LS = 1.14            # カード内はコンパクトな行送り(2行サブメッセージ+呼吸を両立)
    has_kicker = any(p.get("kicker") for p in points)
    kx_off = 2.7 if has_kicker else 0.34
    tw = w - kx_off - 0.34
    line_head = TS["section_heading"] / 72.0 * LS * 1.22
    line_body = TS["body"] / 72.0 * LS * 1.22
    def _content_h(p):
        hl = _text_lines(p.get("heading", ""), tw, TS["section_heading"], 600)
        bl = _text_lines(p.get("body", ""), tw, TS["body"], 400)
        return hl * line_head + PAIR_GAP_PT / 72.0 + bl * line_body
    # カードは帯シェア(領域をN等分した高さ)まで使い切る(fill)。中身は MIDDLE アンカー
    # なので余剰高はそのまま上下対称の呼吸になる。内容+最小パッドがシェアを超える場合に
    # カードは伸ばせない(領域が上限) — それはコピー過多であり、スペック側で本文を削る。
    content = max(_content_h(p) for p in points) if points else 0.7
    share = (h - GAP * (n - 1)) / n
    card_h = min(1.95, share)
    if content + PAD * 2 > card_h:
        # 内容がシェアに収まらない = コピー過多のシグナル。描画は続けるが余白は保証されない
        card_h = min(share, content + PAD * 2)
    total = card_h * n + GAP * (n - 1)
    y = y0 + max(0.0, (h - total) * 0.5)
    for pt_ in points:
        add_rect(slide, x, y, w, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        add_rect(slide, x, y, 0.055, card_h, C["primary"])
        kx = x + 0.34
        if pt_.get("kicker"):
            add_text(slide, kx, y, 2.2, card_h,
                     [[(pt_["kicker"], TS["body"], 700, C["primary_deep"])]], anchor=MSO_ANCHOR.MIDDLE)
            kx = x + 2.7
        add_text(slide, kx, y + 0.06, x + w - kx - 0.34, card_h - 0.12,
                 [[(pt_.get("heading", ""), TS["section_heading"], 600, C["ink"])],
                  [(pt_.get("body", ""), TS["body"], 400, C["ink_subtle"])]],
                 line_spacing=LS, space_after_pt=PAIR_GAP_PT, anchor=MSO_ANCHOR.MIDDLE)
        y += card_h + GAP


def p_kpi_dashboard(slide, spec, deck):
    kpis = spec.get("kpis", [])
    if not kpis:
        return
    y0, h = body_region(spec)
    n = len(kpis)
    cols = 4 if n >= 7 else (3 if n >= 5 else min(n, 4) or 1)
    rows = math.ceil(n / cols)
    gut = LAY["gutter_in"]
    cw = (CONTENT_W - gut * (cols - 1)) / cols
    # 1行×4枚以下はヒーロースケール — まばらなスライドは装飾でなくサイズで埋める
    vsize = TS["kpi_value_hero"] if rows == 1 and n <= 4 else (TS["kpi_value"] if rows == 1 else 27)
    pad = OPT["inset_in"]
    gap = OPT["gap_in"]
    text_w = cw - 2 * pad
    any_focal_kpi = any(kk.get("focal") for kk in kpis)

    def card_blocks(k):
        """1枚ぶんの縦積み。値の上のラベルと下の注記は同じ gap で挟む(値の上下を対称に)。
        デルタは値の従属行(metric subline)なので、より詰めた gap で値に寄せる。"""
        vcolor = C["primary_deep"] if (k.get("focal") or not any_focal_kpi) else C["ink"]
        vparts = [(k.get("value", ""), vsize, 700, vcolor)]
        if k.get("unit"):
            vparts.append(_unit_part(k["unit"], vsize * 0.45))
        blocks = [
            {"parts": [(k.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])],
             "size": TS["kpi_label"], "kind": "text"},
            {"parts": vparts, "size": vsize, "kind": "numeral",
             "gap_before": gap["value_meta"], "gap_name": "value_meta"},
        ]
        if k.get("delta"):
            dtext = _kpi_delta_text(k)
            blocks.append({
                "parts": [(dtext, TS["kpi_sub"], 600,
                           _delta_color(k.get("delta_dir"), k.get("positive_is_good", True)))],
                "size": TS["kpi_sub"], "kind": "text", "gap_before": gap["metric_subline"],
                "gap_name": "metric_subline"})
        if k.get("note"):
            blocks.append({
                "parts": [(k["note"], TS["kpi_sub"], 400, C["ink_faint"])],
                "size": TS["kpi_sub"], "kind": "text", "gap_before": gap["value_meta"],
                "gap_name": "value_meta"})
        return blocks

    stacks = [card_blocks(k) for k in kpis]
    inner_hs = [stack_optical_height(b, text_w) for b in stacks]
    content_h = max(inner_hs) + 2 * pad
    rh = min((h - gut * (rows - 1)) / rows, max(1.30, content_h))
    # 器と中身のサイズを釣り合わせる: 多段ダッシュボードは本文領域を充填して浮いた帯を防ぐ。
    # 1段(4枚以下)は内容高のカードのまま、行そのものを本文内で中央寄せし、余白はカード内では
    # なく行の外側に置く(背の高い空カードを作らない)
    if rows > 1:
        rh = max(rh, min((h - gut * (rows - 1)) / rows, (0.58 * h - gut * (rows - 1)) / rows))
    else:
        rh = max(rh, min(h - 0.20, 0.50 * h))
    block_h = rows * rh + (rows - 1) * gut
    y_base = y0 + 0.05 + max(0.0, (h - block_h) * 0.44)
    for i, k in enumerate(kpis):
        r_, c_ = divmod(i, cols)
        x = MX + c_ * (cw + gut)
        y = y_base + r_ * (rh + gut)
        # focal: the title's protagonist KPIs get the pale-brand card (主役を立てる)
        add_rect(slide, x, y, cw, rh,
                 C["primary_pale"] if k.get("focal") else C["surface_tint"],
                 radius_pt=LAY["card"]["radius_pt"])
        # インクの隙間で組む: 値の上下(ラベル・注記)が同じ余白で座り、上下インセットも対称
        stack_block(slide, x + pad, y + pad, text_w, rh - 2 * pad, stacks[i])


def _focal_bar_anchor(cx, cw, cspec):
    """Rough x-center of the focal category's bar inside a rendered chart frame."""
    n = max(1, len(cspec.get("categories", [])))
    i = cspec.get("focal_category")
    if i is None:  # 明示的な null もここに落ちる(get のデフォルトでは拾えない)
        i = n - 1
    plot_x, plot_w = cx + 0.15, cw - 0.35
    return plot_x + plot_w * (i + 0.5) / n


def add_chart_annotations(slide, cx, cy, cw, ch, cspec):
    """Exemplar-style annotations: YoY badge above the focal bar, optional trend arrow."""
    ann = cspec.get("annotation") or {}
    if not ann:
        return
    bx = _focal_bar_anchor(cx, cw, cspec)
    if ann.get("yoy") or ann.get("badge"):
        bw, bh = (1.5, 0.52) if ann.get("badge") else (1.32, 0.56)
        ox = min(bx - bw / 2, cx + cw - bw - 0.05)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ox), Inches(cy - 0.06), Inches(bw), Inches(bh))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["canvas"]
        sh.line.color.rgb = C["primary"]
        sh.line.width = Pt(1.5)
        sh.shadow.inherit = False
        _strip_style(sh)
        parts = [(ann["badge"], TS["chart_label"], 700, C["primary_deep"])] if ann.get("badge") else [
            ("YoY ", TS["chart_axis"], 400, C["ink_subtle"]), (ann["yoy"], TS["chart_label"], 700, C["primary_deep"])]
        add_text(slide, ox, cy - 0.06, bw, bh, [parts],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    if ann.get("trend_arrow"):
        aw = min(1.15, cw * 0.16)
        sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(bx - aw - 0.5), Inches(cy + 0.5), Inches(aw), Inches(0.3))
        sh.adjustments[0] = 0.55
        sh.adjustments[1] = 0.55
        sh.rotation = -20
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["primary"]
        sh.line.fill.background()
        sh.shadow.inherit = False
        _strip_style(sh)


# Image-asset track: objects the native pptx engine cannot faithfully draw are rendered as
# Act-styled deterministic images (matplotlib/Graphviz via act_assets) and embedded. The router
# below only ESCALATES to an image; native stays the default. See data-and-diagram-rules.md.
_ASSET_ROOT = None  # set in build(): dir where generated assets + asset-manifest.json are cached
IMAGE_ASSET_KINDS = frozenset({
    "combo", "area", "radar", "scatter", "bubble", "waterfall", "line_multi",
    "org_tree", "node_graph", "ring", "funnel", "pyramid", "venn", "matrix"})


def _asset_kind(obj):
    """Return the image-asset kind if `obj` must route to the image backend, else None.
    Pure membership test — does NOT import matplotlib, so native-only decks stay light."""
    if isinstance(obj, dict):
        k = obj.get("kind")
        # membership のみで判定する。render:"image" だけでは未知 kind を画像経路へ
        # 通さない — act_assets.render_asset は固定の kind 集合しか受けず、通すと
        # ビルドが ValueError で落ちる(validate_spec が組み合わせ自体をエラーにする)
        if k in IMAGE_ASSET_KINDS:
            return k
    return None


def place_asset(slide, spec, x, y, w, h):
    """Render an Act-styled image asset for `spec` and embed it at the box (inches).
    The figure is rendered at the box's inch dimensions so aspect matches — no distortion."""
    import act_assets  # lazy: only pulled when an image object is actually present
    root = _ASSET_ROOT or Path(".")
    png = act_assets.render_asset(spec, root, size_in=(max(1.0, w), max(1.0, h)))
    slide.shapes.add_picture(str(png), Inches(x), Inches(y), Inches(w), Inches(h))


def p_chart_insight(slide, spec, deck):
    y0, h = body_region(spec)
    layout = spec.get("layout", "chart_left")
    takeaways = spec.get("takeaways", [])
    chart = spec["chart"]
    akind = _asset_kind(chart)
    if layout == "chart_full" or not takeaways:
        cx, cw = grid(0, 12)
        if akind:
            place_asset(slide, chart, cx, y0 + 0.12, cw, h - 0.24)
        else:
            add_act_chart(slide, cx, y0 + 0.25, cw, h - 0.3, chart)
            add_unit_note(slide, cx, y0 + 0.25, chart.get("unit"))
            add_chart_annotations(slide, cx, y0 + 0.25, cw, h - 0.3, chart)
        return
    cx, cw = grid(0, 7)
    tx, tw = grid(7, 5)
    chart_y, chart_h = y0 + 0.25, h - 0.3
    if akind:
        place_asset(slide, chart, cx, chart_y, cw, chart_h)
    else:
        add_act_chart(slide, cx, chart_y, cw, chart_h, chart)
        add_unit_note(slide, cx, chart_y, chart.get("unit"))
        add_chart_annotations(slide, cx, chart_y, cw, chart_h, chart)
    _takeaways_rail(slide, spec, takeaways, tx, tw, y0, chart_y, chart_h)


# 要点レール(図表の横に置く解釈列)も、カードと同じく「インクで積む」。箱の高さを見積もって
# 積むと、見積りと実描画の差がそのままブロック間の空白の不揃いになり、1行ぶんの空きや重なりに
# 見える(字数近似なら尚更ずれる)。行数は _text_lines、位置は stack_optical — 数える実装も
# 積む実装も1つずつしか持たない。
RAIL_HEAD_PT = 15


def _rail_blocks(items):
    """見出し+本文の並びを、インク積みのブロック列にする。見出しと本文は1つの意味のかたまり
    (かたまり内の間隔 < かたまり間の間隔)。"""
    gap = OPT["gap_in"]
    blocks = []
    for i, t in enumerate(items):
        head = t.get("heading") if isinstance(t, dict) else None
        body = t.get("body") if isinstance(t, dict) else str(t)
        if head:
            blocks.append({"parts": [(head, RAIL_HEAD_PT, 600, C["ink"])], "size": RAIL_HEAD_PT,
                           "kind": "text", "gap_before": gap["item"] if i else 0.0,
                           "gap_name": "item" if i else None})
        if body:
            blocks.append({"parts": [(body, TS["body"], 400, C["ink_subtle"])], "size": TS["body"],
                           "kind": "text",
                           "gap_before": gap["heading_body"] if head else (gap["item"] if i else 0.0),
                           "gap_name": "heading_body" if head else ("item" if i else None)})
    return blocks


def _takeaways_rail(slide, spec, takeaways, tx, tw, y0, region_y, region_h):
    """Interpretation rail beside a chart/diagram: heading + rule + heading/body blocks,
    vertically centred on the exhibit region. Shared by chart_insight and diagram."""
    text_w = tw - 0.04
    blocks = _rail_blocks(takeaways)
    stack_h = stack_optical_height(blocks, text_w)
    group_h = 0.5 + _stack_drawn_h(blocks, text_w)      # 塊の高さは「描かれる高さ」で測る
    ty = max(y0 + 0.05, region_y + (region_h - group_h) / 2)
    add_text(slide, tx, ty, tw, 0.3, [[(spec.get("takeaways_heading", "要点"), 16, 600, C["primary_deep"])]])
    add_line(slide, tx, ty + 0.34, tx + tw, ty + 0.34, C["rule"], 0.75)
    stack_block(slide, tx, ty + 0.5, text_w, stack_h, blocks)


def p_diagram(slide, spec, deck):
    """Relationship/schematic diagram rendered via the image-asset backend (org tree, node
    graph, ring, funnel, pyramid, Venn, matrix). Optional interpretation rail like chart_insight."""
    y0, h = body_region(spec)
    dspec = spec.get("diagram") or {}
    takeaways = spec.get("takeaways", [])
    if not takeaways:
        cx, cw = grid(0, 12)
        place_asset(slide, dspec, cx, y0 + 0.12, cw, h - 0.24)
        return
    cx, cw = grid(0, 7)
    tx, tw = grid(7, 5)
    place_asset(slide, dspec, cx, y0 + 0.12, cw, h - 0.24)
    _takeaways_rail(slide, spec, takeaways, tx, tw, y0, y0 + 0.12, h - 0.24)


def p_chart_grid(slide, spec, deck):
    """Small-multiples: 2-4 coordinated NATIVE charts sharing one claim, tiled in a row so each
    stays editable. Each cell carries its own chart spec (with the emphasis knobs — focal_category,
    annotation.badge/trend_arrow — so per-chart CAGR/latest-bar emphasis works). Native-first: use
    this rather than an image when the sub-charts are plain; escalate a cell to the image track
    only if that chart genuinely needs it."""
    y0, h = body_region(spec)
    charts = spec.get("charts", [])[:4]
    n = max(1, len(charts))
    span = max(1, 12 // n)
    for i, cell in enumerate(charts):
        cx, cw = grid(i * span, span)
        ch = cell.get("chart", cell)
        cy = y0 + 0.05
        sub = cell.get("title")
        if sub:
            add_text(slide, cx, cy, cw, 0.5, [[(sub, 15, 600, C["ink"])]],
                     align=PP_ALIGN.CENTER, wrap=True)
            cy += 0.52
        akind = _asset_kind(ch)
        if akind:
            # セル単位の image-kind エスカレーション(deck-spec に明記)をネイティブ経路へ
            # 流さない — `kind` は `type` ではないため、add_act_chart に渡すと既定の
            # column として誤描画されるかデータ形状が失われる
            place_asset(slide, ch, cx, cy + 0.10, cw, y0 + h - (cy + 0.10) - 0.05)
        else:
            add_unit_note(slide, cx, cy, ch.get("unit"))
            ch_y = cy + 0.22
            ch_h = y0 + h - ch_y - 0.05
            add_act_chart(slide, cx, ch_y, cw, ch_h, ch)
            add_chart_annotations(slide, cx, ch_y, cw, ch_h, ch)


def p_market_sizing(slide, spec, deck):
    """TAM/SAM/SOM horizontal funnel bars with per-stage description column."""
    stages = spec.get("stages", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(stages))
    row_h = min(1.7, (h - 0.2) / n)
    y0 += max(0.0, (h - n * row_h) * 0.4)
    bar_max_w = w * 0.52
    vals = [float(s.get("numeric", 0) or 0) for s in stages]
    vmax = max(vals) if any(vals) else 1.0
    # 桁違いのファネル(TAM 35兆 vs SOM 663億)では実比率だと下段が潰れて段差が読めない。
    # 実比率(√スケール)と均等ステップ(1.0→0.4)の大きい方を採り、隣接段には必ず差を残す
    fracs = []
    for i in range(n):
        raw = ((vals[i] / vmax) ** 0.5) if vmax else 1.0
        even = 1.0 - (0.6 * i / max(1, n - 1))
        f = max(raw, even)
        if i > 0:
            f = min(f, fracs[i - 1] - 0.12)
        fracs.append(max(f, 0.2))
    for i, s in enumerate(stages):
        y = y0 + 0.1 + i * row_h
        bar_h = row_h - 0.28
        bw = bar_max_w * fracs[i]
        fill = [C["chart_gray"], C["primary"], C["primary_deep"]][i % 3]
        txt = C["ink"] if i % 3 == 0 else C["canvas"]  # 白文字はグレー地でコントラスト不足
        add_rect(slide, x, y, bw, bar_h, fill, radius_pt=4)
        # オブジェクトに載る/併記するテキストは、オブジェクト高に対して垂直中央に置く
        add_text(slide, x + 0.18, y, bw - 0.3, bar_h,
                 [[(s.get("label", ""), TS["chart_axis"], 600, txt)],
                  [(s.get("value", ""), 24, 700, txt)]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after_pt=2)
        dx = x + bar_max_w + 0.35
        add_text(slide, dx, y, x + w - dx, bar_h,
                 [[(s.get("name", s.get("label", "")), 16, 600, C["ink"])],
                  [(s.get("desc", ""), TS["body"], 400, C["ink_subtle"])]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.26, space_after_pt=3)


# 表を閉じる罫(最終行の下)。行間のヘアライン(0.5pt)より太くして、表の下端を締める
TABLE_CLOSING_RULE_PT = 1.5


RULE_TOKEN = "__rule__"   # bottom_hex の既定: rule トークンを使う(色の直値を持たない)
NO_RULE = ""              # 下罫を引かない。None は「既定=rule トークン」の意味なので使わない


def _cell_borders(cell, bottom_hex: str = RULE_TOKEN, width_pt: float = 0.5,
                  dash: str | None = None):
    """Kill vertical/top borders; keep a thin bottom rule (banker table look).
    width_pt controls the bottom rule weight (header uses a heavier accent line).
    dash: prstDash 値(例 "dash")。行グループ境界のセパレータなど、通常のヘアラインと
    役割の違う罫を破線で区別するときに使う。

    bottom_hex: 既定(RULE_TOKEN)は rule トークンの色 — 直値を持つとパレット変更が
    テーブルに届かない。下罫を消したいときは NO_RULE を渡す(None ではない)。"""
    if bottom_hex == RULE_TOKEN:
        bottom_hex = str(C["rule"])
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        el = tcPr.find(f"{A_NS}{tag}")
        if el is not None:
            tcPr.remove(el)
    els = []
    for tag in ("lnL", "lnR", "lnT"):
        ln = etree.SubElement(tcPr, f"{A_NS}{tag}")
        etree.SubElement(ln, f"{A_NS}noFill")
        els.append(ln)
    lnB = etree.SubElement(tcPr, f"{A_NS}lnB")
    if bottom_hex:
        lnB.set("w", str(int(width_pt * 12700)))
        fill = etree.SubElement(lnB, f"{A_NS}solidFill")
        etree.SubElement(fill, f"{A_NS}srgbClr").set("val", bottom_hex)
        if dash:
            etree.SubElement(lnB, f"{A_NS}prstDash").set("val", dash)
    else:
        etree.SubElement(lnB, f"{A_NS}noFill")
    els.append(lnB)
    for i, el in enumerate(els):
        tcPr.remove(el)
        tcPr.insert(i, el)


def _table_font(cell, text, size, weight, color, align=PP_ALIGN.LEFT, width_in=None):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if width_in:
        # セルも表示テキスト。列幅で機械的に割ると「導入費＋固定利用/料」のように語が割れる
        text = display_wrap_text(str(text), width_in, size, weight)
    _add_script_runs(p, text, size, weight, color)


def _cell_text_w(col_w_in: float) -> float:
    """セル内で本文が使える幅(in)。左右マージン(0.06×2)と折返しの安全余白。
    行高の見積りと折返しの判断は同じ幅を見ること — 片方だけ広く見ると、見積りより1行増える。"""
    return max(0.4, col_w_in - 0.18)


def add_act_table(slide, x, y, w, h, tspec: dict):
    """tspec: {headers:[], rows:[[...]], col_widths:[fractions], align:[l|r|c per col],
    emphasis_col: int|None, emphasis_row: int|None}"""
    headers = tspec.get("headers", [])
    rows = tspec.get("rows", [])
    nrows, ncols = len(rows) + 1, len(headers)
    # --- 列幅(inch)。レイアウトと折返し行数推定の両方に使う ---
    if tspec.get("col_widths"):
        fr = tspec["col_widths"]
        total_fr = sum(fr) or 1.0
        col_w_in = [w * f / total_fr for f in fr]
    else:
        col_w_in = [w / ncols] * ncols
    # --- 内容量ドリブンの行高(面をストレッチで埋めず、折返し行数から実測して算出) ---
    # ヘッダーも本文行と同じ級差ユニットで組み、本文行が複数行で高いときはヘッダーも
    # 比例して高くする(級差のバランス)。表全体は body 内で垂直中央に置く。
    cell_pt, hdr_pt = TS["table_cell"], TS["table_header"]
    line_h = cell_pt / 72.0 * 1.34 * 1.13   # 1 行あたりの高さ(和文行ボックス補正込み)
    v_pad = 0.30                            # 1 行の内容最小余白(上下対称)。痩せると文字が罫に近接
    emph_col, emph_row = tspec.get("emphasis_col"), tspec.get("emphasis_row")

    def _cell_weight(ri: int, ci: int) -> int:
        """そのセルを実際に描く重み。行数見積りは描画と同じ重みで測る(強調列は太い)。"""
        if ci == 0:
            return 700
        focal = (emph_col is not None and ci == emph_col) or (emph_row is not None and ri == emph_row)
        return 600 if focal else 400

    def _row_lines(cells, ri: int):
        ln = 1
        for ci in range(ncols):
            txt = str(cells[ci]) if ci < len(cells) else ""
            ln = max(ln, _text_lines(txt, _cell_text_w(col_w_in[ci]), cell_pt, _cell_weight(ri, ci)))
        return ln
    row_min = [_row_lines(r, ri) * line_h + v_pad for ri, r in enumerate(rows)] or [line_h + v_pad]
    # ヘッダーも本文と同じく文節で折り返す — 1行ぶんで高さを取ると、折返したぶんだけ表が
    # 計画より下へ伸びる(レンダラが行を足す)。ヘッダーの実測行数から高さを取る
    hdr_rows = max((_text_lines(str(h_), _cell_text_w(col_w_in[ci]), hdr_pt, 700)
                    for ci, h_ in enumerate(headers)), default=1)
    hdr_min = hdr_rows * (hdr_pt / 72.0 * 1.34) + v_pad * 0.7
    n = len(row_min)
    sum_min = sum(row_min)
    # 余白を詰めて本文領域を埋める(全体の高さを出す)。追加高は主にデータ行へ配分し、
    # ヘッダーはデータ行より低めに保つ。1 行が過度に伸びないよう上限も設ける。
    FILL, HDR_RATIO, ROW_CAP = 0.90, 0.60, 1.9
    target = h * FILL
    if sum_min + hdr_min >= target:         # 既に埋まる/溢れる → 実測高そのまま
        hdr_h_in, body_heights = hdr_min, list(row_min)
    else:                                   # 余白ぶんをデータ行主体に配ってフィル
        hdr_h_in = max(hdr_min, HDR_RATIO * target / (n + HDR_RATIO))
        remain = target - hdr_h_in
        body_heights = [min(remain * (rm / sum_min), rm * ROW_CAP) for rm in row_min]
        hdr_h_in = min(hdr_h_in, min(body_heights) * 0.9)   # ヘッダーはデータ行より低く
        hdr_h_in = max(hdr_h_in, hdr_min)
    total = hdr_h_in + sum(body_heights)
    if total <= h:                          # 残余はわずかな上下マージンとして中央へ
        y = y + (h - total) / 2.0
    else:                                   # 超えるなら比率維持で全行を縮小して収める
        k = h / total
        hdr_h_in *= k
        body_heights = [bh * k for bh in body_heights]
        total = h  # スケール後の実寸。旧 total のままだとフレームだけ縦に溢れフッターへ食い込む
    g = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(total))
    table = g.table
    table.first_row = False
    table.horz_banding = False
    tbl = g._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # "no style" table style
    for ci, cwin in enumerate(col_w_in):
        table.columns[ci].width = Emu(int(Inches(cwin)))
    aligns = tspec.get("align", ["l"] * ncols)
    amap = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER}
    color_neg = bool(tspec.get("color_negatives"))  # △/▲/▼/- 始まりの数値セルを danger 色に
    table.rows[0].height = Inches(hdr_h_in)
    for ri in range(1, nrows):
        table.rows[ri].height = Inches(body_heights[ri - 1])
    # ヘッダーは濃緑の塗り帯はやめつつ、遠目でも構造が読めるよう淡い面(surface_tint)で
    # 行を定義し、primary_deep の太字テキスト+1.75pt の primary_deep 下線で締める
    # (面で押さず、淡面+濃い罫+級差で階層を作る)。本文行はゼブラを外しヘアライン区切り、
    # 文字色は #2D332E(ink)に固定して役割を明確にする(Qwen: ヘッダー弱い/本文色を固定)。
    for ci, htxt in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["surface_tint"]
        _table_font(cell, str(htxt), TS["table_header"], 700, C["primary_deep"],
                    amap.get(aligns[ci], PP_ALIGN.LEFT), width_in=_cell_text_w(col_w_in[ci]))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _cell_borders(cell, str(C["primary_deep"]), 1.75)  # ヘッダー下線(濃いアクセント罫)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            focal = (emph_col is not None and ci == emph_col) or (emph_row is not None and ri - 1 == emph_row)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["primary_pale"] if focal else C["canvas"]
            weight = 700 if ci == 0 else (600 if focal else 400)
            color = C["ink"]  # 本文は #2D332E に固定(ラベル/値の別は級差で表す)
            if color_neg and ci != 0 and str(val).lstrip().startswith(("△", "▲", "▼", "-", "−")):
                color = C["danger"]  # 減益/減少など負値の下方向を色でも冗長に伝える
            _table_font(cell, str(val), TS["table_cell"], weight, color,
                        amap.get(aligns[ci], PP_ALIGN.LEFT), width_in=_cell_text_w(col_w_in[ci]))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 最終行の下だけは表を閉じる罫なので、行間のヘアラインより少し太くする
            _cell_borders(cell, width_pt=TABLE_CLOSING_RULE_PT if ri == nrows - 1 else 0.5)
    # 行グループの見出しセルを縦結合(例: 「電機関連」が売上高/利益の2行にまたがる)。
    # col0_spans: [[start_data_row, length], ...] — 0-based。ラベルは各グループ先頭行に置き、
    # 残り行は空文字にしておく(結合セルは先頭セルのテキストを保持)。編集可能なネイティブ表のまま。
    for span in tspec.get("col0_spans", []):
        st, ln = span[0], span[1]
        if ln > 1 and 1 + st + ln - 1 < nrows:
            origin = table.cell(1 + st, 0)
            origin.merge(table.cell(1 + st + ln - 1, 0))
            origin.vertical_anchor = MSO_ANCHOR.MIDDLE
    # group_dividers: 第1列グループの切れ目が追いにくい表(行数が多い/他の視覚手掛かりが
    # ない)に限り、各グループ末行の下に破線グレー(chart_gray)のセパレータを引く。
    # 常時適用はしない — 判断はスペック側(group_dividers: true)で明示する。
    # 実装はセル罫線の prstDash ではなくオーバーレイのコネクタ線: テーブルセル罫線の破線は
    # LibreOffice では描画されるが PowerPoint 系ビューアで無視されることがあり、
    # シェイプ(コネクタ)の破線はどのビューアでも確実に描画される(pptx-pitfalls 参照)。
    if tspec.get("group_dividers"):
        ends = {sp[0] + sp[1] - 1 for sp in tspec.get("col0_spans", [])}
        yy_l = y + hdr_h_in
        for ri, bh in enumerate(body_heights):
            yy_l += bh
            if ri in ends and ri < len(body_heights) - 1:
                # 破線は「役割の違い」を示す線であり主張ではない — 通常ヘアラインと
                # 同じ 0.5pt(色と破線だけで役割を区別し、太さでは主張しない)
                add_line(slide, x, yy_l, x + w, yy_l, C["chart_gray"], 0.5, dash="dash")
    return g


def p_comparison_table(slide, spec, deck):
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    add_act_table(slide, x, y0 + 0.08, w, h - 0.16, spec["table"])


def _free_below(placed, x, y, w, h, gap=0.04, limit=None):
    """すでに置いたラベルと重なるなら、重ならない高さまで下げた y を返す。

    プロットの位置はデータが決めるので、近い2点のラベルは自然に重なる。下へ逃がして両方を
    読めるようにするが、逃げ先は図の中に留める(limit)。留まれないほど密なら、それはコピーでは
    なくデータの見せ方の問題なので、元の位置に置いて検証の警告に委ねる。"""
    y0 = y
    for _ in range(8):
        hit = next((b for b in placed
                    if x < b[2] and b[0] < x + w and y < b[3] and b[1] < y + h), None)
        if not hit:
            return y
        y = hit[3] + gap
        if limit is not None and y + h > limit:
            return y0
    return y if (limit is None or y + h <= limit) else y0


def p_competitive_landscape(slide, spec, deck):
    """2x2 positioning map. spec: {x_axis:{low,high}, y_axis:{low,high},
    players:[{name, x, y, focal}]} coords in 0..1."""
    y0, h = body_region(spec)
    size = min(h - 0.35, 5.0)
    px, pw = grid(0, 8)
    cx = px + 1.05
    cy = y0 + 0.15
    cw_ = pw - 2.05
    add_rect(slide, cx, cy, cw_, size, C["surface_tint"], radius_pt=0)
    add_line(slide, cx, cy + size / 2, cx + cw_, cy + size / 2, C["chart_gray"], 0.75)
    add_line(slide, cx + cw_ / 2, cy, cx + cw_ / 2, cy + size, C["chart_gray"], 0.75)
    ax, ay = spec.get("x_axis", {}), spec.get("y_axis", {})
    add_text(slide, cx - 1.02, cy + size / 2 - 0.10, 0.94, 0.34, [[(ax.get("low", ""), TS["chart_axis"], 400, C["ink_faint"])]], align=PP_ALIGN.RIGHT)
    add_text(slide, cx + cw_ + 0.08, cy + size / 2 - 0.10, 0.92, 0.34, [[(ax.get("high", ""), TS["chart_axis"], 400, C["ink_faint"])]])
    add_text(slide, cx + cw_ / 2 - 1.0, cy - 0.30, 2.0, 0.30, [[(ay.get("high", ""), TS["chart_axis"], 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)
    add_text(slide, cx + cw_ / 2 - 1.0, cy + size + 0.05, 2.0, 0.30, [[(ay.get("low", ""), TS["chart_axis"], 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)
    placed: list[tuple[float, float, float, float]] = []
    for pl in spec.get("players", []):
        dot = 0.5 if pl.get("focal") else 0.3
        dx = cx + pl["x"] * cw_ - dot / 2
        dy = cy + (1 - pl["y"]) * size - dot / 2
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(dot), Inches(dot))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["primary"] if pl.get("focal") else C["chart_gray"]
        sh.line.fill.background()
        sh.shadow.inherit = False
        _strip_style(sh)
        # 枠は文字ぶんの大きさを持つ — 幅は実測、高さは折返し行数から。余った幅を枠に含めると、
        # 近くのプロットのラベルと枠だけが重なる(文字は離れていても、編集で掴み違える)
        lab_pt = TS["body_small"]
        weight = 600 if pl.get("focal") else 400
        lab_w = _label_w(pl["name"], lab_pt, weight)          # 実測幅 = 1行に収まり、余分な枠を持たない
        lab_h = lab_pt / 72.0 * leading(lab_pt) + 0.02
        lx, ly = dx + dot / 2 - lab_w / 2, dy + dot + 0.02
        # 逃げ先は図の中(下の軸ラベルの手前)まで
        ly = _free_below(placed, lx, ly, lab_w, lab_h, limit=cy + size - 0.02)
        placed.append((lx, ly, lx + lab_w, ly + lab_h))
        add_text(slide, lx, ly, lab_w, lab_h,
                 [[(pl["name"], lab_pt, weight,
                    C["ink"] if pl.get("focal") else C["ink_subtle"])]], align=PP_ALIGN.CENTER)
    notes = spec.get("notes", [])
    if notes:
        nx, nw = grid(8, 4)
        blocks = _rail_blocks(notes)
        stack_h = stack_optical_height(blocks, nw)
        # 右の要点列は左の2x2マップ(cy..cy+size)の垂直中心に揃える(塊は描かれる高さで測る)
        ny = max(cy, cy + (size - (0.5 + _stack_drawn_h(blocks, nw))) / 2)
        add_text(slide, nx, ny, nw, 0.3, [[(spec.get("notes_heading", "Positioning"), 16, 600, C["primary_deep"])]])
        add_line(slide, nx, ny + 0.34, nx + nw, ny + 0.34, C["rule"], 0.75)
        stack_block(slide, nx, ny + 0.5, nw, stack_h, blocks)


def p_financial_summary(slide, spec, deck):
    y0, h = body_region(spec)
    if spec.get("chart") and spec.get("table"):
        tx, tw = grid(0, 7)
        cx, cw = grid(7, 5)
        add_act_table(slide, tx, y0 + 0.08, tw - 0.15, h - 0.16, spec["table"])
        chart = spec["chart"]
        if _asset_kind(chart):
            # image-kind(combo 等)はネイティブ経路に流さない — chart_insight / chart_grid
            # と同じルーティング(native は series を要求するため KeyError/誤描画になる)
            place_asset(slide, chart, cx + 0.15, y0 + 0.15, cw - 0.15, h - 0.30)
        else:
            add_unit_note(slide, cx + 0.15, y0 + 0.02, chart.get("unit"))
            add_act_chart(slide, cx + 0.15, y0 + 0.28, cw - 0.15, h - 0.36, chart)
    elif spec.get("table"):
        p_comparison_table(slide, spec, deck)
    else:
        p_chart_insight(slide, spec, deck)


def p_waterfall(slide, spec, deck):
    """Shape-composed waterfall. items: [{label, value, kind: start|delta|end}]"""
    items = spec.get("items", [])
    if not items:
        return
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    chart_h = h - 1.0
    base_y = y0 + 0.42
    running, levels = 0.0, []
    for it in items:
        v = float(it["value"])
        if it.get("kind") == "start":
            running = v
            levels.append((0.0, v))
        elif it.get("kind") == "end":
            levels.append((0.0, v))
            running = v
        else:
            levels.append((running, running + v))
            running += v
    # 負の累計(赤字転落ブリッジ)も帯内に収める: 0 の基準線を span 内に置いてスケールする
    hi_max = max(max(a, b) for a, b in levels)
    lo_min = min(min(a, b) for a, b in levels)
    top = max(hi_max, 0.0) * 1.15 or 1.0
    bot = min(lo_min, 0.0) * 1.15
    span = (top - bot) or 1.0
    zero_y = base_y + chart_h * top / span
    n = len(items)
    slot = w / n
    bar_w = slot * 0.62
    add_line(slide, x, zero_y, x + w, zero_y, C["chart_gray"], 0.75)
    for i, (it, (lo, hi)) in enumerate(zip(items, levels)):
        v = float(it["value"])
        kind = it.get("kind", "delta")
        y_top = base_y + chart_h * (top - max(lo, hi)) / span
        bh = max(0.03, chart_h * abs(hi - lo) / span)
        bx = x + i * slot + (slot - bar_w) / 2
        if kind in ("start", "end"):
            fill = C["primary_deep"]
        else:
            fill = C["primary"] if v >= 0 else C["danger"]
        if it.get("forecast"):
            # 計画・予想のバーは実績と同じソリッド塗りにしない(IBCS)
            add_rect(slide, bx, y_top, bar_w, bh, C["primary_pale"],
                     line=C["primary_deep"], line_w_pt=1.5)
        else:
            add_rect(slide, bx, y_top, bar_w, bh, fill)
        if i > 0:
            edge = base_y + chart_h * (top - lo) / span if kind == "delta" else y_top
            add_line(slide, bx - (slot - bar_w), edge, bx, edge, C["chart_gray"], 0.5, dash="dash")
        if it.get("display") is not None:
            label_v = it["display"]
        elif kind == "delta":
            mag = f"{abs(v):,.1f}".rstrip("0").rstrip(".")
            label_v = ("△" if v < 0 else "+") + mag  # IR 慣行: 負値は - でなく △
        else:
            label_v = f"{v:,.1f}".rstrip("0").rstrip(".")
        add_text(slide, bx - slot * 0.19, y_top - 0.26, slot, 0.22,
                 [[(label_v, TS["chart_label"], 600, C["ink"])]], align=PP_ALIGN.CENTER)
        add_text(slide, x + i * slot + 0.02, base_y + chart_h + 0.08, slot - 0.04, 0.55,
                 [[(it.get("label", ""), TS["chart_axis"], 600, C["ink_subtle"])]], align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_unit_note(slide, x, y0 + 0.0, spec.get("unit"))


def p_roadmap(slide, spec, deck):
    """Phase roadmap: chevron phase headers + workstream rows with milestone text."""
    phases = spec.get("phases", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(phases))
    label_w = 0.0
    rows = spec.get("rows", [])
    if rows:
        label_w = 1.7
    pw = (w - label_w) / n
    def _phase_text(ph):
        t = ph.get("label", "")
        if ph.get("period"):
            t += "  " + ph["period"]
        return t
    notch_ratio = 0.6
    ph_h = 0.52
    inner_w = pw - 0.06 - ph_h * notch_ratio * 2
    phase_pt = TS["body_small"]
    period_pt = TS["kpi_sub"]
    two_line = any(_ja_len(_phase_text(ph)) * phase_pt / 72.0 > inner_w for ph in phases)
    if two_line:
        # 矢羽高はラベル+期間の実測行数から導出する(固定 0.78 では 3 行以上が溢れる)
        head_rows = max(_text_lines(ph.get("label", ""), inner_w, phase_pt, 600)
                        + (_text_lines(ph["period"], inner_w, period_pt, 400) if ph.get("period") else 0)
                        for ph in phases)
        ph_h = max(0.86, 0.18 + head_rows * 0.30)
    gap = LAY["gutter_in"]
    for i, ph in enumerate(phases):
        px = x + label_w + i * pw
        # 矢羽の間に小さな隙間を空け、矢先を次段へオーバーハングさせて「流れ」を明確にする
        # (Qwen: フローの視認性)。矢羽色はデフォルトで primary_deep、非注目段はやや淡い teal。
        # 最終段の右端は下のコンテンツ列の右端(セル幅 pw-0.06)とエッジロックする —
        # 最終段だけ gap ぶん短くすると帯と列の右端が揃わず「幅が足りない」欠陥に見える
        chev_w = pw if i < n - 1 else pw - 0.06
        focal_ph = spec.get("focal_phase")
        chev_fill = C["primary_deep"] if (focal_ph is None or i == focal_ph) else C["primary"]
        add_chevron(slide, px, y0 + 0.05, chev_w, ph_h, chev_fill)
        paras = [[(ph.get("label", ""), phase_pt, 600, C["canvas"])]]
        if ph.get("period"):
            if two_line:
                paras.append([(ph["period"], period_pt, 400, C["canvas"])])
            else:
                paras[0].append(("  " + ph["period"], period_pt, 400, C["canvas"]))
        tip = ph_h * notch_ratio if not two_line else 0.52 * notch_ratio
        add_text(slide, px + 0.12, y0 + 0.05, chev_w - tip - 0.12, ph_h, paras,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False, line_spacing=1.1)
    ry = y0 + ph_h + 0.25
    if rows:
        # 行高は等分割ではなく各行の実測行数(全セルの最大折返し行数)から導出する。
        # 等分割は文字数の多い行だけセル内の上下余白が痩せ、行間の呼吸が不揃いになる。
        # 余りは全行へ均等に配り、収まらない場合のみ等率で縮小する。
        avail = h - ph_h - 0.35
        cell_w = pw - 0.34
        line_cell = TS["body_small"] / 72.0 * 1.2 * 1.13  # 描画 line_spacing=1.2 + 和文行ボックス補正
        V_PAD = 0.30                                       # セル内の上下余白合計(対称)
        slot_gap = 0.12                                    # 行スロット間ギャップ
        mins = []
        for row in rows:
            ln = 1
            for c_i in range(n):
                cell = (row.get("cells") or [""] * n)[c_i] if c_i < len(row.get("cells", [])) else ""
                if cell:
                    ln = max(ln, _text_lines(str(cell), cell_w, TS["body_small"], 400))
            mins.append(ln * line_cell + V_PAD + slot_gap)
        extra = (avail - sum(mins)) / max(1, len(mins))
        if extra >= 0:
            slots = [m + extra for m in mins]
        else:
            k = avail / sum(mins)
            slots = [m * k for m in mins]
        yy = ry
        for r_i, row in enumerate(rows):
            slot = slots[r_i]
            add_text(slide, x, yy, label_w - 0.15, slot - slot_gap,
                     [[(row.get("label", ""), TS["body_small"], 600, C["ink"])]], anchor=MSO_ANCHOR.MIDDLE)
            for c_i in range(n):
                cell = (row.get("cells") or [""] * n)[c_i] if c_i < len(row.get("cells", [])) else ""
                cx_ = x + label_w + c_i * pw
                add_rect(slide, cx_, yy, pw - 0.06, slot - slot_gap, C["surface_tint"], radius_pt=4)
                if cell:
                    add_text(slide, cx_ + 0.14, yy, cell_w, slot - slot_gap,
                             [[(str(cell), TS["body_small"], 400, C["ink_subtle"])]],
                             line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
            yy += slot
    else:
        row_h = h - ph_h - 0.35
        for i, ph in enumerate(phases):
            px = x + i * pw
            add_rect(slide, px, ry, pw - 0.06, row_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
            # 箇条書き枠の上下インセットは対称にする。上端をカードに合わせて高さだけ
            # 削ると、枠の中心がカードの中心より上にずれ、MIDDLE 寄せの結果も上に浮く
            pad = 0.1
            add_bullets(slide, px + 0.2, ry + pad, pw - 0.44, row_h - 2 * pad,
                        ph.get("items", []), TS["body_small"], C["ink_subtle"],
                        line_spacing=1.25, space_after_pt=7, anchor=MSO_ANCHOR.MIDDLE)


def p_two_column(slide, spec, deck):
    y0, h = body_region(spec)
    # card height follows the measured content of the fuller side (container ≒ content),
    # then the block is centered — a full-height card over short content reads as neglect
    HEAD_H, CONTENT_TOP, ITEM_GAP = 0.54, 0.80, 0.22  # ヘッダー帯と本文に上下の余裕を与える
    # 行高の見積りは「描画と同じ line_spacing × 和文行ボックス補正 1.22(LibreOffice実測)」
    # で行う。補正が甘いと不足分がすべてカード下側の内側余白から差し引かれ、下縁だけ
    # 文字が近接する(上下余白の非対称は欠陥)。見出し行も同係数+space_after 実寸で積む。
    LINE_H = (TS["body"] / 72.0) * 1.28 * 1.22
    HEAD_LINE = (TS["body_small"] / 72.0) * 1.28 * 1.22 + 3 / 72.0
    _, w_probe = grid(0, 6)
    def _blocks(blk, w):
        out = []
        for it in blk.get("items", []):
            head = it.get("heading") if isinstance(it, dict) else None
            body = it.get("body") if isinstance(it, dict) else str(it)
            # 行数見積もりは、その item を実際に描く経路と同じ実効幅で行う:
            #   見出しあり → add_text(幅 w - 0.4)
            #   見出しなし → add_bullets(箱 w - 0.48、本文はさらに字下げ分だけ狭い)
            # 字下げ幅を直値で持たないこと — 変えたときにここだけ古い前提が残る
            # (p_process_flow で実際に起きた)
            text_w = (w - 0.4) if head else (w - 0.48 - BULLET_INDENT_IN)
            out.append((HEAD_LINE if head else 0.0)
                       + _text_lines(body, text_w, TS["body"], 400) * LINE_H + 0.06)
        return out
    side_blocks = {s: _blocks(spec.get(s, {}), w_probe) for s in ("left", "right")}
    # 対比プリミティブの item は左右で行単位に比較される。item 数が同じ場合は各行の
    # 高さを max(左, 右) に揃えて両カラムを同じ y から描く — 揃えないと行の対応が
    # 崩れ、片側だけ下へずれた「配置ズレ」に見える。数が違う場合のみ独立に流す。
    # カードの下側内側余白は上側(CONTENT_TOP - HEAD_H)と対称にする。内容実測高の
    # 末尾 ITEM_GAP を除いた上で同値の余白を積む(下側だけ痩せた非対称は欠陥)
    BOT_PAD = CONTENT_TOP - HEAD_H
    aligned = bool(side_blocks["left"]) and len(side_blocks["left"]) == len(side_blocks["right"])
    if aligned:
        rows_h = [max(a, b) for a, b in zip(side_blocks["left"], side_blocks["right"])]
        # ブロックの行送り(ピッチ)は可能な限り均一にする。行ごとの max 高のままだと、
        # 折返し見積りの境界揺れがそのまま「1-2 個目だけ間隔が広い」不揃いに見える。
        # 均一ピッチ(全ブロック中の最大高)が領域に収まるときは常に優先し、
        # 収まらないときのみ行別高へフォールバックする。
        pitch = max(rows_h)
        n_rows = len(rows_h)
        if CONTENT_TOP + pitch * n_rows + ITEM_GAP * (n_rows - 1) + BOT_PAD <= h - 0.1:
            rows_h = [pitch] * n_rows
        content_h = sum(rows_h) + n_rows * ITEM_GAP
    else:
        content_h = max((sum(b) + len(b) * ITEM_GAP) for b in side_blocks.values()) or 0.4
    card_h = min(h - 0.1, max(1.9, CONTENT_TOP + content_h - ITEM_GAP + BOT_PAD))
    y_top = y0 + 0.05 + max(0.0, (h - 0.1 - card_h) * 0.44)
    for side, (c0, span) in (("left", (0, 6)), ("right", (6, 6))):
        blk = spec.get(side, {})
        x, w = grid(c0, span)
        add_rect(slide, x, y_top, w, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        mark = blk.get("mark")
        if mark == "cross":
            hd_fill = C["chart_gray"]
        elif mark == "check" or blk.get("focal"):
            hd_fill = C["primary_deep"]
        else:
            hd_fill = C["navy"]
        add_rect(slide, x, y_top, w, HEAD_H, hd_fill, radius_pt=LAY["card"]["radius_pt"])
        add_rect(slide, x, y_top + HEAD_H - 0.19, w, 0.19, hd_fill)
        head_txt = blk.get("heading", side)
        if mark == "cross":
            head_txt = "× " + head_txt
        elif mark == "check":
            head_txt = "○ " + head_txt
        add_text(slide, x + 0.16, y_top, w - 0.32, HEAD_H,
                 [[(head_txt, TS["body_small"], 600, C["canvas"])]], anchor=MSO_ANCHOR.MIDDLE)
        # 領域キャップでカードが実測内容より低い場合は、下側余白(上側と対称)を犠牲に
        # せず item 間ギャップを床値(0.10)まで圧縮して収める。床値でも収まらない場合は
        # コピー量の問題 — スペック側で本文を削る(タイポグラフィは縮小しない)。
        # aligned 時は行送り(rows)を左右共通にし、各 item を同じ y から描く。
        items = blk.get("items", [])
        blocks = side_blocks[side]
        rows = rows_h if aligned else blocks
        gap_eff = ITEM_GAP
        avail = card_h - CONTENT_TOP - BOT_PAD
        if len(rows) > 1 and sum(rows) + (len(rows) - 1) * ITEM_GAP > avail:
            gap_eff = max(0.10, (avail - sum(rows)) / (len(rows) - 1))
        yy = y_top + CONTENT_TOP
        for it, block_h, row_h in zip(items, blocks, rows):
            head = it.get("heading") if isinstance(it, dict) else None
            body = it.get("body") if isinstance(it, dict) else str(it)
            if head:
                add_text(slide, x + 0.2, yy, w - 0.4, block_h,
                         [[(head, TS["body_small"], 600, C["ink"])], [(body, TS["body"], 400, C["ink_subtle"])]],
                         line_spacing=1.28, space_after_pt=3)
            else:
                add_bullets(slide, x + 0.24, yy, w - 0.48, block_h,
                            [body], TS["body"], C["ink_subtle"], line_spacing=1.28, space_after_pt=0)
            yy += row_h + gap_eff


def p_process_flow(slide, spec, deck):
    """Horizontal step flow with descriptions under each step. Card height follows
    the fullest step's measured content; the whole block centers in the body region.
    箇条書きは ●+ぶら下げインデント+段落後スペース、矢羽の矢先はカード右端から
    少しはみ出させて前進感を作る(参照モック準拠)。"""
    steps = spec.get("steps", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(steps))
    sw = w / n
    gut = LAY["gutter_in"]
    overhang = gut  # 矢先をカード右端より前へ出す(次の矢羽の左辺には触れるだけ)
    # 矢羽の見出しは折返す(狭い列に長いラベルが入ると2行以上になる)。高さを直値で固定すると、
    # 折返したぶんが箱からあふれてカードへ食い込む — ラベルの実測行数から高さを決める。
    # ただし矢先(tip)は見出し高に比例するので、高さが伸びるとテキスト幅は逆に狭くなる。
    # 最小高の tip で測ると「描く箱より広い幅」で行数を数えることになり、実際には1行増えて
    # 折返し予算を外す(=文節で割れず自然折返しへ落ちる)。高さと幅を収束させてから決める
    # 高さが伸びると矢先も伸びてテキスト幅は狭くなる = 行数がまた増える、という正のループ。
    # 上限を置かないと発散する(幅が負になり、スライドの外へ出る巨大な矢羽が生まれる)。
    # 矢羽の見出しは3行までとし、それでも入らないなら組版ではなくコピーの問題として
    # 自然折返し + verify の警告に委ねる
    HEAD_PT, HEAD_MIN_H, TIP_RATIO, HEAD_MAX_ROWS, HEAD_MIN_W = 15, 0.52, 0.6, 3, 0.6
    head_line = (HEAD_PT / 72.0) * 1.30 * 1.22
    head_h, head_rows = HEAD_MIN_H, 1
    for _ in range(4):                                # 高さは単調に増えるだけなので数回で収束
        head_w = max(HEAD_MIN_W, sw - gut + overhang - head_h * TIP_RATIO - 0.12)
        head_rows = min(HEAD_MAX_ROWS,
                        max((_text_lines(st.get("label", ""), head_w, HEAD_PT, 600)
                             for st in steps), default=1))
        grown = max(HEAD_MIN_H, head_rows * head_line + 0.20)
        if grown <= head_h + 1e-6:
            break
        head_h = grown
    line_h = (TS["body"] / 72.0) * 1.30 * 1.22  # 1.22 = 和文フォントの行ボックス補正(実測)
    sp_h = 8 / 72.0
    card_pad = 0.56  # カード内左右の余白合計に相当する差し引き(gutter 別)
    # 箇条書きはぶら下げインデント(BULLET_INDENT_IN)の分だけ折返し幅が狭い。行数見積もり
    # は描画箱幅からインデントと安全余白を引いた実効幅で行う(引き忘れると card_h が足りず、
    # 折り返した項目が outcome の罫線に食い込む)。字下げ幅を直値で持たないこと — 記号や
    # インデントを変えたときに、ここだけ古い前提のまま残る
    bullet_text_w = sw - LAY["gutter_in"] - card_pad - BULLET_INDENT_IN - 0.20
    outcome_text_w = sw - LAY["gutter_in"] - card_pad
    outcome_pt = TS["section_heading"]
    content_h = 0.0
    has_outcome = any(st.get("outcome") for st in steps)
    for st in steps:
        used = 0.34 if st.get("desc") else 0.0
        for b in st.get("items", []):
            used += _text_lines(b, bullet_text_w, TS["body"], 400) * line_h + sp_h
        content_h = max(content_h, used)
    outcome_h = 0.0
    if has_outcome:
        # +0.46 = 罫線とテキストの間のゾーン(罫→▼→テキスト 0.26)+箇条書きとの分離 0.20。
        # outcome のテキスト高は全カードで共有し、罫・▼・結論の y を横一列に揃える
        outcome_h = max(_text_lines(st.get("outcome", ""), outcome_text_w, outcome_pt, 700)
                        for st in steps) * (outcome_pt / 72.0) * 1.20 * 1.22 + 0.46
    # カードの上下内側余白は対称(0.26/0.26)。下側だけ痩せると outcome が下縁に近接する
    card_h = min(h - head_h - 0.35, max(1.5, 0.26 + content_h + outcome_h + 0.26))
    # フィルルール: 帯(シェブロン+カード)で本文領域を ~82% まで埋め、上下余白を詰める
    # (中身はカード内で垂直中央に置く)
    card_h = max(card_h, min(h - head_h - 0.35, 0.82 * h - head_h - 0.20))
    y_top = y0 + 0.05 + max(0.05, (h - head_h - 0.25 - card_h - 0.1) * 0.42)
    # 彩度は主役ステップに配給する。デフォルトは到達点(最終ステップ)だが、タイトルの
    # 主張を担うステップが途中にある場合は focal_step で指定する — タイトルの核心と
    # 強調色の位置がずれると、主張と視覚焦点の不一致になる
    try:
        focal_i = int(spec.get("focal_step", len(steps) - 1))
    except (TypeError, ValueError):
        focal_i = len(steps) - 1  # 非数値は既定(最終ステップ)へフォールバック — ビルドを落とさない
    focal_i = min(max(focal_i, 0), max(0, len(steps) - 1))  # 整数の範囲外はクランプ
    for i, st in enumerate(steps):
        sx = x + i * sw
        head_fill = C["primary"] if i == focal_i else C["primary_pale"]
        label_color = C["canvas"] if i == focal_i else C["primary_deep"]
        add_chevron(slide, sx, y_top, sw - gut + overhang, head_h, head_fill)
        tip = head_h * 0.6  # PENTAGON: 左辺フラット、右にだけ矢先
        add_text(slide, sx + 0.12, y_top, sw - gut + overhang - tip - 0.12, head_h,
                 [[(st.get("label", ""), HEAD_PT, 600, label_color)]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False)
        cy = y_top + head_h + 0.15
        add_rect(slide, sx, cy, sw - gut, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        # このカード自身の内容高を実測し、items+outcome を1つの意味グループとして
        # 垂直中央に置く。outcome をカード下端へ固定すると、短い箇条書きとの間に
        # 大きな空白帯が生まれる。
        own_h = (0.34 if st.get("desc") else 0.0)
        for b in st.get("items", []):
            own_h += _text_lines(b, bullet_text_w, TS["body"], 400) * line_h + sp_h
        # 箇条書きは上端固定(上インセット0.26)、結論ストリップ(罫+▼+outcome)はカード
        # 下端から対称インセットで bottom edge-lock する。行高見積りの誤差はすべて
        # 中間の余白に吸収され、結論が下縁に沈む/浮くことがない。outcome テキスト高は
        # 全カード共有なので、罫・▼・結論の y は横一列に揃う。
        iy = cy + 0.26
        if st.get("desc"):
            add_text(slide, sx + 0.28, iy, sw - gut - 0.56, 0.28,
                     [[(st["desc"], 14.5, 600, C["ink"])]])
            iy += 0.34
        if st.get("outcome"):
            out_text_h = outcome_h - 0.46
            ty = cy + card_h - 0.26 - out_text_h
            bullet_h = max(0.24, min(own_h, ty - 0.46 - iy))
            add_bullets(slide, sx + 0.30, iy, sw - gut - card_pad, bullet_h,
                        st.get("items", []), TS["body"], C["ink_subtle"],
                        line_spacing=1.30, space_after_pt=8)
            add_line(slide, sx + 0.28, ty - 0.26, sx + sw - gut - 0.28, ty - 0.26, C["rule"], 0.5)
            add_text(slide, sx + 0.28, ty - 0.23, sw - gut - 0.56, 0.18,
                     [[("▼", 9, 400, C["primary"])]], align=PP_ALIGN.CENTER)
            add_text(slide, sx + 0.28, ty, sw - gut - 0.56, out_text_h,
                     [[(st["outcome"], outcome_pt, 700, C["ink"])]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)
        else:
            add_bullets(slide, sx + 0.30, iy, sw - gut - card_pad, max(0.24, own_h),
                        st.get("items", []), TS["body"], C["ink_subtle"],
                        line_spacing=1.30, space_after_pt=8)


def p_quote_or_statement(slide, spec, deck):
    """Flexible closing statement. center_hero is ceremonial; split_evidence is an
    editorial close with proof metrics attached."""
    y0, h = body_region(spec)
    recap = spec.get("recap", [])
    stmt = spec.get("statement", "")
    variant = spec.get("variant", spec.get("layout", "center_hero"))
    if variant in ("evidence_strip", "thesis_strip", "closing_grid") and recap:
        x, w = grid(0, 12)
        lead, support = split_message(spec)
        s_pt = 31 if _ja_len(lead) <= 40 else 28
        sup_pt = TS["subtitle"]
        lead_lines, lead_w = shape_message(lead, w, s_pt, 700, max_lines=2)
        # 支える一文は「文章」— 行を埋め、語だけ割らずに流す。行長は見出しに寄り添う幅に取り、
        # 右に余白を残す(占有バランス)
        sup_w = min(w, max(lead_w, w * MESSAGE_SUPPORT_MEASURE))
        sup_lines = (wrap_prose(support, sup_w, sup_pt, 400).split("\n") if support else [])
        lead_h = len(lead_lines) * (s_pt / 72.0) * leading(s_pt, "statement") * 1.15 + 0.10
        sup_h = (len(sup_lines) * (sup_pt / 72.0) * leading(sup_pt) * 1.22 + 0.08) if sup_lines else 0.0
        SUP_GAP = 0.18 if sup_lines else 0.0
        stmt = "\n".join(lead_lines)
        text_h = lead_h + SUP_GAP + sup_h
        text_w = max(lead_w, sup_w)
        n = min(4, len(recap))
        gap = LAY["gutter_in"]
        heading_h = 0.34
        card_h = 1.42 if n <= 3 else 1.30
        group_h = 0.06 + 0.36 + text_h + 0.40 + heading_h + 0.16 + card_h
        gy = y0 + max(0.0, (h - group_h) * 0.44)
        add_rect(slide, x, gy, 0.72, 0.055, C["primary"])
        # 整形済み(行が決まっている)テキストは、もう一度折り返さない
        add_text(slide, x, gy + 0.34, lead_w, lead_h,
                 [[(stmt, s_pt, 700, C["ink"])]], role="statement", display_wrap=False)
        if sup_lines:
            add_text(slide, x, gy + 0.34 + lead_h + SUP_GAP, sup_w, sup_h,
                     [[("\n".join(sup_lines), sup_pt, 400, C["ink_subtle"])]],
                     display_wrap=False)
        hy = gy + 0.34 + text_h + 0.40
        add_text(slide, x, hy, w, heading_h,
                 [[(spec.get("recap_heading", "確認指標"), TS["section_heading"], 700, C["primary_deep"])]])
        add_line(slide, x, hy + heading_h + 0.05, x + w, hy + heading_h + 0.05, C["rule"], 0.75)
        card_y = hy + heading_h + 0.20
        cw = (w - gap * (n - 1)) / n
        any_focal = any(m.get("focal") for m in recap[:n])
        for i, m in enumerate(recap[:n]):
            cx = x + i * (cw + gap)
            focal = m.get("focal") or not any_focal
            add_rect(slide, cx, card_y, cw, card_h,
                     C["primary_pale"] if focal else C["surface_tint"],
                     radius_pt=LAY["card"]["radius_pt"])
            add_text(slide, cx + 0.22, card_y + 0.18, cw - 0.44, 0.28,
                     [[(m.get("label", ""), TS["body"], 600, C["ink_subtle"])]])
            vcolor = C["primary_deep"] if focal else C["ink"]
            vparts = [(str(m.get("value", "")), 34, 700, vcolor)]
            if m.get("unit"):
                vparts.append(_unit_part(m["unit"], 16))
            add_text(slide, cx + 0.22, card_y + 0.52, cw - 0.44, 0.58, [vparts],
                     anchor=MSO_ANCHOR.MIDDLE)
            if m.get("note"):
                add_text(slide, cx + 0.22, card_y + 1.08, cw - 0.44, 0.24,
                         [[(m["note"], TS["kpi_sub"], 400, C["ink_faint"])]])
        return
    if variant in ("split_evidence", "editorial_split") and recap:
        lx, lw = grid(0, 7)
        rx, rw = grid(8, 4)
        s_pt = TS["statement"] if _ja_len(stmt) <= 60 else max(30, TS["statement"] - 6)
        lines = _text_lines(stmt, lw, s_pt, 700)
        # Split-evidence closers use a large editorial statement; mixed JP/ASCII numerals
        # often wrap one line earlier in the rendered font than the simple JA estimate.
        text_h = (lines + 0.35) * (s_pt / 72.0) * 1.52 * 1.15 + 0.16
        attr_h = 0.40 if spec.get("attribution") else 0.0
        group_h = 0.36 + text_h + attr_h
        rail_h = min(h - 0.30, max(2.75, 0.55 + len(recap) * 0.72 + 0.28))
        gy = y0 + max(0.0, (h - max(group_h, rail_h)) * 0.45)
        add_rect(slide, lx, gy, 0.82, 0.055, C["primary"])
        add_text(slide, lx, gy + 0.32, lw, text_h, [[(stmt, s_pt, 700, C["ink"])]],
                 line_spacing=1.34, align=PP_ALIGN.LEFT)
        if spec.get("attribution"):
            add_text(slide, lx, gy + 0.32 + text_h + 0.10, lw, 0.32,
                     [[(spec["attribution"], TS["body_small"], 400, C["ink_faint"])]])
        ry = gy
        add_rect(slide, rx, ry, rw, rail_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        add_text(slide, rx + 0.25, ry + 0.24, rw - 0.5, 0.30,
                 [[(spec.get("recap_heading", "結論を支える指標"), TS["section_heading"], 700, C["primary_deep"])]],
                 wrap=False)
        cy = ry + 0.72
        any_focal = any(m.get("focal") for m in recap)
        for i, m in enumerate(recap):
            if i > 0:
                add_line(slide, rx + 0.25, cy - 0.10, rx + rw - 0.25, cy - 0.10, C["rule"], 0.75)
            label_w = rw * 0.34
            add_text(slide, rx + 0.25, cy, label_w, 0.36,
                     [[(m.get("label", ""), TS["body"], 600, C["ink_subtle"])]],
                     anchor=MSO_ANCHOR.MIDDLE)
            vcolor = C["primary_deep"] if (m.get("focal") or not any_focal) else C["ink"]
            vparts = [(str(m.get("value", "")), 34, 700, vcolor)]
            if m.get("unit"):
                vparts.append(_unit_part(m["unit"], 16))
            value_x = rx + 0.25 + label_w + 0.10
            add_text(slide, value_x, cy - 0.08, rx + rw - 0.25 - value_x, 0.64, [vparts],
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            cy += 0.72
        return

    x, w = grid(1, 10)
    eff_w = w * 0.94  # 中央揃えの実効幅
    s_pt = TS["statement"] if _ja_len(stmt) <= 60 else max(30, TS["statement"] - 6)
    stmt_lines = _statement_lines(stmt, eff_w, s_pt)  # 、で節分割し末尾の孤立行を防ぐ
    # 節が実効幅を超えて折返す間はフォントを1ptずつ縮小(28pt下限)。節自体が長すぎて
    # 28ptでも1行に収まらない場合は validate_spec が短縮を警告する(縮小より短文化が原則)。
    while s_pt > 28 and any(_text_lines(ln, eff_w, s_pt, 700) > 1 for ln in stmt_lines):
        s_pt -= 1
        stmt_lines = _statement_lines(stmt, eff_w, s_pt)
    lines = sum(max(1, _text_lines(ln, eff_w, s_pt, 700)) for ln in stmt_lines)  # 実効幅で行数見積
    text_h = lines * (s_pt / 72.0) * 1.5 * 1.15 + 0.06  # 1.15 = 和文行ボックス補正
    rule_h = 0.30
    attr_h = 0.44 if spec.get("attribution") else 0.0
    recap_h = 1.37 if recap else 0.0
    group_h = rule_h + text_h + attr_h + recap_h
    gy = y0 + max(0.0, (h - group_h) * 0.45)
    add_rect(slide, x + w / 2 - 0.35, gy, 0.7, 0.05, C["primary"])
    ty = gy + rule_h
    add_text(slide, x, ty, w, text_h, [[("\n".join(stmt_lines), s_pt, 600, C["ink"])]],
             line_spacing=1.5, align=PP_ALIGN.CENTER)
    ay = ty + text_h + 0.14
    if spec.get("attribution"):
        add_text(slide, x, ay, w, 0.3, [[(spec["attribution"], 13, 400, C["ink_faint"])]],
                 align=PP_ALIGN.CENTER)
        ay += attr_h
    if recap:
        ry = ay + 0.30
        n = len(recap)
        rw_total = min(w, 3.1 * n)
        rx0 = x + (w - rw_total) / 2
        rw = rw_total / n
        for i, m in enumerate(recap):
            rx = rx0 + i * rw
            if i > 0:
                add_line(slide, rx, ry + 0.04, rx, ry + 0.92, C["rule"], 0.75)
            add_text(slide, rx, ry, rw, 0.28, [[(m.get("label", ""), 14, 600, C["ink_faint"])]],
                     align=PP_ALIGN.CENTER)
            # 主役は1つに絞る: focal だけ primary_deep、残りは ink(本文色)で支える(Qwen)
            any_focal = any(mm.get("focal") for mm in recap)
            vcolor = C["primary_deep"] if (m.get("focal") or not any_focal) else C["ink"]
            vparts = [(str(m.get("value", "")), 34, 700, vcolor)]
            if m.get("unit"):
                vparts.append(_unit_part(m["unit"], 17))
            add_text(slide, rx, ry + 0.36, rw, 0.55, [vparts], align=PP_ALIGN.CENTER)




def _unit_part(unit: str, size: float, weight: int = 600):
    pad = "" if unit in ("%", "pt", "x", "倍") else " "
    return (pad + unit, max(size, TS["kpi_sub"]), weight, C["ink_subtle"])


def _metric_pair(slide, x, y, w, m, value_size=24):
    """label (small) / big value + small unit / colored delta — exemplar big-number grammar."""
    add_text(slide, x, y, w, 0.26, [[(m.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])]])
    vparts = [(str(m.get("value", "")), value_size, 700, C["primary_deep"])]
    if m.get("unit"):
        vparts.append(_unit_part(m["unit"], value_size * 0.45))
    add_text(slide, x, y + 0.26, w, value_size / 60, [vparts])
    dy = y + 0.26 + value_size / 60 + VALUE_DELTA_GAP_IN
    if m.get("delta"):
        add_text(slide, x, dy, w, 0.26,
                 [[(m["delta"], TS["kpi_sub"], 600, _delta_color(m.get("delta_dir"), m.get("positive_is_good", True)))]])
        dy += 0.26
    return dy


def p_financial_highlights(slide, spec, deck):
    """Earnings highlight board: hero KPIs up top, supporting metrics in an evidence strip."""
    groups = spec.get("groups", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    hero_items = []
    support_items = []
    notes = []
    for g in groups:
        metrics = g.get("metrics", [])
        if not metrics:
            continue
        hero = next((m for m in metrics if m.get("hero") or m.get("focal")), metrics[0])
        hero_items.append((g, hero))
        for m in metrics:
            if m is not hero:
                support_items.append((g, m))
        if g.get("note"):
            notes.append(g["note"])
    if not hero_items and not support_items:
        return

    # ヒーローカードは3枚まで。4グループ目以降の主指標は黙って落とさず補助指標
    # ストリップの先頭へ回す。support_h / heading / group_h の算出より前に回すこと —
    # 後から足すと「溢れ分で初めて support が生まれる」ケースで高さ0のストリップに
    # 描画されレイアウトが潰れる。validate_spec も groups > 3 を警告する
    hero_n = min(3, len(hero_items))
    if len(hero_items) > hero_n:
        support_items = list(hero_items[hero_n:]) + support_items
        hero_items = hero_items[:hero_n]

    gap = LAY["gutter_in"]
    hero_h = min(2.48, max(2.28, h * 0.44))
    support_h = min(1.66, max(1.50, h * 0.30)) if support_items else 0.0
    support_heading_h = 0.34 if support_items else 0.0
    note_h = 0.26 if notes else 0.0
    group_h = hero_h + (0.40 + support_heading_h + 0.16 + support_h if support_items else 0) + note_h
    y = y0 + max(0.03, (h - group_h) * 0.42)

    hero_w = (w - gap * (hero_n - 1)) / hero_n
    any_focal = any(m.get("focal") for _, m in hero_items[:hero_n])
    for i, (g, m) in enumerate(hero_items[:hero_n]):
        hx = x + i * (hero_w + gap)
        focal = m.get("focal") or (i == 0 and not any_focal)
        fill = C["primary_pale"] if focal else C["surface_tint"]
        add_rect(slide, hx, y, hero_w, hero_h, fill, radius_pt=LAY["card"]["radius_pt"])
        add_rect(slide, hx + 0.24, y + 0.22, 0.10, 0.10, C["primary"])
        add_text(slide, hx + 0.40, y + 0.15, hero_w - 0.64, 0.30,
                 [[(g.get("label", ""), TS["kpi_label"], 700, C["ink"])]])
        if g.get("claim"):
            add_text(slide, hx + 0.24, y + 0.54, hero_w - 0.48, 0.48,
                     [[(g["claim"], TS["section_heading"], 700, C["ink"])]], line_spacing=1.13)
        vcolor = C["primary_deep"] if focal else C["ink"]
        vparts = [(str(m.get("value", "")), TS["kpi_value_hero"], 700, vcolor)]
        if m.get("unit"):
            vparts.append(_unit_part(m["unit"], 17))
        value_y = y + 1.02
        value_h = 0.74
        add_text(slide, hx + 0.24, value_y, hero_w - 0.48, value_h, [vparts],
                 anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        metric_line = m.get("label", "")
        if m.get("delta"):
            metric_line = f"{metric_line}  {m['delta']}"
        metric_y = value_y + value_h + VALUE_DELTA_GAP_IN
        add_text(slide, hx + 0.24, metric_y, hero_w - 0.48, 0.28,
                 [[(metric_line, TS["body"], 700, _delta_color(m.get("delta_dir"), m.get("positive_is_good", True)))]],
                 wrap=False)

    if support_items:
        sy = y + hero_h + 0.40
        add_text(slide, x, sy, w, support_heading_h,
                 [[(spec.get("support_heading", "補助指標"), TS["section_heading"], 700, C["primary_deep"])]])
        add_line(slide, x, sy + support_heading_h + 0.04, x + w, sy + support_heading_h + 0.04, C["rule"], 0.75)
        card_y = sy + support_heading_h + 0.16
        n = min(4, len(support_items))
        cw = (w - gap * (n - 1)) / n
        for i, (g, m) in enumerate(support_items[:n]):
            cx = x + i * (cw + gap)
            add_rect(slide, cx, card_y, cw, support_h, C["canvas"], line=C["rule"],
                     radius_pt=LAY["card"]["radius_pt"])
            add_text(slide, cx + 0.18, card_y + 0.16, cw - 0.36, 0.25,
                     [[(m.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])]], wrap=False)
            vparts = [(str(m.get("value", "")), 34, 700, C["primary_deep"])]
            if m.get("unit"):
                vparts.append(_unit_part(m["unit"], 15))
            value_y = card_y + 0.47
            value_h = 0.50
            add_text(slide, cx + 0.18, value_y, cw - 0.36, value_h, [vparts],
                     anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            if m.get("delta"):
                delta_y = value_y + value_h + VALUE_DELTA_GAP_IN
                add_text(slide, cx + 0.18, delta_y, cw - 0.36, 0.28,
                         [[(m["delta"], TS["kpi_sub"], 700, _delta_color(m.get("delta_dir"), m.get("positive_is_good", True)))]],
                         wrap=False)
        if len(support_items) > n:
            more = " / ".join(f"{m.get('label','')}: {m.get('value','')}{m.get('unit','')}" for _, m in support_items[n:])
            add_text(slide, x, card_y + support_h + 0.08, w, 0.22,
                     [[(more, TS["body_small"], 400, C["ink_faint"])]], wrap=False)

    if notes:
        add_text(slide, x, y + group_h - note_h + 0.04, w, note_h,
                 [[(" / ".join(notes), TS["body_small"], 400, C["ink_subtle"])]], wrap=False)


def p_metrics_rows(slide, spec, deck):
    """Exemplar archetype: hairline rows of label | big value+unit | YoY (1-2 column groups)."""
    columns = spec.get("columns", [])
    y0, h = body_region(spec)
    n = max(1, len(columns))
    gut = 0.6
    gw = (CONTENT_W - gut * (n - 1)) / n
    max_rows = max((len(g.get("rows", [])) for g in columns), default=1)
    has_heading = any(g.get("heading") for g in columns)
    row_h_all = min(0.95, (h - 0.13 - (0.36 if has_heading else 0.0)) / max(1, max_rows))
    block_h = (0.36 if has_heading else 0.0) + max_rows * row_h_all
    y_shift = max(0.0, (h - 0.08 - block_h) * 0.44)
    for gi, g in enumerate(columns):
        gx = MX + gi * (gw + gut)
        gy = y0 + 0.08 + y_shift
        if g.get("heading"):
            add_text(slide, gx, gy, gw, 0.30, [[(g["heading"], TS["body"], 600, C["ink"])]], align=PP_ALIGN.CENTER)
            gy += 0.38
        rows = g.get("rows", [])
        row_h = row_h_all
        for r in rows:
            emph = r.get("emphasis")
            if emph:
                add_rect(slide, gx - 0.08, gy + 0.02, gw + 0.16, row_h - 0.06, C["primary_pale"], radius_pt=4)
            add_text(slide, gx, gy, gw * 0.34, row_h,
                     [[(r.get("label", ""), TS["kpi_label"], 600, C["ink"] if emph else C["ink_subtle"])]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
            vparts = [(str(r.get("value", "")), 25, 700, C["primary_deep"] if emph else C["ink"])]
            if r.get("unit"):
                vparts.append(_unit_part(r["unit"], TS["kpi_sub"], 400))
            add_text(slide, gx + gw * 0.34, gy, gw * 0.38, row_h, [vparts],
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            if r.get("delta"):
                add_text(slide, gx + gw * 0.76, gy, gw * 0.24, row_h,
                         [[(r["delta"], TS["kpi_sub"], 600, _delta_color(r.get("delta_dir", "up"), r.get("positive_is_good", True)))]],
                         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            add_line(slide, gx, gy + row_h - 0.02, gx + gw, gy + row_h - 0.02, C["rule"], 0.5)
            gy += row_h


def p_driver_decomposition(slide, spec, deck):
    """Exemplar archetype: KPI factor cards joined by operators (数 × 単価 = GMV)."""
    factors = spec.get("factors", [])
    y0, h = body_region(spec)
    n = max(1, len(factors))
    op_w = 0.55
    gut = 0.18
    cw = (CONTENT_W - op_w * (n - 1) - gut * 2 * (n - 1)) / n
    has_delta = any(f.get("delta") for f in factors)
    factor_note_pt = TS["kpi_sub"]
    note_line_h = (factor_note_pt / 72.0) * 1.22 * 1.13
    note_w = max(0.4, cw - 0.3)     # 因数が増えるほどカードは痩せる — 幅は正のまま床を持つ
    note_lines = max((_text_lines(f.get("note", ""), note_w, factor_note_pt, 400) for f in factors),
                     default=0)
    card_h = min(h - 0.75, max(1.7, 0.30 + 0.34 + 0.66 + (0.46 if has_delta else 0.0)
                               + (0.12 + note_lines * note_line_h if note_lines else 0.0) + 0.26))
    cy = y0 + 0.15 + max(0.0, (h - 0.55 - card_h) * 0.44)
    ops = spec.get("operators") or (["×"] * (n - 2) + ["="] if n >= 2 else [])
    for i, f in enumerate(factors):
        x = MX + i * (cw + op_w + gut * 2)
        emphasized = f.get("focal", i == n - 1)
        add_rect(slide, x, cy, cw, card_h, C["primary_pale"] if emphasized else None,
                 line=None if emphasized else C["rule"], line_w_pt=1.0, radius_pt=8)
        add_text(slide, x + 0.15, cy + 0.22, cw - 0.3, 0.32,
                 [[(f.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])]], align=PP_ALIGN.CENTER)
        vparts = [(str(f.get("value", "")), 34, 700, C["primary_deep"])]
        if f.get("unit"):
            vparts.append(_unit_part(f["unit"], 15))
        vy = cy + 0.58
        add_text(slide, x + 0.15, vy, cw - 0.3, 0.55, [vparts], align=PP_ALIGN.CENTER)
        if f.get("delta"):
            add_text(slide, x + 0.15, vy + 0.55 + VALUE_DELTA_GAP_IN, cw - 0.3, 0.30,
                     [[(f["delta"], TS["kpi_sub"], 600, _delta_color(f.get("delta_dir", "up"), f.get("positive_is_good", True)))]],
                     align=PP_ALIGN.CENTER)
        if f.get("note"):
            add_text(slide, x + 0.15, cy + card_h - 0.16 - note_lines * note_line_h, cw - 0.3, note_lines * note_line_h + 0.08,
                     [[(f["note"], factor_note_pt, 400, C["ink_faint"])]], align=PP_ALIGN.CENTER, line_spacing=1.18)
        if i < n - 1:
            add_text(slide, x + cw + gut, cy, op_w, card_h,
                     [[(ops[i] if i < len(ops) else "×", 26, 600, C["ink_subtle"])]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    caption = spec.get("formula_note")
    if caption is None:
        labels = [f.get("label", "").split("(")[0] for f in factors]
        parts = []
        for i, lb in enumerate(labels):
            parts.append(lb)
            if i < len(labels) - 1:
                parts.append(ops[i] if i < len(ops) else "×")
        caption = " ".join(parts)
    if caption:
        add_text(slide, MX, cy + card_h + 0.22, CONTENT_W, 0.26,
                 [[(caption, TS["body_small"], 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)




def p_guidance_progress(slide, spec, deck):
    """Guidance progress. With comparable bars, draw a comparison chart; without them,
    draw a current-position gauge so the engine never forces a single vertical bar."""
    bars = spec.get("bars", [])
    cur = spec.get("current", {})
    y0, h = body_region(spec)
    x, w = grid(0, 8)
    sx, sw = grid(8, 4)
    # guidance_progress is a proof-field + facts-rail composition. A normal grid gutter is
    # too tight for two different roles, so add a section gap while preserving the 8+4 role map.
    section_gap_extra = 0.16
    w -= section_gap_extra
    sx += section_gap_extra
    sw -= section_gap_extra
    chart_h = h - 0.75
    base_y = y0 + 0.30
    g_high = float(cur.get("guidance_high", cur.get("guidance_low", 0)) or 0)
    actual = float(cur.get("actual", 0))

    def add_guidance_side(y: float, side_h: float) -> None:
        side = spec.get("side", [])
        if not side:
            return
        sy = y + 0.04
        add_text(slide, sx, sy, sw, 0.32, [[(spec.get("side_heading", "進捗"), TS["section_heading"], 600, C["primary_deep"])]])
        add_line(slide, sx, sy + 0.34, sx + sw, sy + 0.34, C["rule"], 0.75)
        sy += 0.52
        step = min(0.68, max(0.54, (side_h - 0.62) / max(1, len(side))))
        for it in side:
            add_text(slide, sx, sy, sw * 0.55, 0.44, [[(it.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])]],
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, sx + sw * 0.55, sy, sw * 0.45, 0.44,
                     [[(str(it.get("value", "")), 20, 700, C["ink"])]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            sy += step

    if not bars:
        panel_h = min(h - 0.35, max(3.35, h * 0.72))
        py = y0 + max(0.05, (h - panel_h) * 0.42)
        add_rect(slide, x, py, w, panel_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        heading = spec.get("progress_heading") or (cur.get("label", "進捗") + " 現在地")
        add_text(slide, x + 0.34, py + 0.28, w - 0.68, 0.34,
                 [[(heading, TS["section_heading"], 700, C["primary_deep"])]],
                 wrap=False)
        progress = actual / g_high if g_high else 0.0
        progress = max(0.0, progress)
        progress_clamped = min(1.0, progress)
        pct_label = cur.get("progress_display") or f"{progress * 100:.0f}%"
        pct_pt = 58 if _ja_len(pct_label) <= 4 else 48
        add_text(slide, x + 0.34, py + 0.86, w * 0.42, 0.88,
                 [[(pct_label, pct_pt, 700, C["primary_deep"])]],
                 anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        supporting = cur.get("actual_display", "")
        if supporting:
            add_text(slide, x + 0.38, py + 1.70, w * 0.46, 0.30,
                     [[(supporting, TS["body"], 700, C["ink"])]], wrap=False)
        range_text = cur.get("range_display", "")
        if range_text:
            add_text(slide, x + w * 0.58, py + 1.02, w * 0.34, 0.34,
                     [[(range_text, TS["section_heading"], 700, C["ink"])]],
                     align=PP_ALIGN.RIGHT, wrap=False)
            add_text(slide, x + w * 0.58, py + 1.42, w * 0.34, 0.26,
                     [[("目標レンジ", TS["body_small"], 400, C["ink_subtle"])]],
                     align=PP_ALIGN.RIGHT, wrap=False)
        rail_x = x + 0.38
        rail_y = py + 2.42
        rail_w = w - 0.76
        rail_h = 0.34
        add_rect(slide, rail_x, rail_y, rail_w, rail_h, C["canvas"],
                 line=C["rule"], radius_pt=6)
        if progress_clamped > 0:
            add_rect(slide, rail_x, rail_y, rail_w * progress_clamped, rail_h,
                     C["primary"], radius_pt=6)
        if progress_clamped < 0.98:
            rem_x = rail_x + rail_w * progress_clamped
            add_rect(slide, rem_x, rail_y, rail_w * (1 - progress_clamped), rail_h, None,
                     line=C["primary_deep"], line_w_pt=1.0, dash="dash", radius_pt=6)
        add_text(slide, rail_x, rail_y + 0.42, rail_w * 0.25, 0.24,
                 [[("0%", TS["chart_axis"], 400, C["ink_subtle"])]], wrap=False)
        add_text(slide, rail_x + rail_w * 0.75, rail_y + 0.42, rail_w * 0.25, 0.24,
                 [[("100%", TS["chart_axis"], 600, C["ink"])]], align=PP_ALIGN.RIGHT, wrap=False)
        if progress_clamped > 0.06:
            marker_x = rail_x + rail_w * progress_clamped
            add_line(slide, marker_x, rail_y - 0.06, marker_x, rail_y + rail_h + 0.06,
                     C["primary_deep"], 1.0)
        add_guidance_side(py, panel_h)
        return

    vals = [float(b.get("value", 0)) for b in bars] + [g_high]
    top = max(vals) * 1.18 if vals and max(vals) > 0 else 1.0
    n = len(bars) + 1
    slot = w / n
    bar_w = slot * 0.55
    add_line(slide, x, base_y + chart_h, x + w, base_y + chart_h, C["chart_gray"], 0.75)
    for i, b in enumerate(bars):
        v = float(b.get("value", 0))
        bh = chart_h * v / top
        bx = x + i * slot + (slot - bar_w) / 2
        add_rect(slide, bx, base_y + chart_h - bh, bar_w, bh, C["chart_gray"])
        add_text(slide, bx - slot * 0.2, base_y + chart_h - bh - 0.26, bar_w + slot * 0.4, 0.22,
                 [[(b.get("display", f"{v:,.0f}"), TS["chart_label"], 400, C["ink_subtle"])]], align=PP_ALIGN.CENTER, wrap=False)
        add_text(slide, x + i * slot, base_y + chart_h + 0.08, slot, 0.3,
                 [[(b.get("label", ""), TS["chart_axis"], 400, C["ink_subtle"])]], align=PP_ALIGN.CENTER)
    bx = x + len(bars) * slot + (slot - bar_w) / 2
    a_h = chart_h * actual / top
    g_h = chart_h * g_high / top
    add_rect(slide, bx, base_y + chart_h - a_h, bar_w, a_h, C["primary"])
    add_text(slide, bx, base_y + chart_h - a_h + 0.06, bar_w, 0.24,
             [[(cur.get("actual_display", f"{actual:,.0f}"), TS["chart_label"], 600, C["canvas"])]], align=PP_ALIGN.CENTER)
    if g_h > a_h + 0.02:
        add_rect(slide, bx, base_y + chart_h - g_h, bar_w, g_h - a_h, None,
                 line=C["primary_deep"], line_w_pt=1.25, dash="dash")
    lab_h = TS["kpi_sub"] / 72.0 * 1.30          # ラベル1行ぶんの高さ(枠は文字ぶんだけ持つ)
    if cur.get("range_display"):
        range_x = max(x, bx - slot * 0.35)
        range_w = min(bar_w + slot * 0.7, x + w - range_x)
        add_text(slide, range_x, base_y + chart_h - g_h - 0.34, range_w, lab_h,
                 [[(cur["range_display"], TS["chart_label"], 700, C["ink"])]], align=PP_ALIGN.CENTER, wrap=False)
    g_low = float(cur.get("guidance_low", 0) or 0)
    if g_low and g_low != g_high:
        low_h = chart_h * g_low / top
        add_line(slide, bx, base_y + chart_h - low_h, bx + bar_w, base_y + chart_h - low_h, C["primary_deep"], 0.75)
        # 上限・下限はデータの位置に紐づくので、レンジが狭いと2つのラベルが重なる。
        # 枠が触れないところまで押し広げる(値の対応は行の高さぶんのズレでは崩れない)
        hi_y = base_y + chart_h - g_h - 0.07
        lo_y = max(base_y + chart_h - low_h + 0.04, hi_y + lab_h + 0.04)
        lab_x = bx + bar_w + 0.08
        hi_txt = "上限 " + f"{g_high:,.0f}".rstrip("0").rstrip(".")
        lo_txt = "下限 " + f"{g_low:,.0f}".rstrip("0").rstrip(".")
        # 枠は文字の幅ぶんだけ持つ。余った幅を枠に含めると、隣(右の要点列)の枠へ食い込む。
        # 幅を足りない値へ切り詰めると、こんどは折り返して2行になる — 実測幅がどちらも防ぐ
        for txt, ly in ((hi_txt, hi_y), (lo_txt, lo_y)):
            add_text(slide, lab_x, ly, _label_w(txt, TS["kpi_sub"], 600), lab_h,
                     [[(txt, TS["kpi_sub"], 600, C["ink_subtle"])]], wrap=False)
    add_text(slide, x + len(bars) * slot, base_y + chart_h + 0.08, slot, 0.3,
             [[(cur.get("label", ""), TS["chart_axis"], 600, C["ink"])]], align=PP_ALIGN.CENTER)
    add_unit_note(slide, x, y0 - 0.02, spec.get("unit"))
    add_guidance_side(base_y + 0.2, chart_h)


PATTERNS = {
    "cover": p_cover,
    "agenda": p_agenda,
    "section_divider": p_section_divider,
    "executive_summary": p_executive_summary,
    "kpi_dashboard": p_kpi_dashboard,
    "chart_insight": p_chart_insight,
    "market_sizing": p_market_sizing,
    "comparison_table": p_comparison_table,
    "competitive_landscape": p_competitive_landscape,
    "financial_summary": p_financial_summary,
    "waterfall": p_waterfall,
    "roadmap": p_roadmap,
    "two_column": p_two_column,
    "process_flow": p_process_flow,
    "statement": p_quote_or_statement,
    "financial_highlights": p_financial_highlights,
    "metrics_rows": p_metrics_rows,
    "driver_decomposition": p_driver_decomposition,
    "guidance_progress": p_guidance_progress,
    "diagram": p_diagram,
    "chart_grid": p_chart_grid,
}

NO_CHROME = {"cover", "section_divider"}


def _strip_visible_periods(obj):
    """スライド表示テキストは体言止め・句点なしで組む。文末の 。/． を安全網として
    取り除く(speaker_notes は口頭原稿なので対象外)。"""
    if isinstance(obj, dict):
        return {k: (v if k == "speaker_notes" else _strip_visible_periods(v))
                for k, v in obj.items()}
    if isinstance(obj, list):
        return [_strip_visible_periods(v) for v in obj]
    if isinstance(obj, str):
        s = obj.rstrip()
        while s[-1:] in ("。", "．"):
            s = s[:-1].rstrip()
        return s
    return obj


def set_template_background(prs, color: RGBColor) -> None:
    """地の色をスライドマスター(テンプレート)の背景として設定する。

    全面を覆う矩形は置かない — 図形で地を塗ると、編集時に本文の下で毎回つかんでしまう
    オブジェクトが1枚ずつ増えるだけである。背景はテンプレート側の属性なので、マスターに
    一度置けば全レイアウト・全スライドが継承し、スライド上にオブジェクトは現れない。
    cSld の子要素は順序が決まっている(bg → spTree → …)ため、bg は先頭に差し込む。"""
    for master in prs.slide_masters:
        cSld = master._element.find(f"{P_NS}cSld")
        for old in cSld.findall(f"{P_NS}bg"):
            cSld.remove(old)
        bg = etree.Element(f"{P_NS}bg")
        bgPr = etree.SubElement(bg, f"{P_NS}bgPr")
        fill = etree.SubElement(bgPr, f"{A_NS}solidFill")
        etree.SubElement(fill, f"{A_NS}srgbClr").set("val", str(color))
        etree.SubElement(bgPr, f"{A_NS}effectLst")
        cSld.insert(0, bg)


def build(spec_path: Path, out_path: Path) -> Path:
    global _ASSET_ROOT
    if not MEASURE_OK:
        # 実測フォントが無いと折返し位置と行数見積りが近似に落ち、同じ deck.json から
        # 別の .pptx が出る(--baseline 比較では原因不明の差分に見える)。黙って作らない
        print("WARN: NotoSansJP-{400,600,700}.ttf が見つからない — 折返しは近似で組む"
              "(実測フォントを入れた環境と改行位置が変わる)")
    _ASSET_ROOT = out_path.resolve().parent  # image assets + manifest cache beside the .pptx
    deck = _strip_visible_periods(json.loads(spec_path.read_text()))
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    set_template_background(prs, C["canvas"])
    blank = prs.slide_layouts[6]
    slides = deck["slides"]
    total = len(slides)
    deck["_chapters"] = [(s.get("number") or i + 1, s.get("title", ""))
                         for i, s in enumerate([sp for sp in slides if sp.get("pattern") == "section_divider"])]
    for idx, spec in enumerate(slides, start=1):
        pattern = spec.get("pattern")
        if pattern not in PATTERNS:
            raise SystemExit(f"slide {idx}: unknown pattern '{pattern}'. valid: {sorted(PATTERNS)}")
        slide = prs.slides.add_slide(blank)
        if pattern in NO_CHROME:
            PATTERNS[pattern](slide, spec, deck)
            if pattern == "section_divider":
                add_footer(slide, spec, idx, total, deck)
        else:
            add_chrome(slide, spec, idx, total, deck)
            PATTERNS[pattern](slide, spec, deck)
            add_insight_bar(slide, spec)
        if spec.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = spec["speaker_notes"]
    prs.save(out_path)
    return out_path


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", type=Path)
    ap.add_argument("-o", "--out", type=Path, default=None)
    args = ap.parse_args()
    out = args.out or args.spec.with_suffix(".pptx")
    build(args.spec, out)
    print(f"built: {out}")


if __name__ == "__main__":
    main()
