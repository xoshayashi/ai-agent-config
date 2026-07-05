#!/usr/bin/env python3
"""Build an Act-styled 16:9 banker-grade .pptx from a deck.json spec.

Usage: build_deck.py <deck.json> [-o out.pptx]

The LLM owns content (deck.json); this script owns geometry, typography,
and color so every deck lands on the same visual system deterministically.
"""
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

TOKENS = json.loads((Path(__file__).resolve().parent.parent / "references" / "tokens.json").read_text())

C = {k: RGBColor.from_string(v) for k, v in TOKENS["colors"].items()}
TS = TOKENS["type_scale_pt"]
LAY = TOKENS["layout"]
CH = TOKENS["chart_style"]
FONTS = TOKENS["fonts"]
SLIDE_W = Inches(TOKENS["slide"]["width_in"])
SLIDE_H = Inches(TOKENS["slide"]["height_in"])
MX = LAY["margin_x_in"]
CONTENT_W = TOKENS["slide"]["width_in"] - 2 * MX
FOOT_Y = LAY["footer"]["y_in"]
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ---------------------------------------------------------------- utilities

# 表示長(_ja_len)と全角→半角(hw)は deck_text.py の単一実装を使う — validate_spec の
# 字数バジェット・lint_render の readback 照合と同一実装であることが契約
from deck_text import hw, ja_len as _ja_len


def _text_lines(text: str, width_in: float, size_pt: float) -> int:
    """Estimated wrapped line count for JA-aware text in a box of width_in."""
    if not text:
        return 0
    chars_per_line = max(4.0, width_in / (size_pt / 72.0))
    return max(1, math.ceil(_ja_len(text) / chars_per_line))


def grid(col_start: int, col_span: int) -> tuple[float, float]:
    """Return (x_in, w_in) for a 12-column grid region (0-indexed start)."""
    ncols = LAY["grid_columns"]
    gut = LAY["gutter_in"]
    colw = (CONTENT_W - gut * (ncols - 1)) / ncols
    x = MX + col_start * (colw + gut)
    w = col_span * colw + (col_span - 1) * gut
    return x, w


def _set_run_fonts(run, size_pt: float, weight: int = 400, color: RGBColor | None = None) -> None:
    """Apply Act typography to a run: Geist + Noto Sans JP, weight via family/bold."""
    f = run.font
    f.size = Pt(size_pt)
    f.color.rgb = color if color is not None else C["ink"]
    if weight >= 700:
        f.name, ea = FONTS["latin"], FONTS["ea"]
        f.bold = True
    elif weight >= 600:
        f.name, ea = FONTS["latin_semibold"], FONTS["ea_semibold"]
        f.bold = False
    else:
        f.name, ea = FONTS["latin"], FONTS["ea"]
        f.bold = False
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find(f"{A_NS}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{A_NS}{tag}")
        el.set("typeface", ea)


def add_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing: float = 1.15, space_after_pt: float = 0.0, wrap=True,
             autosize_off=True):
    """Add a textbox. `runs` is a list of paragraphs; each paragraph is a list of
    (text, size_pt, weight, color) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    expanded = []
    for para in runs:
        if len(para) == 1 and "\n" in para[0][0]:
            text, size_pt, weight, color = para[0]
            expanded.extend([(ln, size_pt, weight, color)] for ln in text.split("\n"))
        else:
            expanded.append(para)
    runs = expanded
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after_pt:
            p.space_after = Pt(space_after_pt)
        for text, size_pt, weight, color in para:
            r = p.add_run()
            r.text = hw(text)
            _set_run_fonts(r, size_pt, weight, color)
    return tb


P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def _strip_style(sh):
    """Remove <p:style> so LibreOffice/PowerPoint don't apply theme effects (shadows)."""
    el = sh._element
    style = el.find(f"{P_NS}style")
    if style is not None:
        el.remove(style)


def add_rect(slide, x, y, w, h, fill: RGBColor | None, *, line: RGBColor | None = None,
             line_w_pt: float = 0.75, radius_pt: float = 0.0, dash: str | None = None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_pt > 0 else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius_pt > 0:
        frac = min(0.5, (radius_pt / 72.0) / max(0.01, min(w, h)))
        sh.adjustments[0] = frac
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w_pt)
        if dash:
            ln = sh.line._get_or_add_ln()
            d = ln.find(f"{A_NS}prstDash")
            if d is None:
                d = etree.SubElement(ln, f"{A_NS}prstDash")
            d.set("val", dash)
    sh.shadow.inherit = False
    _strip_style(sh)
    return sh


def add_line(slide, x1, y1, x2, y2, color: RGBColor, w_pt: float = 0.75, dash: str | None = None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w_pt)
    ln.shadow.inherit = False
    _strip_style(ln)
    if dash:
        el = ln.line._get_or_add_ln()
        d = etree.SubElement(el, f"{A_NS}prstDash")
        d.set("val", dash)
    return ln


def add_chevron(slide, x, y, w, h, fill: RGBColor):
    # 矢羽は常に左辺フラットの PENTAGON(ホームベース形)。左に窪みのある CHEVRON は
    # 「前工程から食い込まれる」形で、Act ではどの位置のステップにも使わない(全パターン共通)
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    _strip_style(sh)
    return sh


def _delta_color(direction: str | None, positive_is_good: bool = True) -> RGBColor:
    if direction == "up":
        return C["success"] if positive_is_good else C["danger"]
    if direction == "down":
        return C["danger"] if positive_is_good else C["success"]
    return C["ink_subtle"]


# ---------------------------------------------------------------- chrome

def title_lines(title: str) -> int:
    return 1 if _ja_len(title) <= TOKENS["text_budget"]["action_title_max_chars_ja"] else 2


def header_offset(spec: dict) -> float:
    return 0.26 if spec.get("kicker") else 0.0


def header_metrics(spec: dict) -> dict:
    """Deterministic header geometry: the accent bar spans exactly from the title
    cap-top to the bottom of the last header line (title or subtitle)."""
    hdr = LAY["header"]
    off = header_offset(spec)
    title = spec.get("title", "")
    lines = title_lines(title)
    # タイトルサイズは行数によらず一定。縮小して収める操作は AutoFit と同種の一貫性破壊
    # (スライド間でタイトルの見た目が揺れる)。2行予算 66 字(tokens.json の
    # action_title_two_line_max_chars_ja)は 25pt × 11.6in の2行に収まる
    t_size = TS["action_title"]
    title_y = hdr["title_y_in"] + off
    title_h = lines * (t_size / 72.0) * 1.16
    sub_y = title_y + title_h + hdr["title_subtitle_gap_in"]
    sub_h = (TS["subtitle"] / 72.0) * 1.25
    block_bottom = sub_y + sub_h if spec.get("subtitle") else title_y + title_h
    return {
        "off": off, "lines": lines, "t_size": t_size,
        "title_y": title_y, "title_h": title_h,
        "sub_y": sub_y, "sub_h": sub_h,
        "bar_y": title_y + 0.045,
        "bar_h": block_bottom - (title_y + 0.045),
        "body_y": block_bottom + hdr["body_gap_in"],
    }


def add_chrome(slide, spec: dict, page_no: int, total: int, deck: dict) -> None:
    """Background, header (optional kicker + accent bar + action title + subtitle), footer."""
    add_rect(slide, 0, 0, TOKENS["slide"]["width_in"], TOKENS["slide"]["height_in"], C["canvas"])
    hdr = LAY["header"]
    m = header_metrics(spec)
    if spec.get("kicker"):
        add_text(slide, hdr["title_x_in"], 0.18, hdr["title_w_in"], 0.26,
                 [[(spec["kicker"], 12, 600, C["ink_faint"])]])
    add_rect(slide, hdr["accent_bar_x_in"], m["bar_y"], hdr["accent_bar_w_in"], m["bar_h"], C["primary"])
    add_text(slide, hdr["title_x_in"], m["title_y"], hdr["title_w_in"], m["title_h"],
             [[(spec.get("title", ""), m["t_size"], 700, C["ink"])]], line_spacing=1.12)
    if spec.get("subtitle"):
        add_text(slide, hdr["title_x_in"], m["sub_y"], hdr["title_w_in"], m["sub_h"],
                 [[(spec["subtitle"], TS["subtitle"], 400, C["ink_subtle"])]])
    add_footer(slide, spec, page_no, total, deck)


def add_footer(slide, spec: dict, page_no: int, total: int, deck: dict) -> None:
    foot = LAY["footer"]
    frags = []
    if spec.get("source"):
        frags.append(("Source: " + spec["source"], TS["footnote"], 400, C["ink_faint"]))
    if spec.get("assumption"):
        pre = "   " if frags else ""
        frags.append((pre + "Assumption: " + spec["assumption"], TS["footnote"], 400, C["ink_faint"]))
    if spec.get("note"):
        pre = "   " if frags else ""
        frags.append((pre + "Note: " + spec["note"], TS["footnote"], 400, C["ink_faint"]))
    if frags:
        add_text(slide, foot["source_x_in"], foot["y_in"], CONTENT_W - 1.2, foot["h_in"], [frags])
    pn_w = 0.9
    add_text(slide, TOKENS["slide"]["width_in"] - foot["page_num_right_in"] - pn_w, foot["y_in"],
             pn_w, foot["h_in"], [[(str(page_no), TS["page_number"], 400, C["ink_faint"])]],
             align=PP_ALIGN.RIGHT)



def body_region(spec: dict) -> tuple[float, float]:
    """(y, h) available for body content, honoring kicker/insight if present."""
    y = header_metrics(spec)["body_y"]
    h = FOOT_Y - 0.12 - y
    if spec.get("insight"):
        h -= 0.74
    return y, h


def add_insight_bar(slide, spec: dict) -> None:
    """Optional bottom band. style 'accent' = selective judgment (pale yellow),
    style 'primary' = conclusion restating the so-what (pale petrol, centered)."""
    if not spec.get("insight"):
        return
    y = FOOT_Y - 0.74
    h = 0.58
    if spec.get("insight_style", "accent") == "primary":
        add_rect(slide, MX, y, CONTENT_W, h, C["primary_pale"], radius_pt=4)
        add_text(slide, MX + 0.18, y, CONTENT_W - 0.36, h,
                 [[(spec["insight"], TS["insight"], 600, C["primary_deep"])]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    else:
        add_rect(slide, MX, y, CONTENT_W, h, C["accent_pale"], line=C["accent_line"], line_w_pt=1.0, radius_pt=4)
        add_text(slide, MX + 0.18, y, CONTENT_W - 0.36, h,
                 [[(spec["insight"], TS["insight"], 600, C["ink"])]], anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- charts

CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
}


def _chart_ea_fonts(chart) -> None:
    """Ensure every defRPr in the chart XML carries latin+ea Act fonts."""
    root = chart._chartSpace
    ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
          "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for defRPr in root.findall(".//a:defRPr", ns):
        for tag in ("latin", "ea", "cs"):
            el = defRPr.find(f"a:{tag}", ns)
            if el is None:
                el = etree.SubElement(defRPr, f"{A_NS}{tag}")
            el.set("typeface", FONTS["latin"] if tag == "latin" else FONTS["ea"])


def add_act_chart(slide, x, y, w, h, cspec: dict):
    """Native chart with Act styling. cspec:
    {type, categories, series:[{name, values, focal|color}], unit, value_labels,
     focal_category (index highlighted for single-series), number_format}
    """
    ctype = cspec.get("type", "column")
    data = CategoryChartData()
    data.categories = [hw(str(c)) for c in cspec["categories"]]
    for s in cspec["series"]:
        data.add_series(hw(s["name"]), tuple(float(v) for v in s["values"]))
    gf = slide.shapes.add_chart(CHART_TYPES[ctype], Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = gf.chart
    chart.has_title = False
    nseries = len(cspec["series"])
    chart.has_legend = bool(cspec.get("legend", nseries > 1 and ctype != "donut"))
    if chart.has_legend:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(TS["chart_axis"])
        chart.legend.font.color.rgb = C["ink_subtle"]
    chart.font.size = Pt(TS["chart_axis"])
    chart.font.color.rgb = C["ink_subtle"]
    chart.font.name = FONTS["latin"]

    palette = [RGBColor.from_string(c) for c in CH["comparison_series_colors"]]
    focal_cat = cspec.get("focal_category")
    for si, (s, ser) in enumerate(zip(cspec["series"], chart.plots[0].series)):
        color = RGBColor.from_string(s["color"]) if s.get("color") else (
            C["chart_gray"] if s.get("focal") is False else palette[si % len(palette)])
        if ctype == "line":
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.25 if s.get("focal", si == 0) else 1.5)
            ser.smooth = False
        elif ctype == "donut":
            for pi, pt_ in enumerate(ser.points):
                pt_.format.fill.solid()
                pt_.format.fill.fore_color.rgb = palette[pi % len(palette)] if focal_cat is None else (
                    C["primary"] if pi == focal_cat else C["chart_gray"])
        else:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
            forecast_from = cspec.get("forecast_from") if nseries == 1 else None
            if nseries == 1 and (focal_cat is not None or forecast_from is not None):
                for pi, pt_ in enumerate(ser.points):
                    if forecast_from is not None and pi >= int(forecast_from):
                        # IBCS-style forecast: pale fill + brand outline (vs solid actuals).
                        # A focal category in the forecast zone stays pale (it is still a
                        # forecast) but gets a heavier deep outline for emphasis.
                        pt_.format.fill.solid()
                        pt_.format.fill.fore_color.rgb = C["primary_pale"]
                        emphasized = focal_cat is not None and pi == focal_cat
                        pt_.format.line.color.rgb = C["primary_deep"] if emphasized else C["primary"]
                        pt_.format.line.width = Pt(1.75 if emphasized else 1.0)
                        continue
                    pt_.format.fill.solid()
                    pt_.format.fill.fore_color.rgb = (
                        color if focal_cat is None or pi == focal_cat else C["chart_gray"])

    if ctype != "donut":
        plot = chart.plots[0]
        if ctype in ("column", "bar", "stacked_column"):
            plot.gap_width = CH["gap_width_pct"]
            if nseries > 1 and ctype != "stacked_column":
                plot.overlap = CH["overlap_pct"]
        cat_ax, val_ax = chart.category_axis, chart.value_axis
        show_labels = cspec.get("value_labels", nseries == 1)
        axisless = cspec.get("axis_less", show_labels and nseries == 1 and ctype in ("column", "bar", "line"))
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = C["chart_gray_light"]
        val_ax.major_gridlines.format.line.width = Pt(CH["gridline_width_pt"])
        val_ax.format.line.fill.background()
        val_ax.major_tick_mark = XL_TICK_MARK.NONE
        val_ax.tick_labels.font.size = Pt(TS["chart_axis"])
        val_ax.tick_labels.font.color.rgb = C["ink_faint"]
        axis_fmt = cspec.get("axis_number_format") or cspec.get("number_format")
        if axis_fmt:
            val_ax.tick_labels.number_format = axis_fmt
            val_ax.tick_labels.number_format_is_linked = False
        if cspec.get("y_max") is not None:
            val_ax.maximum_scale = float(cspec["y_max"])
        elif cspec.get("annotation"):
            top = max(v for s in cspec["series"] for v in s["values"])
            val_ax.maximum_scale = top * 1.28
        if cspec.get("y_min") is None and all(
                v >= 0 for s in cspec["series"] for v in s["values"]):
            val_ax.minimum_scale = 0.0
        if cspec.get("y_min") is not None:
            val_ax.minimum_scale = float(cspec["y_min"])
        if axisless or cspec.get("hide_value_axis"):
            val_ax.visible = False
            val_ax.has_major_gridlines = False
        cat_ax.format.line.color.rgb = RGBColor.from_string(CH["axis_line_color"])
        cat_ax.format.line.width = Pt(0.75)
        cat_ax.major_tick_mark = XL_TICK_MARK.NONE
        cat_ax.tick_labels.font.size = Pt(TS["chart_axis"])
        cat_ax.tick_labels.font.color.rgb = C["ink_subtle"]
        if show_labels:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(TS["chart_label"])
            dl.font.color.rgb = C["ink"]
            dl.number_format = cspec.get("number_format", "#,##0.0")
            dl.number_format_is_linked = False
            if ctype in ("column", "bar"):
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            elif ctype == "line":
                dl.position = XL_LABEL_POSITION.ABOVE
    _chart_ea_fonts(chart)
    return gf


def add_unit_note(slide, x, y, unit: str | None):
    if unit:
        add_text(slide, x, y, 3.0, 0.2, [[("(" + unit + ")", TS["chart_axis"], 400, C["ink_faint"])]])


# ---------------------------------------------------------------- patterns

def p_cover(slide, spec, deck):
    add_rect(slide, 0, 0, TOKENS["slide"]["width_in"], TOKENS["slide"]["height_in"], C["canvas"])
    m = deck.get("meta", {})
    add_rect(slide, MX, 2.28, 0.5, 0.055, C["primary"])
    add_text(slide, MX, 2.62, CONTENT_W, 1.6,
             [[(spec.get("title", m.get("title", "")), TS["cover_title"], 700, C["ink"])]],
             line_spacing=1.18)
    if spec.get("subtitle"):
        add_text(slide, MX, 3.95, CONTENT_W, 0.6,
                 [[(spec["subtitle"], TS["cover_subtitle"], 400, C["ink_subtle"])]])
    meta_bits = [b for b in [spec.get("date", m.get("date")), spec.get("author", m.get("author"))] if b]
    if meta_bits:
        add_text(slide, MX, 6.55, CONTENT_W, 0.3,
                 [[("  |  ".join(meta_bits), 12, 400, C["ink_faint"])]])
    if m.get("confidential"):
        add_text(slide, TOKENS["slide"]["width_in"] - MX - 3.0, 0.42, 3.0, 0.3,
                 [[(m["confidential"], 10, 600, C["ink_faint"])]], align=PP_ALIGN.RIGHT)
    foot = LAY["footer"]
    add_text(slide, TOKENS["slide"]["width_in"] - foot["page_num_right_in"] - 0.9, foot["y_in"],
             0.9, foot["h_in"], [[("1", TS["page_number"], 400, C["ink_faint"])]], align=PP_ALIGN.RIGHT)


def p_agenda(slide, spec, deck):
    items = spec.get("items", [])
    y0, h = body_region(spec)
    x, w = grid(1, 10)
    row_h = min(0.95, h / max(1, len(items)))
    y0 += max(0.0, (h - len(items) * row_h) * 0.42)
    for i, it in enumerate(items):
        label = it if isinstance(it, str) else it.get("label", "")
        desc = None if isinstance(it, str) else it.get("desc")
        y = y0 + 0.2 + i * row_h
        add_text(slide, x, y, 0.6, 0.4,
                 [[(f"{i + 1:02d}", 18, 600, C["primary"])]])
        add_text(slide, x + 0.85, y - 0.02, w - 3.0, 0.4,
                 [[(label, 16, 600, C["ink"])]])
        if desc:
            add_text(slide, x + 0.85, y + 0.30, w - 1.0, 0.3,
                     [[(desc, 12.5, 400, C["ink_subtle"])]])
        if i < len(items) - 1:
            add_line(slide, x, y + row_h - 0.12, x + w, y + row_h - 0.12, C["rule"], 0.5)


def p_section_divider(slide, spec, deck):
    """章扉: 巨大数字(96pt、タイトル比3:1)が主役のスケールコントラスト。右は
    surface_tint の面パネル+章ナビ(現在章スポットライト)で、扉をサムネイルでも
    識別できるマクロ形状にする。左ブロックは光学中心(50%-0.25in リフト)に置く。"""
    W, H = TOKENS["slide"]["width_in"], TOKENS["slide"]["height_in"]
    add_rect(slide, 0, 0, W, H, C["canvas"])
    chapters = deck.get("_chapters") or []
    panel_x = 8.83
    if len(chapters) >= 2:
        add_rect(slide, panel_x, 0, W - panel_x, H, C["surface_tint"])
    num_pt = 96
    has_num = bool(spec.get("number"))
    # number 省略時は数字の段をレイアウトから除く(96pt 分の空洞をバーがブラケットしない)
    num_h = num_pt / 72.0 * 1.12 if has_num else 0.0
    num_gap = 0.08 if has_num else 0.0
    title_h = 0.55
    desc_h = 0.55 if spec.get("desc") else 0.0
    block_h = num_h + num_gap + title_h + desc_h
    by = H / 2 - block_h / 2 - 0.25
    block_cy = by + block_h / 2
    # 縦バーは先頭行のキャップトップから最終行の下端までをブラケットする(固定長にしない)
    bar_top = by + num_h * 0.16 if has_num else by + 0.06
    bar_bottom = by + num_h + num_gap + title_h + desc_h - 0.10
    add_rect(slide, MX - 0.1, bar_top, 0.075, bar_bottom - bar_top, C["primary"])
    if has_num:
        add_text(slide, MX + 0.2, by, 3.4, num_h, [[(str(spec["number"]).zfill(2), num_pt, 600, C["primary"])]])
    add_text(slide, MX + 0.2, by + num_h + num_gap, panel_x - MX - 0.7, title_h,
             [[(spec.get("title", ""), TS["divider_title"], 700, C["ink"])]])
    if spec.get("desc"):
        add_text(slide, MX + 0.2, by + num_h + num_gap + title_h, panel_x - MX - 0.7, desc_h,
                 [[(spec["desc"], 14, 400, C["ink_subtle"])]], line_spacing=1.3)
    if len(chapters) >= 2:
        row_h = 0.5
        nav_h = len(chapters) * row_h
        nx = panel_x + 0.5
        ny = block_cy - nav_h / 2
        for num, title in chapters:
            cur = title == spec.get("title", "")
            add_text(slide, nx, ny, 0.6, 0.32,
                     [[(str(num).zfill(2), 14, 600, C["primary"] if cur else C["ink_faint"])]])
            add_text(slide, nx + 0.6, ny + 0.01, W - nx - 1.1, 0.32,
                     [[(title, 14, 700 if cur else 400, C["ink"] if cur else C["ink_faint"])]])
            ny += row_h


def p_executive_summary(slide, spec, deck):
    points = spec.get("points", [])
    y0, h = body_region(spec)
    row_h = min(1.55, h / max(1, len(points)))
    x, w = grid(0, 12)
    y0 += max(0.0, (h - row_h * len(points)) * 0.4)
    for i, pt_ in enumerate(points):
        y = y0 + i * row_h + 0.05
        card_h = row_h - 0.14
        add_rect(slide, x, y, w, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        add_rect(slide, x, y, 0.055, card_h, C["primary"])
        kx = x + 0.25
        if pt_.get("kicker"):
            add_text(slide, kx, y, 2.1, card_h,
                     [[(pt_["kicker"], 12, 600, C["primary_deep"])]], anchor=MSO_ANCHOR.MIDDLE)
            kx = x + 2.45
        add_text(slide, kx, y + 0.08, x + w - kx - 0.25, card_h - 0.16,
                 [[(pt_.get("heading", ""), 14, 600, C["ink"])],
                  [(pt_.get("body", ""), TS["body"], 400, C["ink_subtle"])]],
                 line_spacing=1.18, space_after_pt=4, anchor=MSO_ANCHOR.MIDDLE)


def p_kpi_dashboard(slide, spec, deck):
    kpis = spec.get("kpis", [])
    if not kpis:
        return
    y0, h = body_region(spec)
    n = len(kpis)
    cols = 4 if n >= 7 else (3 if n >= 5 else min(n, 4) or 1)
    rows = math.ceil(n / cols)
    gut = LAY["gutter_in"]
    cw = (CONTENT_W - gut * (cols - 1)) / cols
    # 1行×4枚以下はヒーロースケール — まばらなスライドは装飾でなくサイズで埋める
    vsize = TS["kpi_value_hero"] if rows == 1 and n <= 4 else (TS["kpi_value"] if rows == 1 else 27)
    pad = 0.18
    inner_hs = []
    for k in kpis:
        ch_ = 0.52 + vsize / 60
        if k.get("delta"):
            ch_ += 0.24
        if k.get("note"):
            ch_ += _text_lines(k["note"], cw - 2 * pad, TS["kpi_sub"]) * 0.19 + 0.04
        inner_hs.append(ch_)
    content_h = max(inner_hs) + 0.16
    rh = min((h - gut * (rows - 1)) / rows, max(1.30, content_h))
    # フィルルール: ブロックが本文領域の58%未満なら「浮いた帯」— カードを伸ばして埋め、
    # 中身はカード内で垂直中央(ブロック間だけでなくブロック内の余白も設計する)
    rh = max(rh, min((h - gut * (rows - 1)) / rows, (0.58 * h - gut * (rows - 1)) / rows))
    block_h = rows * rh + (rows - 1) * gut
    y_base = y0 + 0.05 + max(0.0, (h - block_h) * 0.44)
    for i, k in enumerate(kpis):
        r_, c_ = divmod(i, cols)
        x = MX + c_ * (cw + gut)
        y = y_base + r_ * (rh + gut)
        # focal: the title's protagonist KPIs get the pale-brand card (主役を立てる)
        add_rect(slide, x, y, cw, rh,
                 C["primary_pale"] if k.get("focal") else C["surface_tint"],
                 radius_pt=LAY["card"]["radius_pt"])
        yin = y + max(0.14, (rh - inner_hs[i]) / 2)
        add_text(slide, x + pad, yin, cw - 2 * pad, 0.28,
                 [[(k.get("label", ""), TS["kpi_label"], 600, C["ink_subtle"])]])
        vparts = [(k.get("value", ""), vsize, 700, C["primary_deep"])]
        if k.get("unit"):
            vparts.append(_unit_part(k["unit"], vsize * 0.45))
        add_text(slide, x + pad, yin + 0.30, cw - 2 * pad, vsize / 60,
                 [vparts])
        by = yin + 0.38 + vsize / 60
        if k.get("delta"):
            arrow = {"up": "▲ ", "down": "▼ "}.get(k.get("delta_dir"), "")
            add_text(slide, x + pad, by, cw - 2 * pad, 0.24,
                     [[(arrow + k["delta"], TS["kpi_sub"], 600,
                        _delta_color(k.get("delta_dir"), k.get("positive_is_good", True)))]])
            by += 0.24
        if k.get("note"):
            add_text(slide, x + pad, by, cw - 2 * pad, max(0.2, rh - (by - y) - 0.08),
                     [[(k["note"], TS["kpi_sub"], 400, C["ink_faint"])]], line_spacing=1.12)


def _focal_bar_anchor(cx, cw, cspec):
    """Rough x-center of the focal category's bar inside a rendered chart frame."""
    n = max(1, len(cspec.get("categories", [])))
    i = cspec.get("focal_category")
    if i is None:  # 明示的な null もここに落ちる(get のデフォルトでは拾えない)
        i = n - 1
    plot_x, plot_w = cx + 0.15, cw - 0.35
    return plot_x + plot_w * (i + 0.5) / n


def add_chart_annotations(slide, cx, cy, cw, ch, cspec):
    """Exemplar-style annotations: YoY badge above the focal bar, optional trend arrow."""
    ann = cspec.get("annotation") or {}
    if not ann:
        return
    bx = _focal_bar_anchor(cx, cw, cspec)
    if ann.get("yoy") or ann.get("badge"):
        bw, bh = (1.5, 0.52) if ann.get("badge") else (1.15, 0.52)
        ox = min(bx - bw / 2, cx + cw - bw - 0.05)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ox), Inches(cy - 0.06), Inches(bw), Inches(bh))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["canvas"]
        sh.line.color.rgb = C["primary"]
        sh.line.width = Pt(1.5)
        sh.shadow.inherit = False
        _strip_style(sh)
        parts = [(ann["badge"], 12.5, 700, C["primary_deep"])] if ann.get("badge") else [
            ("YoY ", 10, 400, C["ink_subtle"]), (ann["yoy"], 13.5, 700, C["primary_deep"])]
        add_text(slide, ox, cy - 0.06, bw, bh, [parts],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    if ann.get("trend_arrow"):
        aw = min(1.15, cw * 0.16)
        sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(bx - aw - 0.5), Inches(cy + 0.5), Inches(aw), Inches(0.3))
        sh.adjustments[0] = 0.55
        sh.adjustments[1] = 0.55
        sh.rotation = -20
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["primary"]
        sh.line.fill.background()
        sh.shadow.inherit = False
        _strip_style(sh)


def p_chart_insight(slide, spec, deck):
    y0, h = body_region(spec)
    layout = spec.get("layout", "chart_left")
    takeaways = spec.get("takeaways", [])
    if layout == "chart_full" or not takeaways:
        cx, cw = grid(0, 12)
        add_unit_note(slide, cx, y0 + 0.02, spec.get("chart", {}).get("unit"))
        add_act_chart(slide, cx, y0 + 0.25, cw, h - 0.3, spec["chart"])
        add_chart_annotations(slide, cx, y0 + 0.25, cw, h - 0.3, spec["chart"])
        return
    cx, cw = grid(0, 7)
    tx, tw = grid(7, 5)
    add_unit_note(slide, cx, y0 + 0.02, spec.get("chart", {}).get("unit"))
    add_act_chart(slide, cx, y0 + 0.25, cw, h - 0.3, spec["chart"])
    add_chart_annotations(slide, cx, y0 + 0.25, cw, h - 0.3, spec["chart"])
    ty = y0 + 0.15
    add_text(slide, tx, ty, tw, 0.3, [[(spec.get("takeaways_heading", "要点"), 14, 600, C["primary_deep"])]])
    add_line(slide, tx, ty + 0.34, tx + tw, ty + 0.34, C["rule"], 0.75)
    ty += 0.5
    chars_per_line = max(8.0, (tw - 0.04) / 0.176)
    blocks = []
    for t in takeaways:
        head = t.get("heading") if isinstance(t, dict) else None
        body = t.get("body") if isinstance(t, dict) else str(t)
        body_lines = math.ceil(_ja_len(body) / chars_per_line) if body else 0
        bh = (0.28 if head else 0.0) + body_lines * 0.235 + 0.05
        blocks.append((head, body, bh))
    total_h = sum(b[2] for b in blocks)
    gap = min(0.42, max(0.16, (h - 0.7 - total_h) / max(1, len(blocks) - 1))) if len(blocks) > 1 else 0
    for head, body, bh in blocks:
        paras = []
        if head:
            paras.append([(head, 13, 600, C["ink"])])
        if body:
            paras.append([(body, TS["body"], 400, C["ink_subtle"])])
        add_text(slide, tx + 0.02, ty, tw - 0.04, bh, paras, line_spacing=1.2, space_after_pt=2)
        ty += bh + gap


def p_market_sizing(slide, spec, deck):
    """TAM/SAM/SOM horizontal funnel bars with per-stage description column."""
    stages = spec.get("stages", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(stages))
    row_h = min(1.7, (h - 0.2) / n)
    y0 += max(0.0, (h - n * row_h) * 0.4)
    bar_max_w = w * 0.52
    vals = [float(s.get("numeric", 0) or 0) for s in stages]
    vmax = max(vals) if any(vals) else 1.0
    for i, s in enumerate(stages):
        y = y0 + 0.1 + i * row_h
        frac = (vals[i] / vmax) if vmax else 1.0
        frac = max(frac, 0.16)
        bw = bar_max_w * (frac ** 0.5)
        fill = [C["chart_gray"], C["primary"], C["primary_deep"]][i % 3]
        txt = C["ink"] if i % 3 == 0 else C["canvas"]  # 白文字はグレー地でコントラスト不足
        add_rect(slide, x, y, bw, row_h - 0.28, fill, radius_pt=4)
        add_text(slide, x + 0.18, y + 0.10, bw - 0.3, 0.26, [[(s.get("label", ""), 12.5, 600, txt)]])
        add_text(slide, x + 0.18, y + 0.34, bw - 0.3, row_h - 0.6,
                 [[(s.get("value", ""), 22, 700, txt)]])
        dx = x + bar_max_w + 0.35
        add_text(slide, dx, y + 0.08, x + w - dx, row_h - 0.3,
                 [[(s.get("name", s.get("label", "")), 14, 600, C["ink"])],
                  [(s.get("desc", ""), TS["body"], 400, C["ink_subtle"])]],
                 line_spacing=1.15, space_after_pt=2)


def _cell_borders(cell, bottom_hex: str | None = "D8D4C9"):
    """Kill vertical/top borders; keep a thin bottom rule (banker table look)."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        el = tcPr.find(f"{A_NS}{tag}")
        if el is not None:
            tcPr.remove(el)
    els = []
    for tag in ("lnL", "lnR", "lnT"):
        ln = etree.SubElement(tcPr, f"{A_NS}{tag}")
        etree.SubElement(ln, f"{A_NS}noFill")
        els.append(ln)
    lnB = etree.SubElement(tcPr, f"{A_NS}lnB")
    if bottom_hex:
        lnB.set("w", "6350")
        fill = etree.SubElement(lnB, f"{A_NS}solidFill")
        etree.SubElement(fill, f"{A_NS}srgbClr").set("val", bottom_hex)
    else:
        etree.SubElement(lnB, f"{A_NS}noFill")
    els.append(lnB)
    for i, el in enumerate(els):
        tcPr.remove(el)
        tcPr.insert(i, el)


def _table_font(cell, text, size, weight, color, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = hw(text)
    _set_run_fonts(r, size, weight, color)


def add_act_table(slide, x, y, w, h, tspec: dict):
    """tspec: {headers:[], rows:[[...]], col_widths:[fractions], align:[l|r|c per col],
    emphasis_col: int|None, emphasis_row: int|None}"""
    headers = tspec.get("headers", [])
    rows = tspec.get("rows", [])
    nrows, ncols = len(rows) + 1, len(headers)
    g = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = g.table
    table.first_row = False
    table.horz_banding = False
    tbl = g._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # "no style" table style
    if tspec.get("col_widths"):
        fr = tspec["col_widths"]
        total_fr = sum(fr) or 1.0
        for ci, fracw in enumerate(fr):
            table.columns[ci].width = Emu(int(Inches(w) * fracw / total_fr))
    aligns = tspec.get("align", ["l"] * ncols)
    amap = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER}
    emph_col = tspec.get("emphasis_col")
    emph_row = tspec.get("emphasis_row")
    hdr_h = Inches(0.32)
    table.rows[0].height = hdr_h
    body_h = max(Inches(0.3), int((Inches(h) - hdr_h) / max(1, len(rows))))
    for ri in range(1, nrows):
        table.rows[ri].height = body_h
    for ci, htxt in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["primary_deep"]
        _table_font(cell, str(htxt), TS["table_header"], 600, C["canvas"], amap.get(aligns[ci], PP_ALIGN.LEFT))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _cell_borders(cell, "004F49")
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            focal = (emph_col is not None and ci == emph_col) or (emph_row is not None and ri - 1 == emph_row)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["primary_pale"] if focal else (
                C["surface_tint"] if ri % 2 == 0 else C["canvas"])
            weight = 600 if (focal or ci == 0) else 400
            color = C["ink"] if (focal or ci == 0) else C["ink_subtle"]
            _table_font(cell, str(val), TS["table_cell"], weight, color, amap.get(aligns[ci], PP_ALIGN.LEFT))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _cell_borders(cell)
    return g


def p_comparison_table(slide, spec, deck):
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    add_act_table(slide, x, y0 + 0.08, w, h - 0.16, spec["table"])


def p_competitive_landscape(slide, spec, deck):
    """2x2 positioning map. spec: {x_axis:{low,high}, y_axis:{low,high},
    players:[{name, x, y, focal}]} coords in 0..1."""
    y0, h = body_region(spec)
    size = min(h - 0.35, 5.0)
    px, pw = grid(0, 8)
    cx = px + 1.05
    cy = y0 + 0.15
    cw_ = pw - 2.05
    add_rect(slide, cx, cy, cw_, size, C["surface_tint"], radius_pt=0)
    add_line(slide, cx, cy + size / 2, cx + cw_, cy + size / 2, C["chart_gray"], 0.75)
    add_line(slide, cx + cw_ / 2, cy, cx + cw_ / 2, cy + size, C["chart_gray"], 0.75)
    ax, ay = spec.get("x_axis", {}), spec.get("y_axis", {})
    add_text(slide, cx - 1.02, cy + size / 2 - 0.10, 0.94, 0.3, [[(ax.get("low", ""), 10.5, 400, C["ink_faint"])]], align=PP_ALIGN.RIGHT)
    add_text(slide, cx + cw_ + 0.08, cy + size / 2 - 0.10, 0.92, 0.3, [[(ax.get("high", ""), 10.5, 400, C["ink_faint"])]])
    add_text(slide, cx + cw_ / 2 - 1.0, cy - 0.28, 2.0, 0.24, [[(ay.get("high", ""), 10.5, 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)
    add_text(slide, cx + cw_ / 2 - 1.0, cy + size + 0.05, 2.0, 0.24, [[(ay.get("low", ""), 10.5, 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)
    for pl in spec.get("players", []):
        dot = 0.42 if pl.get("focal") else 0.3
        dx = cx + pl["x"] * cw_ - dot / 2
        dy = cy + (1 - pl["y"]) * size - dot / 2
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(dot), Inches(dot))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C["primary"] if pl.get("focal") else C["chart_gray"]
        sh.line.fill.background()
        sh.shadow.inherit = False
        _strip_style(sh)
        add_text(slide, dx - 0.55, dy + dot + 0.02, dot + 1.1, 0.22,
                 [[(pl["name"], 11, 600 if pl.get("focal") else 400,
                    C["ink"] if pl.get("focal") else C["ink_subtle"])]], align=PP_ALIGN.CENTER)
    notes = spec.get("notes", [])
    if notes:
        nx, nw = grid(8, 4)
        ny = y0 + 0.15
        add_text(slide, nx, ny, nw, 0.3, [[(spec.get("notes_heading", "Positioning"), 14, 600, C["primary_deep"])]])
        add_line(slide, nx, ny + 0.3, nx + nw, ny + 0.3, C["rule"], 0.75)
        ny += 0.45
        chars_per_line = max(8.0, nw / 0.176)
        for t in notes:
            head = t.get("heading") if isinstance(t, dict) else None
            body = t.get("body") if isinstance(t, dict) else str(t)
            paras = ([[(head, 13, 600, C["ink"])]] if head else []) + ([[(body, TS["body"], 400, C["ink_subtle"])]] if body else [])
            body_lines = math.ceil(_ja_len(body) / chars_per_line) if body else 0
            block_h = (0.28 if head else 0.0) + body_lines * 0.235 + 0.05
            add_text(slide, nx, ny, nw, block_h, paras, line_spacing=1.2, space_after_pt=2)
            ny += block_h + 0.18


def p_financial_summary(slide, spec, deck):
    y0, h = body_region(spec)
    if spec.get("chart") and spec.get("table"):
        tx, tw = grid(0, 7)
        cx, cw = grid(7, 5)
        add_act_table(slide, tx, y0 + 0.08, tw - 0.15, h - 0.16, spec["table"])
        add_unit_note(slide, cx + 0.15, y0 + 0.02, spec["chart"].get("unit"))
        add_act_chart(slide, cx + 0.15, y0 + 0.28, cw - 0.15, h - 0.36, spec["chart"])
    elif spec.get("table"):
        p_comparison_table(slide, spec, deck)
    else:
        p_chart_insight(slide, spec, deck)


def p_waterfall(slide, spec, deck):
    """Shape-composed waterfall. items: [{label, value, kind: start|delta|end}]"""
    items = spec.get("items", [])
    if not items:
        return
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    chart_h = h - 1.0
    base_y = y0 + 0.42
    running, levels = 0.0, []
    for it in items:
        v = float(it["value"])
        if it.get("kind") == "start":
            running = v
            levels.append((0.0, v))
        elif it.get("kind") == "end":
            levels.append((0.0, v))
            running = v
        else:
            levels.append((running, running + v))
            running += v
    # 負の累計(赤字転落ブリッジ)も帯内に収める: 0 の基準線を span 内に置いてスケールする
    hi_max = max(max(a, b) for a, b in levels)
    lo_min = min(min(a, b) for a, b in levels)
    top = max(hi_max, 0.0) * 1.15 or 1.0
    bot = min(lo_min, 0.0) * 1.15
    span = (top - bot) or 1.0
    zero_y = base_y + chart_h * top / span
    n = len(items)
    slot = w / n
    bar_w = slot * 0.62
    add_line(slide, x, zero_y, x + w, zero_y, C["chart_gray"], 0.75)
    for i, (it, (lo, hi)) in enumerate(zip(items, levels)):
        v = float(it["value"])
        kind = it.get("kind", "delta")
        y_top = base_y + chart_h * (top - max(lo, hi)) / span
        bh = max(0.03, chart_h * abs(hi - lo) / span)
        bx = x + i * slot + (slot - bar_w) / 2
        if kind in ("start", "end"):
            fill = C["primary_deep"]
        else:
            fill = C["primary"] if v >= 0 else C["danger"]
        if it.get("forecast"):
            # 計画・予想のバーは実績と同じソリッド塗りにしない(IBCS)
            add_rect(slide, bx, y_top, bar_w, bh, C["primary_pale"],
                     line=C["primary_deep"], line_w_pt=1.5)
        else:
            add_rect(slide, bx, y_top, bar_w, bh, fill)
        if i > 0:
            edge = base_y + chart_h * (top - lo) / span if kind == "delta" else y_top
            add_line(slide, bx - (slot - bar_w), edge, bx, edge, C["chart_gray"], 0.5, dash="dash")
        if it.get("display") is not None:
            label_v = it["display"]
        elif kind == "delta":
            mag = f"{abs(v):,.1f}".rstrip("0").rstrip(".")
            label_v = ("△" if v < 0 else "+") + mag  # IR 慣行: 負値は - でなく △
        else:
            label_v = f"{v:,.1f}".rstrip("0").rstrip(".")
        add_text(slide, bx - slot * 0.19, y_top - 0.26, slot, 0.22,
                 [[(label_v, TS["chart_label"], 600, C["ink"])]], align=PP_ALIGN.CENTER)
        add_text(slide, x + i * slot + 0.02, base_y + chart_h + 0.08, slot - 0.04, 0.55,
                 [[(it.get("label", ""), TS["chart_axis"], 400, C["ink_subtle"])]], align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_unit_note(slide, x, y0 + 0.0, spec.get("unit"))


def p_roadmap(slide, spec, deck):
    """Phase roadmap: chevron phase headers + workstream rows with milestone text."""
    phases = spec.get("phases", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(phases))
    label_w = 0.0
    rows = spec.get("rows", [])
    if rows:
        label_w = 1.7
    pw = (w - label_w) / n
    def _phase_text(ph):
        t = ph.get("label", "")
        if ph.get("period"):
            t += "  " + ph["period"]
        return t
    notch_ratio = 0.6
    ph_h = 0.52
    inner_w = pw - 0.06 - ph_h * notch_ratio * 2
    two_line = any(_ja_len(_phase_text(ph)) * 12.5 / 72.0 > inner_w for ph in phases)
    if two_line:
        ph_h = 0.78
    for i, ph in enumerate(phases):
        px = x + label_w + i * pw
        add_chevron(slide, px, y0 + 0.05, pw - 0.06, ph_h, C["primary_deep"] if i == spec.get("focal_phase") else C["primary"])
        paras = [[(ph.get("label", ""), 12.5, 600, C["canvas"])]]
        if ph.get("period"):
            if two_line:
                paras.append([(ph["period"], 10.5, 400, C["canvas"])])
            else:
                paras[0].append(("  " + ph["period"], 10.5, 400, C["canvas"]))
        tip = ph_h * notch_ratio if not two_line else 0.52 * notch_ratio
        add_text(slide, px + 0.12, y0 + 0.05, pw - 0.06 - tip - 0.12, ph_h, paras,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False, line_spacing=1.1)
    ry = y0 + ph_h + 0.25
    if rows:
        row_h = (h - ph_h - 0.35) / max(1, len(rows))
        for r_i, row in enumerate(rows):
            yy = ry + r_i * row_h
            add_text(slide, x, yy, label_w - 0.15, row_h - 0.12,
                     [[(row.get("label", ""), 12, 600, C["ink"])]], anchor=MSO_ANCHOR.MIDDLE)
            for c_i in range(n):
                cell = (row.get("cells") or [""] * n)[c_i] if c_i < len(row.get("cells", [])) else ""
                cx_ = x + label_w + c_i * pw
                add_rect(slide, cx_, yy, pw - 0.06, row_h - 0.12, C["surface_tint"], radius_pt=4)
                if cell:
                    add_text(slide, cx_ + 0.14, yy, pw - 0.34, row_h - 0.12,
                             [[(str(cell), TS["body_small"], 400, C["ink_subtle"])]],
                             line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
    else:
        row_h = h - ph_h - 0.35
        for i, ph in enumerate(phases):
            px = x + i * pw
            add_rect(slide, px, ry, pw - 0.06, row_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
            bullets = ph.get("items", [])
            paras = [[("・" + b, TS["body_small"], 400, C["ink_subtle"])] for b in bullets]
            add_text(slide, px + 0.16, ry, pw - 0.38, row_h - 0.2, paras,
                     line_spacing=1.25, space_after_pt=5, anchor=MSO_ANCHOR.MIDDLE)


def p_two_column(slide, spec, deck):
    y0, h = body_region(spec)
    # card height follows the measured content of the fuller side (container ≒ content),
    # then the block is centered — a full-height card over short content reads as neglect
    def _side_h(blk, w):
        chars_w = w - 0.44
        used = 0.0
        for it in blk.get("items", []):
            head = it.get("heading") if isinstance(it, dict) else None
            body = it.get("body") if isinstance(it, dict) else str(it)
            lines = _text_lines(body, chars_w, TS["body"])
            used += (0.26 if head else 0.0) + lines * (TS["body"] / 72.0) * 1.3 + 0.04 + 0.16
        return used
    _, w_probe = grid(0, 6)
    content_h = max(_side_h(spec.get(s, {}), w_probe) for s in ("left", "right"))
    card_h = min(h - 0.1, max(1.8, 0.61 + content_h + 0.08))
    y_top = y0 + 0.05 + max(0.0, (h - 0.1 - card_h) * 0.44)
    for side, (c0, span) in (("left", (0, 6)), ("right", (6, 6))):
        blk = spec.get(side, {})
        x, w = grid(c0, span)
        add_rect(slide, x, y_top, w, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        mark = blk.get("mark")
        if mark == "cross":
            hd_fill = C["chart_gray"]
        elif mark == "check" or blk.get("focal"):
            hd_fill = C["primary_deep"]
        else:
            hd_fill = C["navy"]
        add_rect(slide, x, y_top, w, 0.42, hd_fill, radius_pt=LAY["card"]["radius_pt"])
        add_rect(slide, x, y_top + 0.23, w, 0.19, hd_fill)
        head_txt = blk.get("heading", side)
        if mark == "cross":
            head_txt = "× " + head_txt
        elif mark == "check":
            head_txt = "○ " + head_txt
        add_text(slide, x + 0.16, y_top, w - 0.32, 0.42,
                 [[(head_txt, 13, 600, C["canvas"])]], anchor=MSO_ANCHOR.MIDDLE)
        yy = y_top + 0.61
        for it in blk.get("items", []):
            head = it.get("heading") if isinstance(it, dict) else None
            body = it.get("body") if isinstance(it, dict) else str(it)
            paras = []
            if head:
                paras.append([(head, 12.5, 600, C["ink"])])
            if body:
                paras.append([(("" if head else "・") + body, TS["body"], 400, C["ink_subtle"])])
            body_lines = _text_lines(body, w - 0.44, TS["body"])
            block_h = (0.26 if head else 0.0) + body_lines * (TS["body"] / 72.0) * 1.3 + 0.04
            add_text(slide, x + 0.2, yy, w - 0.4, block_h, paras, line_spacing=1.22, space_after_pt=3)
            yy += block_h + 0.16


def p_process_flow(slide, spec, deck):
    """Horizontal step flow with descriptions under each step. Card height follows
    the fullest step's measured content; the whole block centers in the body region."""
    steps = spec.get("steps", [])
    y0, h = body_region(spec)
    x, w = grid(0, 12)
    n = max(1, len(steps))
    sw = w / n
    head_h = 0.5
    line_h = (TS["body_small"] / 72.0) * 1.25 + 5 / 72.0
    content_h = 0.0
    has_outcome = any(st.get("outcome") for st in steps)
    for st in steps:
        used = 0.30 if st.get("desc") else 0.0
        for b in st.get("items", []):
            used += _text_lines("・" + b, sw - 0.40, TS["body_small"]) * line_h
        content_h = max(content_h, used)
    outcome_h = 0.0
    if has_outcome:
        outcome_h = max(_text_lines(st.get("outcome", ""), sw - 0.40, TS["body_small"])
                        for st in steps) * 0.24 + 0.42
    card_h = min(h - head_h - 0.35, max(1.3, 0.14 + content_h + outcome_h + 0.16))
    # フィルルール: 帯(シェブロン+カード)が本文領域の60%未満だと「浮いた帯」になる。
    # カードを伸ばして埋め、中身はカード内で垂直中央に置く
    card_h = max(card_h, min(h - head_h - 0.35, 0.60 * h - head_h - 0.15))
    y_top = y0 + 0.05 + max(0.0, (h - head_h - 0.25 - card_h - 0.1) * 0.44)
    gut = LAY["gutter_in"]
    n_last = len(steps) - 1
    for i, st in enumerate(steps):
        sx = x + i * sw
        # 彩度は到達点に配給する: 途中ステップは淡テール、終端だけソリッド
        head_fill = C["primary"] if i == n_last else C["primary_pale"]
        label_color = C["canvas"] if i == n_last else C["primary_deep"]
        add_chevron(slide, sx, y_top, sw - gut, head_h, head_fill)
        tip = head_h * 0.6  # PENTAGON: 左辺フラット、右にだけ矢先
        add_text(slide, sx + 0.12, y_top, sw - gut - tip - 0.12, head_h,
                 [[(st.get("label", ""), 13, 600, label_color)]],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False)
        cy = y_top + head_h + 0.15
        add_rect(slide, sx, cy, sw - gut, card_h, C["surface_tint"], radius_pt=LAY["card"]["radius_pt"])
        paras = [[("・" + b, TS["body_small"], 400, C["ink_subtle"])] for b in st.get("items", [])]
        if st.get("desc"):
            paras.insert(0, [(st["desc"], 12.5, 600, C["ink"])])
        add_text(slide, sx + 0.16, cy + 0.14, sw - gut - 0.32, card_h - outcome_h - 0.28,
                 paras, line_spacing=1.25, space_after_pt=5, anchor=MSO_ANCHOR.MIDDLE)
        if st.get("outcome"):
            oy = cy + card_h - outcome_h + 0.06
            add_line(slide, sx + 0.16, oy, sx + sw - gut - 0.16, oy, C["rule"], 0.5)
            add_text(slide, sx + 0.16, oy + 0.04, sw - gut - 0.32, 0.18,
                     [[("▼", 9, 400, C["primary"])]], align=PP_ALIGN.CENTER)
            add_text(slide, sx + 0.16, oy + 0.24, sw - gut - 0.32, outcome_h - 0.30,
                     [[(st["outcome"], TS["body_small"], 600, C["ink"])]], line_spacing=1.2)


def p_quote_or_statement(slide, spec, deck):
    """Hero statement: 28pt(長文24pt)/1.5行送り。ルール+本文+出典+recap を実測高で
    1グループに組み、光学中心(残余の45%)へ置く — 固定枠に流し込むと「小さな塊が
    余白に浮く」最頻出欠陥になる。ルールはテキスト実高だけをブラケットする。"""
    y0, h = body_region(spec)
    x, w = grid(1, 10)
    recap = spec.get("recap", [])
    stmt = spec.get("statement", "")
    s_pt = TS["statement"] if _ja_len(stmt) <= 60 else 24
    lines = _text_lines(stmt, w, s_pt)
    text_h = lines * (s_pt / 72.0) * 1.5 + 0.06
    attr_h = 0.58 if spec.get("attribution") else 0.0
    recap_h = 1.05 if recap else 0.0
    group_h = text_h + attr_h + recap_h
    gy = y0 + max(0.0, (h - group_h) * 0.45)
    bar_w = 0.125
    add_rect(slide, x - 0.16 - bar_w, gy - 0.10, bar_w, text_h + 0.20, C["primary"])
    add_text(slide, x, gy, w, text_h, [[(stmt, s_pt, 400, C["ink"])]], line_spacing=1.5)
    ay = gy + text_h + 0.30
    if spec.get("attribution"):
        add_text(slide, x, ay, w, 0.28, [[("— " + spec["attribution"], 12.5, 400, C["ink_subtle"])]])
        ay += 0.58
    if recap:
        ry = ay + 0.05
        n = len(recap)
        rw = w / n
        add_line(slide, x, ry - 0.12, x + w, ry - 0.12, C["rule"], 0.75)
        for i, m in enumerate(recap):
            rx = x + i * rw
            add_text(slide, rx, ry + 0.05, rw - 0.2, 0.22, [[(m.get("label", ""), 11, 600, C["ink_subtle"])]])
            vparts = [(str(m.get("value", "")), 24, 700, C["primary_deep"])]
            if m.get("unit"):
                vparts.append(_unit_part(m["unit"], 12))
            add_text(slide, rx, ry + 0.28, rw - 0.2, 0.4, [vparts])




def _unit_part(unit: str, size: float, weight: int = 600):
    pad = "" if unit in ("%", "pt", "x", "倍") else " "
    return (pad + unit, size, weight, C["ink_subtle"])


def _metric_pair(slide, x, y, w, m, value_size=24):
    """label (small) / big value + small unit / colored delta — exemplar big-number grammar."""
    add_text(slide, x, y, w, 0.22, [[(m.get("label", ""), 11.5, 600, C["ink_subtle"])]])
    vparts = [(str(m.get("value", "")), value_size, 700, C["primary_deep"])]
    if m.get("unit"):
        vparts.append(_unit_part(m["unit"], value_size * 0.45))
    add_text(slide, x, y + 0.24, w, value_size / 60, [vparts])
    dy = y + 0.24 + value_size / 60
    if m.get("delta"):
        add_text(slide, x, dy, w, 0.2,
                 [[(m["delta"], 11, 600, _delta_color(m.get("delta_dir"), m.get("positive_is_good", True)))]])
        dy += 0.2
    return dy


def p_financial_highlights(slide, spec, deck):
    """Exemplar archetype: dot-labelled groups, bold claim, big-number metric pairs."""
    groups = spec.get("groups", [])
    y0, h = body_region(spec)
    n = max(1, len(groups))
    gut = 0.5
    gw = (CONTENT_W - gut * (n - 1)) / n
    est_h = 0.30
    for g in groups:
        gh = 0.30
        if g.get("claim"):
            gh += 0.34 * (1 if _ja_len(g["claim"]) <= gw / 0.24 else 2) + 0.12
        metrics = g.get("metrics", [])
        mcols = 2 if len(metrics) >= 2 else 1
        for row in range(math.ceil(len(metrics) / mcols)):
            vsize = 30 if row == 0 else 24
            gh += 0.24 + vsize / 60 + 0.2 + 0.28
        if g.get("note"):
            gh += 0.08 + _text_lines(g["note"], gw - 0.2, TS["body_small"]) * 0.24
        est_h = max(est_h, gh)
    y_shift = max(0.0, (h - 0.1 - est_h) * 0.44)
    for gi, g in enumerate(groups):
        gx = MX + gi * (gw + gut)
        gy = y0 + 0.1 + y_shift
        add_rect(slide, gx, gy + 0.03, 0.09, 0.09, C["primary"])
        add_text(slide, gx + 0.18, gy - 0.03, gw - 0.2, 0.24, [[(g.get("label", ""), 12.5, 600, C["ink"])]])
        gy += 0.30
        if g.get("claim"):
            add_text(slide, gx, gy, gw, 0.62,
                     [[(g["claim"], 17, 700, C["ink"])]], line_spacing=1.2)
            claim_lines = 1 if _ja_len(g["claim"]) <= gw / 0.24 else 2
            gy += 0.34 * claim_lines + 0.12
        metrics = g.get("metrics", [])
        mcols = 2 if len(metrics) >= 2 else 1
        mw = (gw - 0.3 * (mcols - 1)) / mcols
        row_y = [gy] * mcols
        for mi, m in enumerate(metrics):
            col = mi % mcols
            hero = m.get("hero", mi == 0)
            row_y[col] = _metric_pair(slide, gx + col * (mw + 0.3), row_y[col], mw, m,
                                      value_size=30 if hero else 24) + 0.28
        gy = max(row_y)
        if g.get("note"):
            add_line(slide, gx, gy, gx + gw - 0.2, gy, C["rule"], 0.5)
            add_text(slide, gx, gy + 0.08, gw - 0.2, 0.5,
                     [[(g["note"], TS["body_small"], 400, C["ink_subtle"])]], line_spacing=1.2)
        if gi < n - 1:
            add_line(slide, gx + gw + gut / 2, y0 + 0.10 + y_shift, gx + gw + gut / 2,
                     y0 + 0.10 + y_shift + est_h, C["rule"], 0.75)


def p_metrics_rows(slide, spec, deck):
    """Exemplar archetype: hairline rows of label | big value+unit | YoY (1-2 column groups)."""
    columns = spec.get("columns", [])
    y0, h = body_region(spec)
    n = max(1, len(columns))
    gut = 0.6
    gw = (CONTENT_W - gut * (n - 1)) / n
    max_rows = max((len(g.get("rows", [])) for g in columns), default=1)
    has_heading = any(g.get("heading") for g in columns)
    row_h_all = min(0.7, (h - 0.13 - (0.36 if has_heading else 0.0)) / max(1, max_rows))
    block_h = (0.36 if has_heading else 0.0) + max_rows * row_h_all
    y_shift = max(0.0, (h - 0.08 - block_h) * 0.44)
    for gi, g in enumerate(columns):
        gx = MX + gi * (gw + gut)
        gy = y0 + 0.08 + y_shift
        if g.get("heading"):
            add_text(slide, gx, gy, gw, 0.26, [[(g["heading"], 13.5, 600, C["ink"])]], align=PP_ALIGN.CENTER)
            gy += 0.36
        rows = g.get("rows", [])
        row_h = row_h_all
        for r in rows:
            emph = r.get("emphasis")
            if emph:
                add_rect(slide, gx - 0.08, gy + 0.02, gw + 0.16, row_h - 0.06, C["primary_pale"], radius_pt=4)
            add_text(slide, gx, gy, gw * 0.34, row_h,
                     [[(r.get("label", ""), 11.5, 600, C["ink"] if emph else C["ink_subtle"])]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
            vparts = [(str(r.get("value", "")), 19, 700, C["primary_deep"] if emph else C["ink"])]
            if r.get("unit"):
                vparts.append(_unit_part(r["unit"], 11, 400))
            add_text(slide, gx + gw * 0.34, gy, gw * 0.38, row_h, [vparts],
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            if r.get("delta"):
                add_text(slide, gx + gw * 0.76, gy, gw * 0.24, row_h,
                         [[(r["delta"], 11.5, 600, _delta_color(r.get("delta_dir", "up"), r.get("positive_is_good", True)))]],
                         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            add_line(slide, gx, gy + row_h - 0.02, gx + gw, gy + row_h - 0.02, C["rule"], 0.5)
            gy += row_h


def p_driver_decomposition(slide, spec, deck):
    """Exemplar archetype: KPI factor cards joined by operators (数 × 単価 = GMV)."""
    factors = spec.get("factors", [])
    y0, h = body_region(spec)
    n = max(1, len(factors))
    op_w = 0.55
    gut = 0.18
    cw = (CONTENT_W - op_w * (n - 1) - gut * 2 * (n - 1)) / n
    has_delta = any(f.get("delta") for f in factors)
    note_lines = max((_text_lines(f.get("note", ""), cw - 0.3, 10.5) for f in factors), default=0)
    card_h = min(h - 0.75, max(1.7, 0.26 + 0.30 + 0.62 + (0.30 if has_delta else 0.0)
                               + (0.10 + note_lines * 0.20 if note_lines else 0.0) + 0.24))
    cy = y0 + 0.15 + max(0.0, (h - 0.55 - card_h) * 0.44)
    ops = spec.get("operators") or (["×"] * (n - 2) + ["="] if n >= 2 else [])
    for i, f in enumerate(factors):
        x = MX + i * (cw + op_w + gut * 2)
        emphasized = f.get("focal", i == n - 1)
        add_rect(slide, x, cy, cw, card_h, C["primary_pale"] if emphasized else None,
                 line=None if emphasized else C["rule"], line_w_pt=1.0, radius_pt=8)
        add_text(slide, x + 0.15, cy + 0.24, cw - 0.3, 0.26,
                 [[(f.get("label", ""), 12, 600, C["ink_subtle"])]], align=PP_ALIGN.CENTER)
        vparts = [(str(f.get("value", "")), 34, 700, C["primary_deep"])]
        if f.get("unit"):
            vparts.append(_unit_part(f["unit"], 14))
        vy = cy + 0.58
        add_text(slide, x + 0.15, vy, cw - 0.3, 0.55, [vparts], align=PP_ALIGN.CENTER)
        if f.get("delta"):
            add_text(slide, x + 0.15, vy + 0.60, cw - 0.3, 0.24,
                     [[(f["delta"], 12, 600, _delta_color(f.get("delta_dir", "up"), f.get("positive_is_good", True)))]],
                     align=PP_ALIGN.CENTER)
        if f.get("note"):
            add_text(slide, x + 0.15, cy + card_h - 0.14 - note_lines * 0.20, cw - 0.3, note_lines * 0.20 + 0.06,
                     [[(f["note"], 10.5, 400, C["ink_faint"])]], align=PP_ALIGN.CENTER, line_spacing=1.15)
        if i < n - 1:
            add_text(slide, x + cw + gut, cy, op_w, card_h,
                     [[(ops[i] if i < len(ops) else "×", 26, 600, C["ink_subtle"])]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    caption = spec.get("formula_note")
    if caption is None:
        labels = [f.get("label", "").split("(")[0] for f in factors]
        parts = []
        for i, lb in enumerate(labels):
            parts.append(lb)
            if i < len(labels) - 1:
                parts.append(ops[i] if i < len(ops) else "×")
        caption = " ".join(parts)
    if caption:
        add_text(slide, MX, cy + card_h + 0.22, CONTENT_W, 0.26,
                 [[(caption, 12, 400, C["ink_faint"])]], align=PP_ALIGN.CENTER)




def p_guidance_progress(slide, spec, deck):
    """Yearly bars where the current year shows actual fill + dashed empty box up to
    the guidance range, with a side stack of progress facts. bars:
    [{label, value, display}] ... past years; current: {label, actual, actual_display,
    guidance_low, guidance_high, range_display}; side: [{label, value}]"""
    bars = spec.get("bars", [])
    cur = spec.get("current", {})
    y0, h = body_region(spec)
    x, w = grid(0, 8)
    sx, sw = grid(8, 4)
    chart_h = h - 0.75
    base_y = y0 + 0.30
    g_high = float(cur.get("guidance_high", cur.get("guidance_low", 0)) or 0)
    vals = [float(b.get("value", 0)) for b in bars] + [g_high]
    top = max(vals) * 1.18 if vals and max(vals) > 0 else 1.0
    n = len(bars) + 1
    slot = w / n
    bar_w = slot * 0.55
    add_line(slide, x, base_y + chart_h, x + w, base_y + chart_h, C["chart_gray"], 0.75)
    for i, b in enumerate(bars):
        v = float(b.get("value", 0))
        bh = chart_h * v / top
        bx = x + i * slot + (slot - bar_w) / 2
        add_rect(slide, bx, base_y + chart_h - bh, bar_w, bh, C["chart_gray"])
        add_text(slide, bx - slot * 0.2, base_y + chart_h - bh - 0.26, bar_w + slot * 0.4, 0.22,
                 [[(b.get("display", f"{v:,.0f}"), TS["chart_label"], 400, C["ink_subtle"])]], align=PP_ALIGN.CENTER, wrap=False)
        add_text(slide, x + i * slot, base_y + chart_h + 0.08, slot, 0.3,
                 [[(b.get("label", ""), TS["chart_axis"], 400, C["ink_subtle"])]], align=PP_ALIGN.CENTER)
    actual = float(cur.get("actual", 0))
    bx = x + len(bars) * slot + (slot - bar_w) / 2
    a_h = chart_h * actual / top
    g_h = chart_h * g_high / top
    add_rect(slide, bx, base_y + chart_h - a_h, bar_w, a_h, C["primary"])
    add_text(slide, bx, base_y + chart_h - a_h + 0.06, bar_w, 0.24,
             [[(cur.get("actual_display", f"{actual:,.0f}"), TS["chart_label"], 600, C["canvas"])]], align=PP_ALIGN.CENTER)
    if g_h > a_h + 0.02:
        add_rect(slide, bx, base_y + chart_h - g_h, bar_w, g_h - a_h, None,
                 line=C["primary_deep"], line_w_pt=1.25, dash="dash")
    if cur.get("range_display"):
        add_text(slide, bx - slot * 0.35, base_y + chart_h - g_h - 0.30, bar_w + slot * 0.7, 0.24,
                 [[(cur["range_display"], TS["chart_label"], 700, C["ink"])]], align=PP_ALIGN.CENTER, wrap=False)
    g_low = float(cur.get("guidance_low", 0) or 0)
    if g_low and g_low != g_high:
        low_h = chart_h * g_low / top
        add_line(slide, bx, base_y + chart_h - low_h, bx + bar_w, base_y + chart_h - low_h, C["primary_deep"], 0.75)
        add_text(slide, bx + bar_w + 0.08, base_y + chart_h - g_h - 0.05, 1.2, 0.2,
                 [[("上限 " + f"{g_high:,.0f}".rstrip("0").rstrip("."), 10.5, 600, C["ink_subtle"])]], wrap=False)
        add_text(slide, bx + bar_w + 0.08, base_y + chart_h - low_h + 0.04, 1.2, 0.2,
                 [[("下限 " + f"{g_low:,.0f}".rstrip("0").rstrip("."), 10.5, 600, C["ink_subtle"])]], wrap=False)
    add_text(slide, x + len(bars) * slot, base_y + chart_h + 0.08, slot, 0.3,
             [[(cur.get("label", ""), TS["chart_axis"], 600, C["ink"])]], align=PP_ALIGN.CENTER)
    add_unit_note(slide, x, y0 - 0.02, spec.get("unit"))
    side = spec.get("side", [])
    if side:
        sy = base_y + 0.2
        add_text(slide, sx, sy, sw, 0.28, [[(spec.get("side_heading", "進捗"), 14, 600, C["primary_deep"])]])
        add_line(slide, sx, sy + 0.32, sx + sw, sy + 0.32, C["rule"], 0.75)
        sy += 0.5
        for it in side:
            add_text(slide, sx, sy, sw * 0.55, 0.4, [[(it.get("label", ""), 12, 600, C["ink_subtle"])]],
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, sx + sw * 0.55, sy, sw * 0.45, 0.4,
                     [[(str(it.get("value", "")), 16, 700, C["ink"])]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            sy += 0.48


PATTERNS = {
    "cover": p_cover,
    "agenda": p_agenda,
    "section_divider": p_section_divider,
    "executive_summary": p_executive_summary,
    "kpi_dashboard": p_kpi_dashboard,
    "chart_insight": p_chart_insight,
    "market_sizing": p_market_sizing,
    "comparison_table": p_comparison_table,
    "competitive_landscape": p_competitive_landscape,
    "financial_summary": p_financial_summary,
    "waterfall": p_waterfall,
    "roadmap": p_roadmap,
    "two_column": p_two_column,
    "process_flow": p_process_flow,
    "statement": p_quote_or_statement,
    "financial_highlights": p_financial_highlights,
    "metrics_rows": p_metrics_rows,
    "driver_decomposition": p_driver_decomposition,
    "guidance_progress": p_guidance_progress,
}

NO_CHROME = {"cover", "section_divider"}


def build(spec_path: Path, out_path: Path) -> Path:
    deck = json.loads(spec_path.read_text())
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]
    slides = deck["slides"]
    total = len(slides)
    deck["_chapters"] = [(s.get("number") or i + 1, s.get("title", ""))
                         for i, s in enumerate([sp for sp in slides if sp.get("pattern") == "section_divider"])]
    for idx, spec in enumerate(slides, start=1):
        pattern = spec.get("pattern")
        if pattern not in PATTERNS:
            raise SystemExit(f"slide {idx}: unknown pattern '{pattern}'. valid: {sorted(PATTERNS)}")
        slide = prs.slides.add_slide(blank)
        if pattern in NO_CHROME:
            PATTERNS[pattern](slide, spec, deck)
            if pattern == "section_divider":
                add_footer(slide, spec, idx, total, deck)
        else:
            add_chrome(slide, spec, idx, total, deck)
            PATTERNS[pattern](slide, spec, deck)
            add_insight_bar(slide, spec)
        if spec.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = spec["speaker_notes"]
    prs.save(out_path)
    return out_path


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", type=Path)
    ap.add_argument("-o", "--out", type=Path, default=None)
    args = ap.parse_args()
    out = args.out or args.spec.with_suffix(".pptx")
    build(args.spec, out)
    print(f"built: {out}")


if __name__ == "__main__":
    main()
