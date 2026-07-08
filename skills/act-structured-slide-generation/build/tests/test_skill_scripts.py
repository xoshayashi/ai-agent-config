"""Tests for build_deck / validate_spec / verify_deck.

Run: python3 -m pytest skills/act-structured-slide-generation/build/tests
Requires python-pptx + Pillow (scripts/requirements.txt); tests skip if absent.
"""
import json
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

SKILL = Path(__file__).resolve().parent.parent.parent
SCRIPTS = SKILL / "build" / "scripts"
SAMPLE = SKILL / "examples" / "sample-deck.json"
SAMPLE_EARNINGS = SKILL / "examples" / "sample-earnings-deck.json"


def run(script, *args):
    return subprocess.run([sys.executable, str(SCRIPTS / script), *map(str, args)],
                          capture_output=True, text=True)


def test_skill_metadata_and_resource_map_are_discoverable():
    skill_md = (SKILL / "SKILL.md").read_text()
    desc = skill_md.split("description:", 1)[1].split("---", 1)[0].strip().strip('"')
    assert len(desc) <= 1024
    assert "grid-and-flex-strategy.md" in skill_md
    assert "slide-decision-engine.md" in skill_md
    assert "corpus-derived-composition-atoms.md" in skill_md
    assert "visual-qa-and-repair-rubric.md" in skill_md
    assert "assets/deck-workspace-template/" in skill_md
    assert (SKILL / "agents" / "openai.yaml").exists()
    assert (SKILL / "build" / "references" / "grid-and-flex-strategy.md").exists()
    assert (SKILL / "build" / "references" / "slide-decision-engine.md").exists()
    assert (SKILL / "build" / "references" / "ir-slide-design-principles.md").exists()
    assert (SKILL / "build" / "references" / "corpus-derived-composition-atoms.md").exists()
    assert (SKILL / "build" / "references" / "visual-qa-and-repair-rubric.md").exists()
    assert (SKILL / "assets" / "deck-workspace-template" / "outline.md").exists()


def test_reference_markdown_is_english_without_japanese_residue():
    import re

    japanese = re.compile(r"[ぁ-んァ-ン一-龯]")
    for path in sorted((SKILL / "build" / "references").glob("*.md")):
        text = path.read_text()
        assert not japanese.search(text), f"Japanese residue in {path.relative_to(SKILL)}"


def test_slide_judgment_requires_grid_and_flex_strategy_fields():
    ref = (SKILL / "build" / "references" / "slide-judgment-system.md").read_text()
    assert "these 22 fields" in ref
    assert "grid_role_map" in ref
    assert "main_axis" in ref
    for token in (
        "reader_question",
        "single_takeaway",
        "focal_object",
        "evidence_strategy",
        "composition_move",
        "density_control",
        "whitespace_role",
        "hierarchy_spine",
        "annotation_policy",
        "rhythm_role",
        "failure_mode",
        "repair_instruction",
    ):
        assert token in ref
    for token in (
        "grid_role_map",
        "column_span_plan",
        "alignment_spine",
        "body_band_plan",
        "edge_lock",
        "cross_slide_consistency",
        "main_axis",
        "cross_axis_align",
        "gap_scale",
        "grow_rule",
        "shrink_guard",
        "wrap_rule",
        "fill_repair",
    ):
        assert token in ref


def test_grid_flex_reference_requires_granular_non_template_contract():
    ref = (SKILL / "build" / "references" / "grid-and-flex-strategy.md").read_text()
    compact_ref = " ".join(ref.split())
    assert "Grid Contract" in ref
    assert "relationships, not coordinates" in ref
    assert "fixed template" in compact_ref
    for token in (
        "grid_role_map",
        "column_span_plan",
        "alignment_spine",
        "body_band_plan",
        "edge_lock",
        "cross_slide_consistency",
        "main_axis",
        "cross_axis_align",
        "gap_scale",
        "grow_rule",
        "shrink_guard",
        "wrap_rule",
        "fill_repair",
        "Fine-grained adjustment loop",
    ):
        assert token in ref


def test_ir_corpus_decision_engine_is_wired_as_judgment_not_templates():
    skill_md = (SKILL / "SKILL.md").read_text()
    decision = (SKILL / "build" / "references" / "slide-decision-engine.md").read_text()
    atoms = (SKILL / "build" / "references" / "corpus-derived-composition-atoms.md").read_text()
    principles = (SKILL / "build" / "references" / "ir-slide-design-principles.md").read_text()
    qa = (SKILL / "build" / "references" / "visual-qa-and-repair-rubric.md").read_text()
    anti = (SKILL / "build" / "references" / "anti-patterns.md").read_text()
    skill_compact = " ".join(skill_md.split())
    atoms_compact = " ".join(atoms.split())

    assert "source understanding -> reader_question -> single_takeaway -> focal_object" in skill_compact
    assert "slide-type template catalog" in decision
    assert "reader_question" in decision and "repair_instruction" in decision
    assert "existing_rule_refinement" in decision and "new_composition_atom" in decision
    assert "These are not slide types" in atoms_compact
    assert "Atom Mixing Rules" in atoms
    assert "Use when / Why it works / Parameters / Failure / Repair" in principles
    assert "Investor Lens" in qa and "Legal / Disclaimer Lens" in qa
    assert "F19" in anti and "F24" in anti


def test_design_contract_keeps_banker_base_and_modern_freshness():
    skill_md = (SKILL / "SKILL.md").read_text()
    principles = (SKILL / "build" / "references" / "design-principles.md").read_text()
    grid_ref = (SKILL / "build" / "references" / "grid-and-flex-strategy.md").read_text()
    rubric = json.loads((SKILL / "build" / "evals" / "rubric.json").read_text())
    assert "banker-grade / strategy-consulting base" in skill_md
    assert "Freshness comes second" in principles
    assert "grid/flex contract" in principles
    assert "investment-bank / strategy-consulting discipline" in grid_ref
    assert "Modern freshness" in json.dumps(rubric)


def test_anti_template_audit_and_no_page_numbers_are_contractual(tmp_path):
    skill_md = (SKILL / "SKILL.md").read_text()
    judgment = (SKILL / "build" / "references" / "slide-judgment-system.md").read_text()
    assert "Anti-Template Audit" in judgment
    assert "Do not create page numbers" in skill_md
    assert "Does the close look fixed" in judgment

    deck = {"meta": {}, "slides": [{
        "pattern": "cover",
        "title": "ページ番号を表示しない表紙",
        "subtitle": "脚注は不要",
    }, {
        "pattern": "statement",
        "title": "結論",
        "variant": "evidence_strip",
        "statement": "本文領域の構図で締める",
        "recap": [{"label": "確認指標", "value": "3", "unit": "点"}],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text.strip())
    assert "1" not in texts and "2" not in texts


def _minimal_deck(**slide_overrides):
    slide = {
        "pattern": "chart_insight",
        "title": "テスト市場は年率10%で拡大し、2030年に1兆円に到達する見込み",
        "chart": {"type": "column", "unit": "億円", "categories": ["2029", "2030"],
                  "series": [{"name": "市場", "values": [1, 2]}]},
        "source": "テスト統計2026",
    }
    slide.update(slide_overrides)
    return {"meta": {"title": "t"}, "slides": [slide]}


@pytest.mark.parametrize("sample", [SAMPLE, SAMPLE_EARNINGS], ids=["proposal", "earnings"])
def test_sample_spec_validates(sample):
    r = run("validate_spec.py", sample)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_unknown_pattern(tmp_path):
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(_minimal_deck(pattern="mystery")))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "unknown pattern" in r.stdout


def test_validate_rejects_off_palette_color(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"][0]["color"] = "FF00FF"
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "outside the Act palette" in r.stdout


def test_validate_rejects_double_accent(tmp_path):
    deck = _minimal_deck(insight="一言インサイト")
    deck["slides"][0]["chart"]["series"][0]["color"] = "ECC85A"
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "Accent" in r.stdout


def test_validate_rejects_waterfall_without_start_end(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "ブリッジは開始と終了で挟まれていなければならないという検証",
        "items": [{"label": "a", "value": 1}, {"label": "b", "value": 2}],
    }]}
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "waterfall" in r.stdout


@pytest.mark.parametrize("sample", [SAMPLE, SAMPLE_EARNINGS], ids=["proposal", "earnings"])
def test_build_and_verify_sample(tmp_path, sample):
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", sample, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert out.exists() and out.stat().st_size > 20000
    n = len(json.loads(sample.read_text())["slides"])
    assert len(pptx.Presentation(out).slides._sldIdLst) == n
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_verify_catches_forbidden_color(tmp_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    out = tmp_path / "black.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run_ = tb.text_frame.paragraphs[0].add_run()
    run_.text = "black text"
    run_.font.color.rgb = RGBColor(0, 0, 0)
    prs.save(out)
    r = run("verify_deck.py", out)
    assert r.returncode == 1 and "#000000" in r.stdout


def test_lint_render_clean_canvas_passes(tmp_path):
    from PIL import Image

    Image.new("RGB", (1466, 825), (0xFF, 0xFD, 0xFC)).save(tmp_path / "deck-01.png")
    r = run("lint_render.py", tmp_path)
    assert r.returncode == 0, r.stdout + r.stderr


def test_lint_render_detects_edge_clipping(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), (0xFF, 0xFD, 0xFC))
    # 下端に部分的にかかる濃色ブロック = 見切れ(全幅ではないので full-bleed 扱いにならない)
    ImageDraw.Draw(im).rectangle([300, 790, 900, 824], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    r = run("lint_render.py", tmp_path)
    assert r.returncode == 1 and "見切れ" in r.stdout


def test_validate_warns_title_line_mixing(tmp_path):
    long_title = "国内SaaS市場は年率13%で成長し、中堅企業セグメントの浸透率拡大が2030年まで続く見込みである"
    deck = _minimal_deck()
    deck["slides"].append({**_minimal_deck()["slides"][0], "title": long_title})
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "タイトル1行と2行が混在" in r.stdout


def test_validate_warns_subtitle_mixing(tmp_path):
    deck = _minimal_deck(subtitle="市場規模の推移")
    deck["slides"].append(_minimal_deck()["slides"][0])  # subtitle なし
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "サブタイトルの有無が混在" in r.stdout


def test_lint_render_detects_internal_gap(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), (0xFF, 0xFD, 0xFC))
    d = ImageDraw.Draw(im)
    d.rectangle([100, 200, 1300, 300], fill=(0x2D, 0x33, 0x2E))   # 上部ブロック
    d.rectangle([100, 700, 1300, 750], fill=(0x2D, 0x33, 0x2E))   # 下部バンド(中抜け大)
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "中抜け" in r.stdout


def test_validate_warns_fat_process_step(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "process_flow",
        "title": "参入判断後は、90日でローンチ準備を完了させる実行体制を組む",
        "steps": [{"label": "Step 1", "items": ["a", "b", "c", "d"]}],
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "3個超" in r.stdout


def test_process_flow_outcome_is_centered_destination_label(tmp_path):
    from pptx.enum.text import PP_ALIGN

    deck = {"meta": {}, "slides": [{
        "pattern": "process_flow",
        "title": "3段階の打ち手で90日以内に改善余地を具体化",
        "steps": [
            {"label": "Step 1", "items": ["対象整理"], "outcome": "改善余地の特定"},
            {"label": "Step 2", "items": ["優先順位付け"], "outcome": "実行順序の確定"},
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    para = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame and "改善余地の特定" in sh.text_frame.text:
                para = sh.text_frame.paragraphs[0]
    assert para is not None
    assert para.alignment == PP_ALIGN.CENTER
    sizes = [run.font.size.pt for run in para.runs if run.text and run.font.size]
    assert max(sizes) >= 18


def test_contact_sheet_header_strip(tmp_path):
    from PIL import Image

    for i in range(1, 4):
        Image.new("RGB", (320, 180), (240, 240, 240)).save(tmp_path / f"deck-{i}.png")
    r = run("contact_sheet.py", tmp_path, "--headers")
    assert r.returncode == 0, r.stdout + r.stderr
    assert (tmp_path / "header-strip.png").exists() and "3 slides" in r.stdout


def test_contact_sheet_builds_grid(tmp_path):
    from PIL import Image

    for i in range(1, 6):
        Image.new("RGB", (160, 90), (200 + i, 200, 200)).save(tmp_path / f"deck-{i:02d}.png")
    (tmp_path / "not-a-slide.png").touch()  # 数字サフィックスなし → 無視される
    r = run("contact_sheet.py", tmp_path, "--cols", 3)
    assert r.returncode == 0, r.stdout + r.stderr
    sheet = tmp_path / "contact-sheet.png"
    assert sheet.exists() and "5 slides" in r.stdout
    im = Image.open(sheet)
    assert im.width > 480 * 3 and im.height > 0


def test_contact_sheet_empty_dir_fails(tmp_path):
    r = run("contact_sheet.py", tmp_path)
    assert r.returncode == 1


def test_validate_outline_mode_prints_titles():
    r = run("validate_spec.py", SAMPLE, "--outline")
    assert r.returncode == 0
    titles = [s.get("title", "") for s in json.loads(SAMPLE.read_text())["slides"] if s.get("title")]
    assert titles[0] in r.stdout and titles[-1] in r.stdout


def test_validate_warns_fullwidth_alnum(tmp_path):
    # １０ = 全角「10」、ＡＲＲ = 全角「ARR」
    deck = _minimal_deck(title="市場は年率１０%で拡大し、ＡＲＲは2030年に1兆円に到達する")
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "全角英数字" in r.stdout


def test_validate_warns_unmarked_forecast_categories(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["forecast_from"] = 1  # "2030" にE等の予想表記がない
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "forecast_from" in r.stdout


def test_build_forecast_styling_stays_on_palette(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["categories"] = ["2029", "2030E"]
    deck["slides"][0]["chart"]["forecast_from"] = 1
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_lint_render_detects_bottom_whitespace(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), (0xFF, 0xFD, 0xFC))
    # 本文帯の上端だけにブロック → 下部に大きな空白(器と中身の不釣り合い)
    ImageDraw.Draw(im).rectangle([100, 185, 1300, 260], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "下部に大きな空白" in r.stdout


def test_lint_render_detects_tiny_center_island(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), (0xFF, 0xFD, 0xFC))
    # 本文帯の中央に小さいブロックだけが浮く = 上下とも広い余白
    ImageDraw.Draw(im).rectangle([420, 410, 1040, 470], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "オブジェクトが小さい" in r.stdout


def test_lint_render_baseline_diff_reports_changed_slides(tmp_path):
    from PIL import Image, ImageDraw

    cur, base = tmp_path / "cur", tmp_path / "base"
    cur.mkdir(); base.mkdir()
    for d in (cur, base):
        Image.new("RGB", (733, 412), (0xFF, 0xFD, 0xFC)).save(d / "deck-01.png")
    im = Image.new("RGB", (733, 412), (0xFF, 0xFD, 0xFC))
    ImageDraw.Draw(im).rectangle([100, 100, 400, 200], fill=(0x00, 0x8A, 0x80))
    im.save(cur / "deck-02.png")
    Image.new("RGB", (733, 412), (0xFF, 0xFD, 0xFC)).save(base / "deck-02.png")
    r = run("lint_render.py", cur, "--baseline", base)
    assert "DIFF: slide 2" in r.stdout and "slide 1" not in r.stdout.replace("2 slides", "")


def test_verify_catches_missing_ea_font(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    out = tmp_path / "noea.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run_ = tb.text_frame.paragraphs[0].add_run()
    run_.text = "日本語テキスト"
    run_.font.name = "Noto Sans JP"
    prs.save(out)
    r = run("verify_deck.py", out)
    assert r.returncode == 1 and "<a:ea>" in r.stdout


def test_validate_rejects_nonint_forecast_from(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["forecast_from"] = "x"
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "forecast_from" in r.stdout and "Traceback" not in r.stderr


def test_validate_warns_forecast_from_on_multiseries(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"].append({"name": "利益", "values": [0.5, 1.0]})
    deck["slides"][0]["chart"]["forecast_from"] = 1
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "単一系列" in r.stdout


def test_validate_rejects_nonnumeric_series_values(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"][0]["values"] = [1, "二"]
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "数値でない" in r.stdout and "Traceback" not in r.stderr


def test_waterfall_auto_label_uses_triangle(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "解約の自動ラベルはIR慣行の三角表記で描画されるという検証",
        "unit": "億円",
        "items": [{"label": "開始", "value": 10, "kind": "start"},
                  {"label": "解約", "value": -2.8},
                  {"label": "終了", "value": 7.2, "kind": "end"}],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "△2.8" in joined and "-2.8" not in joined


def test_kpi_signed_delta_is_not_prefixed_with_arrow(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "kpi_dashboard",
        "title": "サブスク売上29.8億円が継続収益の中核を形成",
        "kpis": [{
            "label": "サブスク売上",
            "value": "29.8",
            "unit": "億円",
            "delta": "+31% YoY",
            "delta_dir": "up",
            "note": "継続収益の中核",
            "focal": True,
        }],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "+31% YoY" in joined and "▲ +31% YoY" not in joined


def test_kpi_comparison_delta_is_not_prefixed_with_arrow(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "kpi_dashboard",
        "title": "NRR114%が継続収益品質を下支え",
        "kpis": [{
            "label": "NRR",
            "value": "114",
            "unit": "%",
            "delta": "前年同期111%",
            "delta_dir": "up",
            "focal": True,
        }],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "前年同期111%" in joined and "▲ 前年同期111%" not in joined


def test_validate_rejects_financial_summary_without_table_or_chart(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "financial_summary",
        "title": "売上高は3年で2.4倍の32.4億円に拡大し、利益率も同時に改善している",
        "source": "テスト",
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "financial_summary" in r.stdout


def test_validate_allows_current_only_guidance_progress(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "guidance_progress",
        "title": "Q2進捗率48%で通期132〜136億円を追走",
        "unit": "%",
        "bars": [],
        "current": {
            "label": "FY2026",
            "actual": 48,
            "actual_display": "Q2時点48%",
            "guidance_low": 100,
            "guidance_high": 100,
            "range_display": "通期100%"
        },
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0, r.stdout + r.stderr


def test_build_current_only_guidance_progress_stays_in_bounds(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "guidance_progress",
        "title": "Q2進捗率48%で通期132〜136億円を追走",
        "unit": "%",
        "bars": [],
        "current": {
            "label": "FY2026",
            "actual": 48,
            "actual_display": "Q2時点48%",
            "guidance_low": 100,
            "guidance_high": 100,
            "range_display": "通期100%"
        },
        "side_heading": "確認指標",
        "side": [
            {"label": "売上高", "value": "132〜136億"},
            {"label": "YoY", "value": "+26〜30%"},
            {"label": "中期ARR目標", "value": "250億円"},
            {"label": "中期利益率目標", "value": "15%"},
        ],
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "FY2026 現在地" in joined
    assert "48%" in joined
    first_side_top = last_side_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "売上高" == text.strip():
                    first_side_top = sh.top
                if "中期利益率目標" == text.strip():
                    last_side_top = sh.top
    assert first_side_top is not None and last_side_top is not None
    assert (last_side_top - first_side_top) / 914400 >= 1.8


def test_value_delta_spacing_has_visible_air(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "driver_decomposition",
        "title": "有料課金社数とARPAがARRを形成",
        "factors": [
            {"label": "有料課金社数", "value": "8,420", "unit": "社", "delta": "+18% YoY"},
            {"label": "ARPA", "value": "152", "unit": "万円", "delta": "+10% YoY"},
            {"label": "ARR", "value": "128", "unit": "億円", "delta": "+30% YoY", "focal": True},
        ],
        "operators": ["×", "="],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    value_top = delta_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "8,420" in text:
                    value_top = sh.top
                if "+18% YoY" in text:
                    delta_top = sh.top
    assert value_top is not None and delta_top is not None
    assert (delta_top - value_top) / 914400 >= 0.72


def test_financial_highlights_hero_yoy_uses_metric_subline_gap(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "financial_highlights",
        "title": "Q2売上32.4億円とARR128億円で成長継続",
        "groups": [{
            "label": "売上",
            "claim": "サブスク売上が成長を牽引",
            "metrics": [{"label": "Q2売上高", "value": "32.4", "unit": "億円", "delta": "+28% YoY", "hero": True}],
        }],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    value_top = delta_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "32.4" in text:
                    value_top = sh.top
                if "Q2売上高" in text and "+28% YoY" in text:
                    delta_top = sh.top
    assert value_top is not None and delta_top is not None
    assert (delta_top - value_top) / 914400 >= 0.90


def test_statement_split_evidence_variant_builds_and_verifies(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "statement",
        "title": "結論",
        "variant": "split_evidence",
        "statement": "ARR128億円、NRR114%、Q2進捗率48%が示す通期ガイダンスへの実行土台",
        "attribution": "テスト株式会社, 2026年7月",
        "recap": [
            {"label": "ARR", "value": "128", "unit": "億円", "focal": True},
            {"label": "NRR", "value": "114", "unit": "%"},
            {"label": "進捗率", "value": "48", "unit": "%"}
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_statement_evidence_strip_variant_builds_and_verifies(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "statement",
        "title": "投資家への結論",
        "variant": "evidence_strip",
        "statement": "ARR128億円、NRR114%、Q2進捗率48%を確認し、通期達成確度を継続評価",
        "recap_heading": "確認すべき3指標",
        "recap": [
            {"label": "ARR", "value": "128", "unit": "億円", "focal": True},
            {"label": "NRR", "value": "114", "unit": "%"},
            {"label": "進捗率", "value": "48", "unit": "%"}
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_player_missing_or_out_of_range_xy(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "competitive_landscape",
        "title": "統合スイート×中堅企業の右上象限は当社のみが占める空白地帯である",
        "x_axis": {"low": "単機能", "high": "統合"},
        "y_axis": {"low": "小規模", "high": "大企業"},
        "players": [{"name": "当社", "y": 0.5}, {"name": "A社", "x": 1.5, "y": 0.2}],
        "source": "テスト",
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1
    assert "がない/数値でない" in r.stdout and "範囲外" in r.stdout


def test_validate_allows_primary_insight_with_accent_series(tmp_path):
    # insight_style "primary" は Accent を消費しない — series の ECC85A 1 回は適法
    deck = _minimal_deck(insight="一言インサイト", insight_style="primary")
    deck["slides"][0]["chart"]["series"][0]["color"] = "ECC85A"
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_too_many_takeaways(tmp_path):
    deck = _minimal_deck(
        takeaways=[{"heading": f"論点{i}", "body": "根拠"} for i in range(1, 6)])
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "takeaways" in r.stdout


def test_validate_warns_long_content_run(tmp_path):
    content = _minimal_deck()["slides"][0]
    deck = {"meta": {"title": "t"},
            "slides": [dict(content, title=f"テスト市場は年率1{i}%で拡大し、2030年に1兆円に到達する見込み")
                       for i in range(7)]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "休符" in r.stdout
    # a divider inside the run resets the counter
    deck["slides"].insert(3, {"pattern": "section_divider", "number": 1, "title": "章"})
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert "休符" not in r.stdout


def test_validate_warns_wide_comparison_table(tmp_path):
    deck = {"meta": {"title": "t"}, "slides": [{
        "pattern": "comparison_table",
        "title": "4社比較でも当社の機能カバレッジ優位は変わらないことを示す",
        "table": {"headers": ["評価軸", "当社", "A社", "B社", "C社"],
                  "rows": [["対応領域", "◎", "○", "○", "△"]]},
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "4列以内" in r.stdout


def test_build_handles_null_focal_category_with_annotation(tmp_path):
    # LLM は「focal なし」を null で表現しがち — dict.get のデフォルトでは拾えない
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["focal_category"] = None
    deck["slides"][0]["chart"]["annotation"] = {"yoy": "+12%"}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_negative_waterfall_stays_inside_slide(tmp_path):
    from pptx.util import Inches

    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "一過性費用12億円の計上で通期は7億円の赤字に転落する見込みである",
        "unit": "億円",
        "items": [{"label": "期首", "value": 5, "kind": "start"},
                  {"label": "一過性費用", "value": -12},
                  {"label": "期末", "value": -7, "kind": "end"}],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.top is not None and sh.height is not None:
                assert sh.top + sh.height <= Inches(7.51), \
                    f"shape extends past slide bottom: {sh.shape_type}"


def test_section_divider_without_number_has_no_void(tmp_path):
    from pptx.util import Inches

    deck = {"meta": {}, "slides": [{"pattern": "section_divider", "title": "市場環境"}]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    tops = [sh.top for slide in pptx.Presentation(out).slides for sh in slide.shapes
            if sh.has_text_frame and "市場環境" in sh.text_frame.text]
    # number を省いたら 96pt 数字分の空洞は入らない(タイトルが光学中心近くに来る)
    assert tops and min(tops) < Inches(3.5)


def test_lint_render_flags_slide_count_mismatch(tmp_path):
    from PIL import Image

    for i in (1, 2):
        Image.new("RGB", (733, 412), (0xFF, 0xFD, 0xFC)).save(tmp_path / f"deck-{i}.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "枚数不一致" in r.stdout


def test_validate_rejects_unsupported_chart_type(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["type"] = "area"
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "未対応" in r.stdout and "Traceback" not in r.stderr


def test_chart_type_lists_stay_in_sync():
    sys.path.insert(0, str(SCRIPTS))
    try:
        import build_deck
        import validate_spec
        assert set(validate_spec.SUPPORTED_CHART_TYPES) == set(build_deck.CHART_TYPES)
    finally:
        sys.path.remove(str(SCRIPTS))


def test_validate_warns_missing_speaker_notes(tmp_path):
    deck = _minimal_deck()  # content slide without speaker_notes
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "speaker_notes" in r.stdout


def test_validate_talk_minutes_mismatch_warns(tmp_path):
    deck = _minimal_deck(speaker_notes="短い一言だけの語り。")
    deck["meta"]["talk_minutes"] = 10
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "talk_minutes=10" in r.stdout and "乖離" in r.stdout


def test_validate_talk_minutes_match_passes(tmp_path):
    deck = _minimal_deck(speaker_notes="あ" * 300)  # 300字 ≈ 1分
    deck["meta"]["talk_minutes"] = 1
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "乖離" not in r.stdout
    assert "talk script" in r.stdout and "target 1分" in r.stdout


def test_validate_rejects_nonnumeric_talk_minutes(tmp_path):
    deck = _minimal_deck(speaker_notes="語り。")
    deck["meta"]["talk_minutes"] = "abc"
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 1 and "talk_minutes" in r.stdout


def test_build_writes_speaker_notes_into_pptx(tmp_path):
    deck = _minimal_deck(speaker_notes="この市場は2年で倍になります。次のスライドで内訳を示します。")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    prs = pptx.Presentation(str(out))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "次のスライドで内訳を示します" in notes
