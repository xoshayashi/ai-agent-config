"""Tests for build_deck / validate_spec / verify_deck.

Run: python3 -m pytest skills/act-structured-slide-generation/tests
Requires python-pptx + Pillow (scripts/requirements.txt); tests skip if absent.
"""
import json
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

SKILL = Path(__file__).resolve().parent.parent
SCRIPTS = SKILL / "scripts"
SAMPLE = SKILL / "examples" / "sample-deck.json"
SAMPLE_EARNINGS = SKILL / "examples" / "sample-earnings-deck.json"

# 合成 PNG の地の色はトークンから引く — 直値で書くと canvas を変えた瞬間に
# lint(トークン基準)と食い違い、テストだけが古い色を主張する
CANVAS_RGB = tuple(
    int(json.loads((SKILL / "references" / "tokens.json").read_text())["colors"]["canvas"][i:i + 2], 16)
    for i in (0, 2, 4)
)


def run(script, *args):
    return subprocess.run([sys.executable, str(SCRIPTS / script), *map(str, args)],
                          capture_output=True, text=True)


def test_skill_metadata_and_resource_map_are_discoverable():
    skill_md = (SKILL / "SKILL.md").read_text()
    desc = skill_md.split("description:", 1)[1].split("---", 1)[0].strip().strip('"')
    assert len(desc) <= 1024
    assert "grid-and-flex-strategy.md" in skill_md
    assert "slide-decision-engine.md" in skill_md
    assert "composition-atoms.md" in skill_md
    assert "visual-qa-and-repair-rubric.md" in skill_md
    assert "assets/deck-workspace-template/" in skill_md
    assert (SKILL / "agents" / "openai.yaml").exists()
    assert (SKILL / "references" / "grid-and-flex-strategy.md").exists()
    assert (SKILL / "references" / "slide-decision-engine.md").exists()
    assert (SKILL / "references" / "ir-slide-design-principles.md").exists()
    assert (SKILL / "references" / "composition-atoms.md").exists()
    assert (SKILL / "references" / "visual-qa-and-repair-rubric.md").exists()
    assert (SKILL / "assets" / "deck-workspace-template" / "outline.md").exists()


def test_reference_markdown_is_english_without_japanese_residue():
    import re

    japanese = re.compile(r"[ぁ-んァ-ン一-龯]")
    for path in sorted((SKILL / "references").glob("*.md")):
        text = path.read_text()
        assert not japanese.search(text), f"Japanese residue in {path.relative_to(SKILL)}"


def test_slide_judgment_requires_grid_and_flex_strategy_fields():
    ref = (SKILL / "references" / "slide-judgment-system.md").read_text()
    assert "these 22 fields" in ref
    assert "grid_role_map" in ref
    assert "main_axis" in ref
    for token in (
        "reader_question",
        "single_takeaway",
        "focal_object",
        "evidence_strategy",
        "composition_move",
        "density_control",
        "whitespace_role",
        "hierarchy_spine",
        "annotation_policy",
        "rhythm_role",
        "failure_mode",
        "repair_instruction",
    ):
        assert token in ref
    for token in (
        "grid_role_map",
        "column_span_plan",
        "alignment_spine",
        "body_band_plan",
        "edge_lock",
        "cross_slide_consistency",
        "main_axis",
        "cross_axis_align",
        "gap_scale",
        "grow_rule",
        "shrink_guard",
        "wrap_rule",
        "fill_repair",
    ):
        assert token in ref


def test_grid_flex_reference_requires_granular_non_template_contract():
    ref = (SKILL / "references" / "grid-and-flex-strategy.md").read_text()
    compact_ref = " ".join(ref.split())
    assert "Grid Contract" in ref
    assert "relationships, not coordinates" in ref
    assert "fixed template" in compact_ref
    for token in (
        "grid_role_map",
        "column_span_plan",
        "alignment_spine",
        "body_band_plan",
        "edge_lock",
        "cross_slide_consistency",
        "main_axis",
        "cross_axis_align",
        "gap_scale",
        "grow_rule",
        "shrink_guard",
        "wrap_rule",
        "fill_repair",
        "Fine-grained adjustment loop",
    ):
        assert token in ref


def test_ir_decision_engine_is_wired_as_judgment_not_templates():
    skill_md = (SKILL / "SKILL.md").read_text()
    decision = (SKILL / "references" / "slide-decision-engine.md").read_text()
    atoms = (SKILL / "references" / "composition-atoms.md").read_text()
    principles = (SKILL / "references" / "ir-slide-design-principles.md").read_text()
    qa = (SKILL / "references" / "visual-qa-and-repair-rubric.md").read_text()
    anti = (SKILL / "references" / "anti-patterns.md").read_text()
    skill_compact = " ".join(skill_md.split())
    atoms_compact = " ".join(atoms.split())

    assert "source understanding -> reader_question -> single_takeaway -> focal_object" in skill_compact
    assert "slide-type template catalog" in decision
    assert "reader_question" in decision and "repair_instruction" in decision
    assert "existing_rule_refinement" in decision and "new_composition_atom" in decision
    assert "These are not slide types" in atoms_compact
    assert "Atom Mixing Rules" in atoms
    assert "Use when / Why it works / Parameters / Failure / Repair" in principles
    assert "Investor Lens" in qa and "Legal / Disclaimer Lens" in qa
    assert "F19" in anti and "F24" in anti


def test_design_contract_keeps_banker_base_and_modern_freshness():
    skill_md = (SKILL / "SKILL.md").read_text()
    principles = (SKILL / "references" / "design-principles.md").read_text()
    grid_ref = (SKILL / "references" / "grid-and-flex-strategy.md").read_text()
    rubric = json.loads((SKILL / "evals" / "rubric.json").read_text())
    assert "banker-grade / strategy-consulting base" in skill_md
    assert "Freshness comes second" in principles
    assert "grid/flex contract" in principles
    assert "investment-bank / strategy-consulting discipline" in grid_ref
    assert "Modern freshness" in json.dumps(rubric)


def test_anti_template_audit_and_no_page_numbers_are_contractual(tmp_path):
    skill_md = (SKILL / "SKILL.md").read_text()
    judgment = (SKILL / "references" / "slide-judgment-system.md").read_text()
    assert "Anti-Template Audit" in judgment
    assert "Do not create page numbers" in skill_md
    assert "Does the close look fixed" in judgment

    deck = {"meta": {}, "slides": [{
        "pattern": "cover",
        "title": "ページ番号を表示しない表紙",
        "subtitle": "脚注は不要",
    }, {
        "pattern": "statement",
        "title": "結論",
        "variant": "evidence_strip",
        "statement": "本文領域の構図で締める",
        "recap": [{"label": "確認指標", "value": "3", "unit": "点"}],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text.strip())
    assert "1" not in texts and "2" not in texts


def _minimal_deck(**slide_overrides):
    # ヘッダー契約: title と subtitle は全スライド必須で、それぞれ必ず1行
    slide = {
        "pattern": "chart_insight",
        "title": "テスト市場は年率10%で拡大し、2030年に1兆円に到達する見込み",
        "subtitle": "テスト市場規模の推移(2024-2030年)",
        "chart": {"type": "column", "unit": "億円", "categories": ["2029", "2030"],
                  "series": [{"name": "市場", "values": [1, 2]}]},
        "source": "テスト統計2026",
    }
    slide.update(slide_overrides)
    return {"meta": {"title": "t"}, "slides": [slide]}


@pytest.mark.parametrize("sample", [SAMPLE, SAMPLE_EARNINGS], ids=["proposal", "earnings"])
def test_sample_spec_validates(sample):
    r = run("validate_spec.py", sample)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_unknown_pattern(tmp_path):
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(_minimal_deck(pattern="mystery")))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "unknown pattern" in r.stdout


def test_validate_rejects_off_palette_color(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"][0]["color"] = "FF00FF"
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "outside the Act palette" in r.stdout


def test_validate_rejects_double_accent(tmp_path):
    deck = _minimal_deck(insight="一言インサイト")
    deck["slides"][0]["chart"]["series"][0]["color"] = "ECC85A"
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "Accent" in r.stdout


def test_validate_rejects_waterfall_without_start_end(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "ブリッジは開始と終了で挟まれていなければならないという検証",
        "items": [{"label": "a", "value": 1}, {"label": "b", "value": 2}],
    }]}
    bad = tmp_path / "bad.json"
    bad.write_text(json.dumps(deck))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "waterfall" in r.stdout


@pytest.mark.parametrize("sample", [SAMPLE, SAMPLE_EARNINGS], ids=["proposal", "earnings"])
def test_build_and_verify_sample(tmp_path, sample):
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", sample, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert out.exists() and out.stat().st_size > 20000
    n = len(json.loads(sample.read_text())["slides"])
    assert len(pptx.Presentation(out).slides._sldIdLst) == n
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_verify_catches_forbidden_color(tmp_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    out = tmp_path / "black.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run_ = tb.text_frame.paragraphs[0].add_run()
    run_.text = "black text"
    run_.font.color.rgb = RGBColor(0, 0, 0)
    prs.save(out)
    r = run("verify_deck.py", out)
    assert r.returncode == 1 and "#000000" in r.stdout


def test_lint_render_clean_canvas_passes(tmp_path):
    from PIL import Image

    Image.new("RGB", (1466, 825), CANVAS_RGB).save(tmp_path / "deck-01.png")
    r = run("lint_render.py", tmp_path)
    assert r.returncode == 0, r.stdout + r.stderr


def test_lint_render_detects_edge_clipping(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), CANVAS_RGB)
    # 下端に部分的にかかる濃色ブロック = 見切れ(全幅ではないので full-bleed 扱いにならない)
    ImageDraw.Draw(im).rectangle([300, 790, 900, 824], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    r = run("lint_render.py", tmp_path)
    assert r.returncode == 1 and "見切れ" in r.stdout


def _write(tmp_path, deck):
    p = tmp_path / "deck.json"
    p.write_text(json.dumps(deck, ensure_ascii=False))
    return p


# --- ヘッダー契約: title と subtitle は全スライド必須・それぞれ必ず1行 -----------------

def test_header_contract_is_declarative_and_derived(tmp_path):
    """契約はデータ(tokens)で宣言し、字数上限は描画幾何から導出する — スライド個別の
    ハードコードを持たない。ゆえに(1)未知のパターンにも既定契約が効き、(2)型スケールを
    変えれば上限も自動追従する。"""
    sys.path.insert(0, str(SKILL / "scripts"))
    from deck_text import header_slots, load_tokens, one_line_chars

    tokens = load_tokens()

    # (1) 契約に個別定義のないパターンでも既定(タイトル+副題、各1行)が効く
    slots = {s["slot"]: s for s in header_slots("some_future_pattern", tokens)}
    assert set(slots) == {"title", "subtitle"}
    assert all(s["lines"] == 1 for s in slots.values())
    assert slots["subtitle"]["field"] == "subtitle"

    # 上限は tokens に直書きされた数値ではなく、幅 ÷ 字送り から導出されている
    hdr_w = tokens["layout"]["header"]["title_w_in"]
    assert slots["title"]["max_chars"] == int(one_line_chars(hdr_w, tokens["type_scale_pt"]["action_title"]))
    assert not any("max_chars" in k for k in tokens["text_budget"] if k.startswith(("action_title", "subtitle", "cover", "divider")))

    # (2) 型スケールを大きくすると1行容量は自動的に縮む(定数の二重管理がない)
    bigger = json.loads(json.dumps(tokens))
    bigger["type_scale_pt"]["action_title"] *= 2
    shrunk = {s["slot"]: s for s in header_slots("some_future_pattern", bigger)}
    assert shrunk["title"]["max_chars"] < slots["title"]["max_chars"]

    # 描画のしかたが違うパターンだけが上書きする(カバーの副題は2行、章扉の副題は desc)
    cover = {s["slot"]: s for s in header_slots("cover", tokens)}
    assert cover["subtitle"]["lines"] == 2
    divider = {s["slot"]: s for s in header_slots("section_divider", tokens)}
    assert divider["subtitle"]["field"] == "desc" and divider["subtitle"]["lines"] == 1


def test_validate_rejects_title_that_would_wrap(tmp_path):
    # 2行に折り返す長さのタイトルは仕様欠陥 — 縮小も折返しもさせず ERROR で突き返す
    long_title = "国内SaaS市場は年率13%で成長し、中堅企業セグメントの浸透率拡大が2030年まで続く見込みである"
    r = run("validate_spec.py", _write(tmp_path, _minimal_deck(title=long_title)))
    assert r.returncode == 1 and "title が1行に収まらない" in r.stdout


def test_validate_rejects_subtitle_that_would_wrap(tmp_path):
    long_sub = "国内SaaS市場規模の推移と予測、および中堅企業セグメントの浸透率と競合ポジションの比較"
    r = run("validate_spec.py", _write(tmp_path, _minimal_deck(subtitle=long_sub)))
    assert r.returncode == 1 and "subtitle が1行に収まらない" in r.stdout


def test_validate_rejects_newline_in_title_or_subtitle(tmp_path):
    r = run("validate_spec.py", _write(tmp_path, _minimal_deck(title="市場は拡大\n2030年に1兆円へ")))
    assert r.returncode == 1 and "title は1行ちょうどで書く" in r.stdout and "改行禁止" in r.stdout
    r = run("validate_spec.py", _write(tmp_path, _minimal_deck(subtitle="市場規模\n推移")))
    assert r.returncode == 1 and "subtitle は1行ちょうどで書く" in r.stdout


def test_validate_requires_subtitle_on_every_slide(tmp_path):
    deck = _minimal_deck()
    del deck["slides"][0]["subtitle"]
    r = run("validate_spec.py", _write(tmp_path, deck))
    assert r.returncode == 1 and "subtitle がない" in r.stdout


def test_validate_requires_title_on_every_pattern_reported_once(tmp_path):
    # タイトル欠落はヘッダー契約チェックだけが報告する(PATTERNS 側と二重に出さない)
    for slide in ({"pattern": "statement", "statement": "結びの一文", "subtitle": "提言の要旨"},
                  {"pattern": "agenda", "items": ["市場", "戦略"], "subtitle": "本日の論点構成"},
                  {"pattern": "cover", "subtitle": "副題1行目\n副題2行目"}):
        r = run("validate_spec.py", _write(tmp_path, {"meta": {}, "slides": [slide]}))
        title_errors = [l for l in r.stdout.splitlines()
                        if l.startswith("ERROR") and ("title" in l or "タイトル" in l)]
        assert r.returncode == 1
        assert len(title_errors) == 1, title_errors
        assert "title がない" in title_errors[0]
        assert "missing required field 'title'" not in r.stdout


def test_validate_requires_exactly_two_line_cover_subtitle(tmp_path):
    def cover(sub):
        return {"meta": {"title": "t"}, "slides": [
            {"pattern": "cover", "title": "国内SaaS参入戦略", "subtitle": sub}]}

    # 1行 → NG(カバーだけは必ず2行)
    r = run("validate_spec.py", _write(tmp_path, cover("市場機会の評価と参入シナリオの提言")))
    assert r.returncode == 1 and "ちょうど2行" in r.stdout
    # 3行 → NG
    r = run("validate_spec.py", _write(tmp_path, cover("一行目\n二行目\n三行目")))
    assert r.returncode == 1 and "ちょうど2行" in r.stdout
    # 各行が折返す長さ → NG
    long_line = "市場機会の評価と参入シナリオの提言および中堅企業セグメントの浸透率と競合ポジションの比較検討"
    r = run("validate_spec.py", _write(tmp_path, cover(f"{long_line}\n二行目")))
    assert r.returncode == 1 and "1行目 が1行に収まらない" in r.stdout
    # ちょうど2行 → OK
    r = run("validate_spec.py", _write(tmp_path, cover("市場機会の評価と参入シナリオの提言\n経理領域を起点とした段階参入")))
    assert r.returncode == 0, r.stdout


def test_validate_section_divider_uses_desc_not_subtitle(tmp_path):
    # 章扉はヘッダー chrome を持たない — subtitle を書いても描画されないので誤りとして弾き、
    # 副題は desc に書かせる(desc も1行)
    base = {"pattern": "section_divider", "number": 1, "title": "市場・競争環境"}
    r = run("validate_spec.py", _write(tmp_path, {"meta": {}, "slides": [dict(base, subtitle="副題", desc="市場規模と競合の位置")]}))
    assert r.returncode == 1 and "描画しない" in r.stdout and "desc" in r.stdout

    r = run("validate_spec.py", _write(tmp_path, {"meta": {}, "slides": [base]}))
    assert r.returncode == 1 and "desc がない" in r.stdout

    r = run("validate_spec.py", _write(tmp_path, {"meta": {}, "slides": [dict(base, desc="市場規模と競合の位置")]}))
    assert r.returncode == 0, r.stdout


def test_build_renders_both_cover_subtitle_lines(tmp_path):
    from pptx.util import Inches

    deck = {"meta": {}, "slides": [{
        "pattern": "cover",
        "title": "国内SaaS参入戦略",
        "subtitle": "市場機会の評価と参入シナリオの提言\n経理領域を起点とした段階参入",
    }]}
    out = tmp_path / "deck.pptx"
    assert run("build_deck.py", _write(tmp_path, deck), "-o", out).returncode == 0
    texts = [sh.text_frame.text for sh in pptx.Presentation(out).slides[0].shapes if sh.has_text_frame]
    sub = next(t for t in texts if "市場機会の評価" in t)
    assert "経理領域を起点とした段階参入" in sub  # 2行目が切り落とされていない
    box = next(sh for sh in pptx.Presentation(out).slides[0].shapes
               if sh.has_text_frame and "市場機会の評価" in sh.text_frame.text)
    assert box.height > Inches(0.6)  # 2行分の高さが確保されている


def test_lint_render_detects_internal_gap(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), CANVAS_RGB)
    d = ImageDraw.Draw(im)
    d.rectangle([100, 200, 1300, 300], fill=(0x2D, 0x33, 0x2E))   # 上部ブロック
    d.rectangle([100, 700, 1300, 750], fill=(0x2D, 0x33, 0x2E))   # 下部バンド(中抜け大)
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "中抜け" in r.stdout


def test_validate_warns_fat_process_step(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "process_flow",
        "title": "参入判断後は90日でローンチ準備を完了させる",
        "subtitle": "ローンチまでの実行ステップ",
        "steps": [{"label": "Step 1", "items": ["a", "b", "c", "d"]}],
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "3個超" in r.stdout


def test_bullets_are_real_paragraph_bullets_centered_by_glyph_design(tmp_path):
    """箇条書きは段落の bullet(buChar)+ぶら下げインデントで組む — 図形の円を本文の横に
    置くと、記号が本文と連動せず編集で取り残される(一般的な箇条書きにならない)。

    記号はベースライン揃えで打たれるので、縦位置は「グリフのインク中心 × buSzPct」で
    決まる。字面中央(0.3805em)に乗るのは和文中黒(中心 0.380em)を等倍で使うときだけ —
    縮小すると中心も比例して下がり、点が本文より低く見える(以前の不具合)。"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    sys.path.insert(0, str(SKILL / "scripts"))
    import build_deck as B

    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    deck = {"meta": {}, "slides": [{
        "pattern": "process_flow",
        "title": "箇条書きの記号は段落の bullet で組む",
        "subtitle": "記号の縦位置と字下げの検証",
        "steps": [{"label": "Step 1", "desc": "検証",
                   "items": ["折返しを起こすために十分に長い項目をここに置いて二行以上にする",
                             "短い項目"]}],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0

    slide = pptx.Presentation(out).slides[0]
    paras = [p for sh in slide.shapes if sh.has_text_frame
             for p in sh.text_frame.paragraphs
             if p._p.find(f"{A_NS}pPr/{A_NS}buChar") is not None]
    assert len(paras) == 2, f"項目数ぶんの箇条書き段落: {len(paras)}"

    for p in paras:
        pPr = p._p.find(f"{A_NS}pPr")
        # 和文の中黒を等倍で: これ以外だと記号の中心が字面中央から下へずれる
        assert pPr.find(f"{A_NS}buChar").get("char") == B.BULLET_CHAR
        assert pPr.find(f"{A_NS}buSzPct").get("val") == "100000"
        assert pPr.find(f"{A_NS}buFont").get("typeface") == B.FONTS["ea"]
        # ぶら下げインデント: 1行目に記号、折返し行は marL に揃う
        marL, indent = int(pPr.get("marL")), int(pPr.get("indent"))
        assert marL == Inches(B.BULLET_INDENT_IN) and indent == -marL

    # 記号を図形で置いていないこと(本文と連動しない浮いた点を作らない)
    def _is_oval(sh):
        try:
            return sh.auto_shape_type == MSO_SHAPE.OVAL
        except (ValueError, AttributeError):
            return False

    assert not any(_is_oval(sh) for sh in slide.shapes), "記号が図形の円で描かれている"


def test_bullet_block_honors_middle_anchor(tmp_path):
    """anchor=MIDDLE の呼び出し(roadmap の items カード)では、箇条書きブロック全体が
    枠の中で縦中央に置かれる(テキストフレームの縦位置指定に委ねる)。"""
    from pptx.enum.text import MSO_ANCHOR
    sys.path.insert(0, str(SKILL / "scripts"))

    deck = {"meta": {}, "slides": [{
        "pattern": "roadmap",
        "title": "フェーズごとの打ち手を3段階で示す",
        "subtitle": "ロードマップ(items 指定)",
        "phases": [
            {"label": "Phase 1", "period": "FY27", "items": ["短い項目"]},
            {"label": "Phase 2", "period": "FY28",
             "items": ["項目A", "項目B", "項目C", "項目D"]},
            {"label": "Phase 3", "period": "FY29", "items": ["短い項目"]},
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0

    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    slide = pptx.Presentation(out).slides[0]
    bullet_boxes = [sh for sh in slide.shapes if sh.has_text_frame
                    and any(p._p.find(f"{A_NS}pPr/{A_NS}buChar") is not None
                            for p in sh.text_frame.paragraphs)]
    assert len(bullet_boxes) == 3, f"フェーズごとに1つの箇条書き枠: {len(bullet_boxes)}"
    for box in bullet_boxes:
        assert box.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "MIDDLE 寄せが効いていない"


def test_table_closing_rule_is_heavier_than_the_row_hairlines(tmp_path):
    """表の最下段の下罫は「表を閉じる罫」なので、行間のヘアラインより太くする。
    罫の色は rule トークンから引く(直値を持つとパレット変更がテーブルに届かない)。"""
    sys.path.insert(0, str(SKILL / "scripts"))
    import build_deck as B

    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    deck = {"meta": {}, "slides": [{
        "pattern": "comparison_table",
        "title": "3行の表で下端の締め罫を確認する",
        "subtitle": "罫の太さの検証",
        "table": {"headers": ["区分", "当社", "他社"],
                  "rows": [["A", "1", "2"], ["B", "3", "4"], ["C", "5", "6"]]},
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0

    table = next(sh.table for sh in pptx.Presentation(out).slides[0].shapes if sh.has_table)

    def bottom_rule(cell):
        lnB = cell._tc.find(f".//{A_NS}lnB")
        clr = lnB.find(f".//{A_NS}srgbClr")
        return int(lnB.get("w")) / 12700, clr.get("val")

    inner_w, inner_hex = bottom_rule(table.cell(1, 0))     # 1行目(行間のヘアライン)
    last_w, last_hex = bottom_rule(table.cell(3, 0))       # 最終行(表を閉じる罫)
    assert inner_w == 0.5, inner_w
    assert last_w == B.TABLE_CLOSING_RULE_PT > inner_w, (last_w, inner_w)
    # 色は両方とも rule トークン(直値のハードコードが復活していないこと)
    assert inner_hex == last_hex == str(B.C["rule"])


def test_ground_color_comes_from_the_template_not_an_object(tmp_path):
    """地の色はテンプレート(スライドマスターの背景)で与える。全面を覆う矩形で塗ると、
    編集時に本文の下で毎回つかんでしまうオブジェクトが1枚ずつ増えるだけになる。"""
    sys.path.insert(0, str(SKILL / "scripts"))
    import build_deck as B

    P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    for sample in (SAMPLE, SAMPLE_EARNINGS):
        out = tmp_path / f"{sample.stem}.pptx"
        assert run("build_deck.py", sample, "-o", out).returncode == 0
        prs = pptx.Presentation(out)

        # 1. マスターの背景に canvas 色が入っている(= テンプレートとして継承される)
        master = prs.slide_masters[0]._element
        cSld = master.find(f"{P_NS}cSld")
        bg = cSld.find(f"{P_NS}bg")
        assert bg is not None, f"{sample.name}: マスターに背景が設定されていない"
        clr = bg.find(f".//{A_NS}srgbClr")
        assert clr is not None and clr.get("val") == str(B.C["canvas"]), (
            f"{sample.name}: マスター背景が canvas 色でない")
        # cSld の子要素は bg が先頭でなければ不正な OOXML になる
        assert cSld[0].tag == f"{P_NS}bg"

        # 2. スライド側には全面を覆うオブジェクトが1つも無い
        sw, sh = prs.slide_width, prs.slide_height
        for i, slide in enumerate(prs.slides, 1):
            full = [s for s in slide.shapes
                    if s.width >= sw * 0.98 and s.height >= sh * 0.98]
            assert not full, f"{sample.name} slide {i}: 全面を覆うオブジェクトがある"


def test_process_flow_outcome_is_centered_destination_label(tmp_path):
    from pptx.enum.text import PP_ALIGN

    deck = {"meta": {}, "slides": [{
        "pattern": "process_flow",
        "title": "3段階の打ち手で90日以内に改善余地を具体化",
        "steps": [
            {"label": "Step 1", "items": ["対象整理"], "outcome": "改善余地の特定"},
            {"label": "Step 2", "items": ["優先順位付け"], "outcome": "実行順序の確定"},
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    para = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame and "改善余地の特定" in sh.text_frame.text:
                para = sh.text_frame.paragraphs[0]
    assert para is not None
    assert para.alignment == PP_ALIGN.CENTER
    sizes = [run.font.size.pt for run in para.runs if run.text and run.font.size]
    assert max(sizes) >= 18


def test_contact_sheet_header_strip(tmp_path):
    from PIL import Image

    for i in range(1, 4):
        Image.new("RGB", (320, 180), (240, 240, 240)).save(tmp_path / f"deck-{i}.png")
    r = run("contact_sheet.py", tmp_path, "--headers")
    assert r.returncode == 0, r.stdout + r.stderr
    assert (tmp_path / "header-strip.png").exists() and "3 slides" in r.stdout


def test_contact_sheet_builds_grid(tmp_path):
    from PIL import Image

    for i in range(1, 6):
        Image.new("RGB", (160, 90), (200 + i, 200, 200)).save(tmp_path / f"deck-{i:02d}.png")
    (tmp_path / "not-a-slide.png").touch()  # 数字サフィックスなし → 無視される
    r = run("contact_sheet.py", tmp_path, "--cols", 3)
    assert r.returncode == 0, r.stdout + r.stderr
    sheet = tmp_path / "contact-sheet.png"
    assert sheet.exists() and "5 slides" in r.stdout
    im = Image.open(sheet)
    assert im.width > 480 * 3 and im.height > 0


def test_contact_sheet_empty_dir_fails(tmp_path):
    r = run("contact_sheet.py", tmp_path)
    assert r.returncode == 1


def test_validate_outline_mode_prints_titles():
    r = run("validate_spec.py", SAMPLE, "--outline")
    assert r.returncode == 0
    titles = [s.get("title", "") for s in json.loads(SAMPLE.read_text())["slides"] if s.get("title")]
    assert titles[0] in r.stdout and titles[-1] in r.stdout


def test_validate_warns_fullwidth_alnum(tmp_path):
    # １０ = 全角「10」、ＡＲＲ = 全角「ARR」
    deck = _minimal_deck(title="市場は年率１０%で拡大し、ＡＲＲは2030年に1兆円に到達する")
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "全角英数字" in r.stdout


def test_validate_warns_unmarked_forecast_categories(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["forecast_from"] = 1  # "2030" にE等の予想表記がない
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "forecast_from" in r.stdout


def test_build_forecast_styling_stays_on_palette(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["categories"] = ["2029", "2030E"]
    deck["slides"][0]["chart"]["forecast_from"] = 1
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_edge_scan_tolerance_lets_pale_full_bleed_panels_through(tmp_path):
    """edge_scan の許容差(TOL 26)は緩いままにする。章扉の淡いパネル(surface_tint)は
    設計として端まで敷くので、これを「端で見切れ」と鳴らしてはいけない。
    balance_scan の tol 8 をここへ持ち込むと、章扉が丸ごと誤検知になる。"""
    from PIL import Image, ImageDraw
    sys.path.insert(0, str(SKILL / "scripts"))
    import lint_render as L

    tint = tuple(int(json.loads((SKILL / "references" / "tokens.json").read_text())
                     ["colors"]["surface_tint"][i:i + 2], 16) for i in (0, 2, 4))
    im = Image.new("RGB", (1466, 825), CANVAS_RGB)
    ImageDraw.Draw(im).rectangle([970, 0, 1465, 824], fill=tint)   # 章扉の右パネル(端まで)

    assert L.edge_scan(im) == [], "淡い全面パネルを見切れとして誤検知している"
    # 濃い要素が端に触れていれば従来どおり鳴る(検知力は落ちていない)
    ImageDraw.Draw(im).rectangle([0, 0, 400, 3], fill=(0x2D, 0x33, 0x2E))
    assert any("top edge" in h for h in L.edge_scan(im))


def test_lint_treats_plain_white_as_ground_not_content(tmp_path):
    """canvas は白に近いが白ではない(FCFCFB で最大差4)。純白を地として扱わないと、
    テンプレート背景が抜けて白く出たレンダーが「一面コンテンツ」に見え、空白系の検査が
    黙って効かなくなる。カード面(surface_tint)は白とも差が大きいので構造のまま残る。"""
    from PIL import Image, ImageDraw

    # 背景が抜けた(純白)レンダー: 本文帯の上端だけにブロック → 下部空白を検出できること
    im = Image.new("RGB", (1466, 825), (0xFF, 0xFF, 0xFF))
    ImageDraw.Draw(im).rectangle([100, 185, 1300, 260], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "下部に大きな空白" in r.stdout, r.stdout


def test_lint_render_detects_bottom_whitespace(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), CANVAS_RGB)
    # 本文帯の上端だけにブロック → 下部に大きな空白(器と中身の不釣り合い)
    ImageDraw.Draw(im).rectangle([100, 185, 1300, 260], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "下部に大きな空白" in r.stdout


def test_lint_render_detects_tiny_center_island(tmp_path):
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (1466, 825), CANVAS_RGB)
    # 本文帯の中央に小さいブロックだけが浮く = 上下とも広い余白
    ImageDraw.Draw(im).rectangle([420, 410, 1040, 470], fill=(0x2D, 0x33, 0x2E))
    im.save(tmp_path / "deck-01.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "オブジェクトが小さい" in r.stdout


def test_lint_render_baseline_diff_reports_changed_slides(tmp_path):
    from PIL import Image, ImageDraw

    cur, base = tmp_path / "cur", tmp_path / "base"
    cur.mkdir(); base.mkdir()
    for d in (cur, base):
        Image.new("RGB", (733, 412), CANVAS_RGB).save(d / "deck-01.png")
    im = Image.new("RGB", (733, 412), CANVAS_RGB)
    ImageDraw.Draw(im).rectangle([100, 100, 400, 200], fill=(0x00, 0x8A, 0x80))
    im.save(cur / "deck-02.png")
    Image.new("RGB", (733, 412), CANVAS_RGB).save(base / "deck-02.png")
    r = run("lint_render.py", cur, "--baseline", base)
    assert "DIFF: slide 2" in r.stdout and "slide 1" not in r.stdout.replace("2 slides", "")


def test_verify_catches_missing_ea_font(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    out = tmp_path / "noea.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run_ = tb.text_frame.paragraphs[0].add_run()
    run_.text = "日本語テキスト"
    run_.font.name = "Noto Sans JP"
    prs.save(out)
    r = run("verify_deck.py", out)
    assert r.returncode == 1 and "<a:ea>" in r.stdout


def test_validate_rejects_nonint_forecast_from(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["forecast_from"] = "x"
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "forecast_from" in r.stdout and "Traceback" not in r.stderr


def test_validate_warns_forecast_from_on_multiseries(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"].append({"name": "利益", "values": [0.5, 1.0]})
    deck["slides"][0]["chart"]["forecast_from"] = 1
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 0 and "単一系列" in r.stdout


def test_validate_rejects_nonnumeric_series_values(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["series"][0]["values"] = [1, "二"]
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "数値でない" in r.stdout and "Traceback" not in r.stderr


def test_waterfall_auto_label_uses_triangle(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "解約の自動ラベルはIR慣行の三角表記で描画されるという検証",
        "unit": "億円",
        "items": [{"label": "開始", "value": 10, "kind": "start"},
                  {"label": "解約", "value": -2.8},
                  {"label": "終了", "value": 7.2, "kind": "end"}],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "△2.8" in joined and "-2.8" not in joined


def test_kpi_signed_delta_is_not_prefixed_with_arrow(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "kpi_dashboard",
        "title": "サブスク売上29.8億円が継続収益の中核を形成",
        "kpis": [{
            "label": "サブスク売上",
            "value": "29.8",
            "unit": "億円",
            "delta": "+31% YoY",
            "delta_dir": "up",
            "note": "継続収益の中核",
            "focal": True,
        }],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "+31% YoY" in joined and "▲ +31% YoY" not in joined


def test_kpi_comparison_delta_is_not_prefixed_with_arrow(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "kpi_dashboard",
        "title": "NRR114%が継続収益品質を下支え",
        "kpis": [{
            "label": "NRR",
            "value": "114",
            "unit": "%",
            "delta": "前年同期111%",
            "delta_dir": "up",
            "focal": True,
        }],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "前年同期111%" in joined and "▲ 前年同期111%" not in joined


def test_validate_rejects_financial_summary_without_table_or_chart(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "financial_summary",
        "title": "売上高は3年で2.4倍の32.4億円に拡大し、利益率も同時に改善している",
        "source": "テスト",
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "financial_summary" in r.stdout


def test_validate_allows_current_only_guidance_progress(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "guidance_progress",
        "title": "Q2進捗率48%で通期132〜136億円を追走",
        "subtitle": "通期ガイダンス進捗(売上高)",
        "unit": "%",
        "bars": [],
        "current": {
            "label": "FY2026",
            "actual": 48,
            "actual_display": "Q2時点48%",
            "guidance_low": 100,
            "guidance_high": 100,
            "range_display": "通期100%"
        },
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0, r.stdout + r.stderr


def test_build_current_only_guidance_progress_stays_in_bounds(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "guidance_progress",
        "title": "Q2進捗率48%で通期132〜136億円を追走",
        "unit": "%",
        "bars": [],
        "current": {
            "label": "FY2026",
            "actual": 48,
            "actual_display": "Q2時点48%",
            "guidance_low": 100,
            "guidance_high": 100,
            "range_display": "通期100%"
        },
        "side_heading": "確認指標",
        "side": [
            {"label": "売上高", "value": "132〜136億"},
            {"label": "YoY", "value": "+26〜30%"},
            {"label": "中期ARR目標", "value": "250億円"},
            {"label": "中期利益率目標", "value": "15%"},
        ],
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr
    texts = []
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    joined = " ".join(texts)
    assert "FY2026 現在地" in joined
    assert "48%" in joined
    first_side_top = last_side_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "売上高" == text.strip():
                    first_side_top = sh.top
                if "中期利益率目標" == text.strip():
                    last_side_top = sh.top
    assert first_side_top is not None and last_side_top is not None
    assert (last_side_top - first_side_top) / 914400 >= 1.8


def test_value_delta_spacing_has_visible_air(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "driver_decomposition",
        "title": "有料課金社数とARPAがARRを形成",
        "factors": [
            {"label": "有料課金社数", "value": "8,420", "unit": "社", "delta": "+18% YoY"},
            {"label": "ARPA", "value": "152", "unit": "万円", "delta": "+10% YoY"},
            {"label": "ARR", "value": "128", "unit": "億円", "delta": "+30% YoY", "focal": True},
        ],
        "operators": ["×", "="],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    value_top = delta_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "8,420" in text:
                    value_top = sh.top
                if "+18% YoY" in text:
                    delta_top = sh.top
    assert value_top is not None and delta_top is not None
    assert (delta_top - value_top) / 914400 >= 0.72


def test_financial_highlights_hero_yoy_uses_metric_subline_gap(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "financial_highlights",
        "title": "Q2売上32.4億円とARR128億円で成長継続",
        "groups": [{
            "label": "売上",
            "claim": "サブスク売上が成長を牽引",
            "metrics": [{"label": "Q2売上高", "value": "32.4", "unit": "億円", "delta": "+28% YoY", "hero": True}],
        }],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    value_top = delta_top = None
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                text = sh.text_frame.text
                if "32.4" in text:
                    value_top = sh.top
                if "Q2売上高" in text and "+28% YoY" in text:
                    delta_top = sh.top
    assert value_top is not None and delta_top is not None
    assert (delta_top - value_top) / 914400 >= 0.90


def test_statement_split_evidence_variant_builds_and_verifies(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "statement",
        "title": "結論",
        "variant": "split_evidence",
        "statement": "ARR128億円、NRR114%、Q2進捗率48%が示す通期ガイダンスへの実行土台",
        "attribution": "テスト株式会社, 2026年7月",
        "recap": [
            {"label": "ARR", "value": "128", "unit": "億円", "focal": True},
            {"label": "NRR", "value": "114", "unit": "%"},
            {"label": "進捗率", "value": "48", "unit": "%"}
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_statement_evidence_strip_variant_builds_and_verifies(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "statement",
        "title": "投資家への結論",
        "variant": "evidence_strip",
        "statement": "ARR128億円、NRR114%、Q2進捗率48%を確認し、通期達成確度を継続評価",
        "recap_heading": "確認すべき3指標",
        "recap": [
            {"label": "ARR", "value": "128", "unit": "億円", "focal": True},
            {"label": "NRR", "value": "114", "unit": "%"},
            {"label": "進捗率", "value": "48", "unit": "%"}
        ],
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    r = run("verify_deck.py", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_player_missing_or_out_of_range_xy(tmp_path):
    deck = {"meta": {}, "slides": [{
        "pattern": "competitive_landscape",
        "title": "統合スイート×中堅企業の右上象限は当社のみが占める空白地帯である",
        "x_axis": {"low": "単機能", "high": "統合"},
        "y_axis": {"low": "小規模", "high": "大企業"},
        "players": [{"name": "当社", "y": 0.5}, {"name": "A社", "x": 1.5, "y": 0.2}],
        "source": "テスト",
    }]}
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1
    assert "がない/数値でない" in r.stdout and "範囲外" in r.stdout


def test_validate_allows_primary_insight_with_accent_series(tmp_path):
    # insight_style "primary" は Accent を消費しない — series の ECC85A 1 回は適法
    deck = _minimal_deck(insight="一言インサイト", insight_style="primary")
    deck["slides"][0]["chart"]["series"][0]["color"] = "ECC85A"
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0, r.stdout + r.stderr


def test_validate_rejects_too_many_takeaways(tmp_path):
    deck = _minimal_deck(
        takeaways=[{"heading": f"論点{i}", "body": "根拠"} for i in range(1, 6)])
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "takeaways" in r.stdout


def test_validate_warns_long_content_run(tmp_path):
    # The wall-of-content rest is a READ-deck concern: only a long, read-oriented deck
    # (>= LONG_DECK_MIN slides) is pushed toward a rest. A short talk deck is meant to flow
    # continuously, so a 6+ content run there is NOT nagged.
    content = _minimal_deck()["slides"][0]

    def mk(i):
        return dict(content, title=f"テスト市場は年率1{i % 9}%で拡大し、2030年に{i + 1}兆円に到達する見込み")

    # short talk deck: a 7-slide content run flows — no rest is pushed
    short = {"meta": {"title": "t"}, "slides": [mk(i) for i in range(7)]}
    spec = tmp_path / "short.json"
    spec.write_text(json.dumps(short, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "休符" not in r.stdout

    # long read deck (>= 18 slides): an unbroken 6+ content run warns
    long_run = {"meta": {"title": "t"}, "slides": [mk(i) for i in range(18)]}
    spec = tmp_path / "long.json"
    spec.write_text(json.dumps(long_run, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "休符" in r.stdout

    # dividers breaking the run below 6 clear the warning even in a long deck
    slides = []
    for i in range(18):
        if i and i % 4 == 0:
            slides.append({"pattern": "section_divider", "number": i // 4,
                           "title": f"章{i // 4}", "desc": "章の論点"})
        slides.append(mk(i))
    broken = {"meta": {"title": "t"}, "slides": slides}
    spec = tmp_path / "broken.json"
    spec.write_text(json.dumps(broken, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert "休符" not in r.stdout


def test_validate_flags_divider_overuse_in_short_deck(tmp_path):
    # A short talk deck chopped by dividers into 1-2 slide chapters is flagged: the deck-level
    # over-use warning fires, and any chapter holding <= 1 content slide is named as thin.
    content = _minimal_deck()["slides"][0]

    def mk(i):
        return dict(content, title=f"テスト市場は年率1{i % 9}%で拡大し、2030年に{i + 1}兆円に到達する見込み")

    slides = [
        {"pattern": "section_divider", "number": 1, "title": "章1", "desc": "章の論点"}, mk(0),
        {"pattern": "section_divider", "number": 2, "title": "章2", "desc": "章の論点"}, mk(1),
        {"pattern": "section_divider", "number": 3, "title": "章3", "desc": "章の論点"}, mk(2), mk(3),
    ]
    deck = {"meta": {"title": "t"}, "slides": slides}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0
    assert "枚数が多い" in r.stdout            # deck-level: reserve dividers for long/read decks
    assert "章が内容1枚以下" in r.stdout        # thin chapters (章1, 章2 hold one content slide)


def test_validate_allows_dividers_in_long_read_deck(tmp_path):
    # A long read deck (>= 18 slides) whose chapters each carry >= 3 proof slides is not nagged.
    content = _minimal_deck()["slides"][0]

    def mk(i):
        return dict(content, title=f"テスト市場は年率1{i % 9}%で拡大し、2030年に{i + 1}兆円に到達する見込み")

    slides = []
    for c in range(4):  # 4 chapters x 4 content slides = 20 slides total
        slides.append({"pattern": "section_divider", "number": c + 1,
                       "title": f"章{c + 1}", "desc": "章の論点"})
        slides += [mk(c * 4 + j) for j in range(4)]
    deck = {"meta": {"title": "t"}, "slides": slides}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0
    assert "枚数が多い" not in r.stdout and "章が内容1枚以下" not in r.stdout


def test_validate_warns_wide_comparison_table(tmp_path):
    deck = {"meta": {"title": "t"}, "slides": [{
        "pattern": "comparison_table",
        "title": "4社比較でも当社の機能カバレッジ優位は不変",
        "subtitle": "機能・価格・サポート体制の比較",
        "table": {"headers": ["評価軸", "当社", "A社", "B社", "C社"],
                  "rows": [["対応領域", "◎", "○", "○", "△"]]},
        "source": "テスト統計2026",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "4列以内" in r.stdout


def test_build_handles_null_focal_category_with_annotation(tmp_path):
    # LLM は「focal なし」を null で表現しがち — dict.get のデフォルトでは拾えない
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["focal_category"] = None
    deck["slides"][0]["chart"]["annotation"] = {"yoy": "+12%"}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr


def test_negative_waterfall_stays_inside_slide(tmp_path):
    from pptx.util import Inches

    deck = {"meta": {}, "slides": [{
        "pattern": "waterfall",
        "title": "一過性費用12億円の計上で通期は7億円の赤字に転落する見込みである",
        "unit": "億円",
        "items": [{"label": "期首", "value": 5, "kind": "start"},
                  {"label": "一過性費用", "value": -12},
                  {"label": "期末", "value": -7, "kind": "end"}],
        "source": "テスト",
    }]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    for slide in pptx.Presentation(out).slides:
        for sh in slide.shapes:
            if sh.top is not None and sh.height is not None:
                assert sh.top + sh.height <= Inches(7.51), \
                    f"shape extends past slide bottom: {sh.shape_type}"


def test_section_divider_without_number_has_no_void(tmp_path):
    from pptx.util import Inches

    deck = {"meta": {}, "slides": [{"pattern": "section_divider", "title": "市場環境"}]}
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    tops = [sh.top for slide in pptx.Presentation(out).slides for sh in slide.shapes
            if sh.has_text_frame and "市場環境" in sh.text_frame.text]
    # number を省いたら 96pt 数字分の空洞は入らない(タイトルが光学中心近くに来る)
    assert tops and min(tops) < Inches(3.5)


def test_lint_render_flags_slide_count_mismatch(tmp_path):
    from PIL import Image

    for i in (1, 2):
        Image.new("RGB", (733, 412), CANVAS_RGB).save(tmp_path / f"deck-{i}.png")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(_minimal_deck(), ensure_ascii=False))
    r = run("lint_render.py", tmp_path, "--spec", spec)
    assert r.returncode == 1 and "枚数不一致" in r.stdout


def test_validate_rejects_unsupported_chart_type(tmp_path):
    deck = _minimal_deck()
    deck["slides"][0]["chart"]["type"] = "area"
    bad = tmp_path / "deck.json"
    bad.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", bad)
    assert r.returncode == 1 and "未対応" in r.stdout and "Traceback" not in r.stderr


def test_chart_type_lists_stay_in_sync():
    sys.path.insert(0, str(SCRIPTS))
    try:
        import build_deck
        import validate_spec
        assert set(validate_spec.SUPPORTED_CHART_TYPES) == set(build_deck.CHART_TYPES)
    finally:
        sys.path.remove(str(SCRIPTS))


def test_validate_warns_missing_speaker_notes(tmp_path):
    deck = _minimal_deck()  # content slide without speaker_notes
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "speaker_notes" in r.stdout


def test_validate_talk_minutes_mismatch_warns(tmp_path):
    deck = _minimal_deck(speaker_notes="短い一言だけの語り。")
    deck["meta"]["talk_minutes"] = 10
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "talk_minutes=10" in r.stdout and "乖離" in r.stdout


def test_validate_talk_minutes_match_passes(tmp_path):
    deck = _minimal_deck(speaker_notes="あ" * 300)  # 300字 ≈ 1分
    deck["meta"]["talk_minutes"] = 1
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 0 and "乖離" not in r.stdout
    assert "talk script" in r.stdout and "target 1分" in r.stdout


def test_validate_rejects_nonnumeric_talk_minutes(tmp_path):
    deck = _minimal_deck(speaker_notes="語り。")
    deck["meta"]["talk_minutes"] = "abc"
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    r = run("validate_spec.py", spec)
    assert r.returncode == 1 and "talk_minutes" in r.stdout


def test_build_writes_speaker_notes_into_pptx(tmp_path):
    deck = _minimal_deck(speaker_notes="この市場は2年で倍になります。次のスライドで内訳を示します。")
    spec = tmp_path / "deck.json"
    spec.write_text(json.dumps(deck))
    out = tmp_path / "deck.pptx"
    r = run("build_deck.py", spec, "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    prs = pptx.Presentation(str(out))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "次のスライドで内訳を示します" in notes


# --- image-asset track + native table merge (skip when optional backends absent) ------------

def test_native_table_col0_spans_merges_and_builds(tmp_path):
    deck = {"slides": [{
        "pattern": "comparison_table", "title": "セグメント別は電機が牽引し+2,559",
        "subtitle": "セグメント別の売上・利益", "source": "Act分析",
        "table": {"headers": ["区分", "科目", "前期", "当期"],
                  "rows": [["電機", "売上", "29,265", "31,824"], ["", "利益", "1,585", "1,370"],
                           ["機械", "売上", "9,189", "9,469"], ["", "利益", "269", "417"]],
                  "align": ["l", "l", "r", "r"], "emphasis_col": 3, "color_negatives": True,
                  "col0_spans": [[0, 2], [2, 2]]}}]}
    spec = tmp_path / "d.json"; spec.write_text(json.dumps(deck))
    assert run("validate_spec.py", spec).returncode == 0
    out = tmp_path / "d.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0
    # the first-column group cells are merged (one origin spans two rows)
    tbl = pptx.Presentation(str(out)).slides[0].shapes[0]
    # locate the graphic-frame table
    from pptx.util import Emu  # noqa: F401
    tables = [sh.table for sh in pptx.Presentation(str(out)).slides[0].shapes if sh.has_table]
    assert tables and tables[0].cell(1, 0).is_merge_origin


def test_image_asset_track_combo_and_diagram(tmp_path):
    import shutil
    pytest.importorskip("matplotlib")
    if shutil.which("dot") is None:
        pytest.skip("Graphviz CLI (dot) not installed")
    deck = {"slides": [
        {"pattern": "chart_insight", "title": "売上と利益率を2軸で示す",
         "subtitle": "売上高と営業利益率の推移", "source": "Act分析",
         "chart": {"kind": "combo", "categories": ["'23", "'24", "'25"],
                   "bar": {"name": "売上", "values": [138, 162, 190], "unit": "億円"},
                   "line": {"name": "利益率", "values": [10.1, 12.2, 13.4], "unit": "%"}}},
        {"pattern": "diagram", "title": "資本関係を一枚で示す",
         "subtitle": "グループの資本構成", "source": "Act分析",
         "diagram": {"kind": "org_tree",
                     "nodes": [{"id": "h", "label": "持株", "focal": True}, {"id": "a", "label": "A社"}],
                     "edges": [{"from": "h", "to": "a", "label": "100%"}]}},
        {"pattern": "chart_insight", "title": "自社は品質と実績で優位",
         "subtitle": "3軸での競合比較", "source": "Act分析",
         "chart": {"kind": "radar", "axes": ["価格", "品質", "実績"],
                   "series": [{"name": "自社", "values": [4, 5, 5]}, {"name": "他社", "values": [3, 3, 3]}]}},
        {"pattern": "diagram", "title": "地域別の製品カバレッジ",
         "subtitle": "地域×製品の展開状況", "source": "Act分析",
         "diagram": {"kind": "matrix", "rows": ["東", "西"], "cols": ["X", "Y"],
                     "cells": {"0,0": True, "1,1": True}}}]}
    spec = tmp_path / "d.json"; spec.write_text(json.dumps(deck))
    assert run("validate_spec.py", spec).returncode == 0
    out = tmp_path / "d.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0
    # assets cached beside the pptx with a manifest and audit sidecars
    assets = tmp_path / "assets"
    assert (assets / "asset-manifest.json").exists()
    pngs = list(assets.glob("*.png"))
    assert len(pngs) == 4 and all(p.with_suffix(".json").exists() for p in pngs)
    # both slides embed a picture (the image asset)
    prs = pptx.Presentation(str(out))
    for sl in prs.slides:
        assert any(sh.shape_type == 13 for sh in sl.shapes)  # 13 = PICTURE


def test_act_theme_reads_one_token_core(tmp_path):
    sys.path.insert(0, str(SCRIPTS))
    import act_theme as t
    assert t.COLORS["primary"].startswith("#") and t.SERIES_PALETTE
    assert len(t.tokens_hash()) == 12


def test_native_chart_grid_small_multiples(tmp_path):
    deck = {"slides": [{
        "pattern": "chart_grid", "title": "3事業とも増収しCAGR11%成長",
        "subtitle": "事業別の売上・利益推移", "source": "Act分析",
        "charts": [
            {"title": "売上", "chart": {"type": "column", "unit": "百万円", "categories": ["24", "25", "26"],
             "series": [{"name": "売上", "values": [118, 125, 130]}], "value_labels": True,
             "focal_category": 2, "annotation": {"badge": "CAGR 5%"}}},
            {"title": "利益", "chart": {"type": "column", "unit": "百万円", "categories": ["24", "25", "26"],
             "series": [{"name": "利益", "values": [27, 29, 37]}], "value_labels": True, "focal_category": 2}}]}]}
    spec = tmp_path / "d.json"; spec.write_text(json.dumps(deck))
    assert run("validate_spec.py", spec).returncode == 0
    out = tmp_path / "d.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0
    # two NATIVE chart graphic-frames tiled on one slide (editable, not an image)
    charts = [sh for sh in pptx.Presentation(str(out)).slides[0].shapes if sh.has_chart]
    assert len(charts) == 2


def test_image_annotations_layer_renders(tmp_path):
    import shutil
    pytest.importorskip("matplotlib")
    if shutil.which("dot") is None:
        pytest.skip("dot not installed")
    deck = {"slides": [{
        "pattern": "chart_insight", "title": "増減を要因別に注記した2軸コンボ",
        "subtitle": "売上と利益率の要因注記", "source": "Act分析",
        "chart": {"kind": "combo", "categories": ["'24", "'25", "'26"],
                  "bar": {"name": "売上", "values": [138, 162, 190], "unit": "億円"},
                  "line": {"name": "利益率", "values": [10.1, 12.2, 13.4], "unit": "%"},
                  "annotations": [{"target": 2, "text": "新製品寄与", "dy": 40},
                                  {"target": 0, "text": "統合完了"}]}}]}
    spec = tmp_path / "d.json"; spec.write_text(json.dumps(deck))
    assert run("validate_spec.py", spec).returncode == 0
    out = tmp_path / "d.pptx"
    assert run("build_deck.py", spec, "-o", out).returncode == 0
    # annotations must not break the deterministic content-addressed cache
    import sys as _sys
    _sys.path.insert(0, str(SCRIPTS))
    import act_assets as A
    aid1 = A.asset_id(deck["slides"][0]["chart"], (8.6, 4.7))
    aid2 = A.asset_id(deck["slides"][0]["chart"], (8.6, 4.7))
    assert aid1 == aid2 and len(aid1) == 16


# ---- review follow-ups (PR #127): direct coverage for the robustness mechanisms ----

def _import_build_deck():
    sys.path.insert(0, str(SCRIPTS))
    import importlib
    import build_deck
    return importlib.reload(build_deck)


def test_statement_lines_packs_clauses_without_orphan_tail():
    bd = _import_build_deck()
    stmt = "中堅企業向けSaaSは今後3年が参入の最終ウィンドウ、経理領域からの段階参入とM&A活用で、5年でARR68億円の事業構築を提言"
    lines = bd._statement_lines(stmt, 8.0, 31)
    # 内容が保存され(節の欠落なし)、複数行に分割される
    assert "".join(lines) == stmt
    assert len(lines) >= 2
    # 各行は節境界(、)で終わるか最終行 — 語中改行の孤立行を作らない
    for ln in lines[:-1]:
        assert ln.endswith("、")
    # 尾行が1-2文字の孤立にならない
    assert bd._ja_len(lines[-1]) > 2


def test_statement_lines_single_clause_passthrough():
    bd = _import_build_deck()
    single = "読点を含まない一文のステートメントはそのまま"
    assert bd._statement_lines(single, 8.0, 31) == [single]
    multiline = "一行目\n二行目"
    assert bd._statement_lines(multiline, 8.0, 31) == ["一行目", "二行目"]


def test_validate_warns_on_long_statement_clause(tmp_path):
    # 節が22字超の中央ヒーロー文 → 警告(読点あり/なしで文言が変わる)
    with_comma = {"slides": [{"pattern": "statement", "title": "結論",
                              "statement": "この読点つき文はとても長い節をひとつだけ含んでおり折返し警告の対象になる、短い節"}]}
    no_comma = {"slides": [{"pattern": "statement", "title": "結論",
                            "statement": "読点を全く含まない非常に長い一文のステートメントは節折返しが効かない旨を警告される"}]}
    for deck, needle in ((with_comma, "節が長い"), (no_comma, "読点のない一文")):
        f = tmp_path / "d.json"
        f.write_text(json.dumps(deck, ensure_ascii=False))
        r = run("validate_spec.py", f)
        assert needle in r.stdout, r.stdout


def test_financial_highlights_overflow_hero_reaches_support(tmp_path):
    # 4グループ目の主指標は落とさず補助ストリップへ回る
    groups = [{"label": f"G{i}", "metrics": [{"label": f"主指標{i}", "value": f"{i}00", "unit": "億円", "hero": True}]}
              for i in range(1, 5)]
    deck = {"slides": [{"pattern": "financial_highlights", "title": "4グループの主要指標は補助帯まで含め全件表示",
                        "groups": groups}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "out.pptx"
    r = run("build_deck.py", f, "-o", out)
    assert r.returncode == 0, r.stderr
    from pptx import Presentation
    texts = []
    for shape in Presentation(str(out)).slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    joined = "\n".join(texts)
    assert "主指標4" in joined, joined
    # 溢れ分で support ストリップが新設されるケース: 見出しと帯が実寸で描画される
    assert "補助指標" in joined, joined
    # validate は groups > 3 を警告する
    rv = run("validate_spec.py", f)
    assert "groups" in rv.stdout and "3" in rv.stdout


def test_oversized_table_frame_stays_inside_slide(tmp_path):
    # 行数過多で縮小スケールが走っても、テーブルのフレーム自体がスライド外へ溢れない
    rows = [[f"項目{i}", "あ" * 38, "い" * 30] for i in range(16)]
    deck = {"slides": [{"pattern": "comparison_table", "title": "16行の過大テーブルでも枠はスライド内に収まることを確認",
                        "table": {"headers": ["項目", "内容", "備考"], "rows": rows}}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "out.pptx"
    r = run("build_deck.py", f, "-o", out)
    assert r.returncode == 0, r.stderr
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation(str(out))
    frames = [s for s in prs.slides[0].shapes if s.has_table]
    assert frames
    for fr in frames:
        assert fr.top + fr.height <= Emu(int(7.5 * 914400)) + Emu(1)


def test_chart_grid_image_cell_routes_to_asset(tmp_path):
    # chart_grid のセル単位 image-kind エスカレーションが picture として埋め込まれる
    pytest.importorskip("matplotlib")
    deck = {"slides": [{"pattern": "chart_grid",
                        "title": "2系列の小型チャート群で売上と利益率の両立を確認",
                        "charts": [
                            {"title": "売上", "chart": {"type": "column", "unit": "億円",
                                                        "categories": ["FY24", "FY25"],
                                                        "series": [{"name": "売上", "values": [10, 12]}]}},
                            {"title": "利益率", "chart": {"kind": "combo", "categories": ["FY24", "FY25"],
                                                          "bar": {"name": "売上", "values": [10, 12], "unit": "億円"},
                                                          "line": {"name": "利益率", "values": [10.0, 12.5], "unit": "%"}}},
                        ]}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "out.pptx"
    r = run("build_deck.py", f, "-o", out)
    assert r.returncode == 0, r.stderr
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes), [s.shape_type for s in shapes]
    assert any(s.has_chart for s in shapes if hasattr(s, "has_chart"))


def test_validate_warns_on_out_of_range_focal_step(tmp_path):
    deck = {"slides": [{"pattern": "process_flow", "title": "3ステップの導線で登録転換までの流れを確認する",
                        "focal_step": 9,
                        "steps": [{"label": f"Step {i}", "items": ["項目"]} for i in range(1, 4)]}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "focal_step" in r.stdout, r.stdout
    out = tmp_path / "out.pptx"
    rb = run("build_deck.py", f, "-o", out)
    assert rb.returncode == 0, rb.stderr


def test_validate_covers_chart_grid_cells(tmp_path):
    # chart_grid のセルにも本体 chart と同じ検査(黙殺アノテーション警告・型チェック)が届く
    deck = {"slides": [{"pattern": "chart_grid",
                        "title": "2セルの小型チャート群で検査カバレッジを確認する",
                        "charts": [
                            {"chart": {"kind": "radar", "categories": ["a", "b", "c"],
                                       "series": [{"name": "s", "values": [1, 2, 3]}],
                                       "annotations": [{"target": 0, "text": "落ちる注記"}]}},
                            {"chart": {"type": "no_such_type", "unit": "億円", "categories": ["a"],
                                       "series": [{"name": "s", "values": [1]}]}},
                        ]}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    out = r.stdout
    assert "cell 1" in out and "int target" in out, out
    assert "cell 2" in out and "no_such_type" in out, out


def test_validate_rejects_incomplete_or_unknown_image_kinds(tmp_path):
    # 必須フィールド欠落の image kind と、非対応 kind の render:"image" 強制は
    # build 前に validate がエラーにする(黙って通すと build が落ちる)
    deck = {"slides": [
        {"pattern": "chart_insight", "title": "レーダー1枚で能力バランスの偏りを確認する",
         "chart": {"kind": "radar", "categories": ["a"], "series": [{"name": "s", "values": [1]}]}},
        {"pattern": "chart_insight", "title": "強制image指定の未知kindは事前に弾かれることを確認",
         "chart": {"render": "image", "kind": "column", "unit": "億円", "categories": ["a"],
                   "series": [{"name": "s", "values": [1]}]}},
    ]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "必須フィールドがない: axes" in r.stdout, r.stdout
    assert "image バックエンド非対応" in r.stdout, r.stdout


def test_financial_summary_image_chart_routes_to_asset(tmp_path):
    pytest.importorskip("matplotlib")
    deck = {"slides": [{"pattern": "financial_summary",
                        "title": "売上60億円と利益率12%の同時達成をテーブルとcomboで確認",
                        "table": {"headers": ["項目", "FY24", "FY25"],
                                  "rows": [["売上", "10", "12"], ["利益", "1", "2"]]},
                        "chart": {"kind": "combo", "categories": ["FY24", "FY25"],
                                  "bar": {"name": "売上", "values": [10, 12], "unit": "億円"},
                                  "line": {"name": "利益率", "values": [10.0, 16.7], "unit": "%"}}}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "out.pptx"
    r = run("build_deck.py", f, "-o", out)
    assert r.returncode == 0, r.stderr
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes)


def test_non_numeric_focal_step_does_not_crash_build(tmp_path):
    deck = {"slides": [{"pattern": "process_flow", "title": "非数値focal_stepでもビルドが完走することを確認",
                        "focal_step": "last",
                        "steps": [{"label": f"Step {i}", "items": ["項目"]} for i in range(1, 4)]}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "focal_step" in r.stdout
    out = tmp_path / "out.pptx"
    rb = run("build_deck.py", f, "-o", out)
    assert rb.returncode == 0, rb.stderr


def test_validate_covers_diagram_specs_and_grid_size(tmp_path):
    # diagram の asset spec も image-kind 検査に通り、chart_grid の5枚超はエラーになる
    deck = {"slides": [
        {"pattern": "diagram", "title": "組織ツリーで意思決定ラインの重複を確認する",
         "diagram": {"kind": "org_tree", "nodes": [{"id": "a", "label": "A"}]}},
        {"pattern": "diagram", "title": "未知kindのダイアグラムは事前に弾かれることを確認",
         "diagram": {"kind": "banana"}},
        {"pattern": "chart_grid", "title": "5セルのグリッドは契約違反としてエラーになることを確認",
         "charts": [{"chart": {"type": "column", "unit": "億円", "categories": ["a"],
                               "series": [{"name": "s", "values": [1]}]}} for _ in range(5)]},
    ]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "必須フィールドがない: edges" in r.stdout, r.stdout
    assert "diagram kind 'banana' は未対応" in r.stdout, r.stdout
    assert "パターン契約は2-4" in r.stdout, r.stdout


def test_validate_talk_script_fidelity_and_register(tmp_path):
    # (1) スライドに無い単位つき数値 (2) 敬体でないメモ書き (3) タイトル逐語読みの冒頭 を検出する
    title = "導入社数120社とNRR112%で初年度の検証を完了する"
    base_kpis = [{"label": "導入社数", "value": "120", "unit": "社"},
                 {"label": "NRR", "value": "112", "unit": "%"}]
    slides = [
        {"pattern": "kpi_dashboard", "title": title, "kpis": base_kpis,
         "speaker_notes": "初年度はこの2指標に絞ります。導入社数は120社、NRRは112%です。ここが立てば拡張投資に進めます。ARRは999億円に達する見込みです。次のスライドで前提を確認します。"},
        {"pattern": "kpi_dashboard", "title": title, "kpis": base_kpis,
         "speaker_notes": "導入社数120社・NRR112%の2指標。初年度はユニットエコノミクス検証。達成なら拡張投資へ。前提条件は次スライドの通り。以上の構成で全体像を先に共有する形とした。"},
        {"pattern": "kpi_dashboard", "title": title, "kpis": base_kpis,
         "speaker_notes": title + "、という結論です。導入社数は120社、NRRは112%を目標にします。この2つが立てば来期は拡張投資に進めますので、次のスライドで前提条件を確認します。"},
    ]
    f = tmp_path / "deck.json"
    f.write_text(json.dumps({"slides": slides}, ensure_ascii=False))
    r = run("validate_spec.py", f)
    out = r.stdout
    assert "スライド上に無い数値: 999億" in out, out
    assert "話し言葉でない" in out, out
    assert "タイトルの逐語読み" in out, out


def test_validate_fidelity_gate_matches_unit_pairs(tmp_path):
    # 数字だけの一致では通らない: スライドが 120社/112% のとき「112社」は単位違いとして検出。
    # 逆に、数値型のチャート値(12)+chart.unit(億円)を話す「12億円」は構造ペアで通る
    slides = [
        {"pattern": "kpi_dashboard", "title": "導入社数120社とNRR112%で初年度の検証を完了する",
         "kpis": [{"label": "導入社数", "value": "120", "unit": "社"},
                  {"label": "NRR", "value": "112", "unit": "%"}],
         "speaker_notes": "初年度の指標です。導入社数は112社を目標にします。この水準で検証が成立しますので、次のスライドで前提を確認します。八十字を超えるようにもう一文だけ補足します。"},
        {"pattern": "chart_insight", "title": "売上は2年で12億円へ拡大し成長率を維持する",
         "chart": {"type": "column", "unit": "億円", "categories": ["FY24", "FY25"],
                   "series": [{"name": "売上", "values": [10, 12]}]},
         "speaker_notes": "売上の推移をご覧ください。FY25には12億円へ届き、前年の10億円からの伸びを維持できるかが焦点になります。単価と件数のどちらが効いているかは次のスライドのドライバー分解で確認しますので、まずは水準感だけ押さえてください。"},
    ]
    f = tmp_path / "deck.json"
    f.write_text(json.dumps({"slides": slides}, ensure_ascii=False))
    r = run("validate_spec.py", f)
    out = r.stdout
    assert "112社" in out and "スライド上に無い数値" in out, out
    # slide 2: 数値型チャート値(10, 12)+chart.unit(億円)の組は構造ペアで通り、警告が出ない
    assert not any("slide 2" in ln and "スライド上に無い数値" in ln for ln in out.splitlines()), out


def test_validate_image_chart_data_shape(tmp_path):
    # combo の values 長さ不一致は build 前にエラーになる
    deck = {"slides": [{"pattern": "chart_insight", "title": "売上と利益率の両立を1枚で確認する",
                        "chart": {"kind": "combo", "categories": ["FY24", "FY25", "FY26"],
                                  "bar": {"name": "売上", "values": [10, 12], "unit": "億円"},
                                  "line": {"name": "利益率", "values": [10.0, 11.0, 12.0], "unit": "%"}}}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "bar.values が 2 件で categories 3 件と不一致" in r.stdout, r.stdout


def test_validate_fidelity_ignores_control_fields_and_scopes_table_units(tmp_path):
    slides = [
        # focal_category:1 は描画制御であり「1億円」の根拠にならない
        {"pattern": "chart_insight", "title": "売上は2年で12億円へ拡大し成長軌道を維持する",
         "chart": {"type": "column", "unit": "億円", "categories": ["FY24", "FY25"],
                   "series": [{"name": "売上", "values": [10, 12]}], "focal_category": 1},
         "speaker_notes": "売上の推移をご覧ください。まず1億円の水準から立ち上がり、FY25には12億円へ届く計画です。ここからの伸びの持続性を次のスライドで確認していきます。"},
        # 複数単位ヘッダーの表では単位は自列に限定される: 10 は億円列の値で、10% は存在しない
        {"pattern": "comparison_table", "title": "売上10億円と成長率5%の組み合わせを1表で確認する",
         "table": {"headers": ["項目", "売上(億円)", "YoY(%)"],
                   "rows": [["A事業", "10", "+5"]]},
         "speaker_notes": "この表のポイントはA事業です。売上は10億円、成長率はプラス5%で、規模と成長の両立ができています。仮に10%と読み間違えると評価を誤りますので、列の対応にご注意ください。"},
    ]
    f = tmp_path / "deck.json"
    f.write_text(json.dumps({"slides": slides}, ensure_ascii=False))
    r = run("validate_spec.py", f)
    out = r.stdout
    assert any("slide 1" in ln and "1億円" in ln for ln in out.splitlines()), out
    assert any("slide 2" in ln and "10%" in ln for ln in out.splitlines()), out
    # 正当な組(12億円 / 10億円 / 5%)は警告されない
    assert "12億円" not in out and "10億円" not in out and " 5%" not in out, out


def test_validate_rejects_zero_only_funnel(tmp_path):
    deck = {"slides": [{"pattern": "chart_insight", "title": "獲得ファネルの歩留まりを3段で確認する",
                        "chart": {"kind": "funnel",
                                  "stages": [{"label": "訪問", "value": 0},
                                             {"label": "登録", "value": 0}]}}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert "全て0以下" in r.stdout, r.stdout


def test_fidelity_corpus_has_delimiters_and_filters_control_keys(tmp_path):
    # キー順で focal_category の「1」と unit の「億円」が隣接しても偽陰性にならない
    deck = {"slides": [{"pattern": "chart_insight", "title": "売上は2年で12億円へ拡大し成長軌道を維持する",
                        "chart": {"focal_category": 1, "unit": "億円", "type": "column",
                                  "categories": ["FY24", "FY25"],
                                  "series": [{"name": "売上", "values": [10, 12]}]},
                        "speaker_notes": "売上の推移をご覧ください。1億円の水準からの立ち上がりを経て、FY25には12億円へ届く計画です。伸びの持続性は次のスライドで確認していきます。"}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert any("1億円" in ln and "スライド上に無い数値" in ln for ln in r.stdout.splitlines()), r.stdout


def test_pyramid_accepts_scalar_tiers(tmp_path):
    pytest.importorskip("matplotlib")
    deck = {"slides": [{"pattern": "diagram", "title": "3層の実行体制で意思決定から現場までを接続する",
                        "diagram": {"kind": "pyramid", "tiers": ["経営", "推進", "現場"]}}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    out = tmp_path / "out.pptx"
    r = run("build_deck.py", f, "-o", out)
    assert r.returncode == 0, r.stderr


def test_fidelity_flags_sign_flip_against_negative_slide_values(tmp_path):
    # スライドでは △5.0(損失)のみの数を、正方向の文脈で「5億円」と語ると警告。
    # 「赤字5億円」のように向きの語があれば通る
    table = {"headers": ["(億円)", "FY26"], "rows": [["営業利益", "△5.0"]]}
    base = {"pattern": "comparison_table", "title": "営業損失5.0億円までの縮小を1表で確認する", "table": table}
    slides = [
        dict(base, speaker_notes="今期の損益です。営業利益は5億円に到達し、収益化の目処が立ちました。この水準を維持できるかを次のスライドで確認します。八十字対策の補足文です。"),
        dict(base, speaker_notes="今期の損益です。営業赤字5億円まで縮小し、黒字化が視野に入りました。改善ドライバーの中身は次のスライドで確認します。八十字対策の補足文です。"),
    ]
    f = tmp_path / "deck.json"
    f.write_text(json.dumps({"slides": slides}, ensure_ascii=False))
    r = run("validate_spec.py", f)
    out = r.stdout
    assert any("slide 1" in ln and "負値" in ln for ln in out.splitlines()), out
    assert not any("slide 2" in ln and "負値" in ln for ln in out.splitlines()), out


def test_sign_flip_not_warned_when_positive_form_also_on_slide(tmp_path):
    # 同じ数がスライド上に正・負の両形で存在するなら、符号なしのナレーションは正当
    table = {"headers": ["(億円)", "FY25", "FY26"], "rows": [["売上", "5.0", "6.0"], ["営業利益", "△5.0", "1.0"]]}
    deck = {"slides": [{"pattern": "comparison_table",
                        "title": "売上5.0億円と損失解消の同時進行を1表で確認する", "table": table,
                        "speaker_notes": "この表の見方です。売上は5億円で前期並みを確保しつつ、営業利益は赤字5億円から黒字1億円へ転換しました。両者が同じ5という数字なので読み違えにご注意ください。"}]}
    f = tmp_path / "deck.json"
    f.write_text(json.dumps(deck, ensure_ascii=False))
    r = run("validate_spec.py", f)
    assert not any("負値" in ln for ln in r.stdout.splitlines()), r.stdout


def test_venn_subsets_derive_from_sizes_and_overlaps():
    pytest.importorskip("matplotlib")
    sys.path.insert(0, str(SCRIPTS))
    import importlib
    import act_assets
    act_assets = importlib.reload(act_assets)
    # 2集合: サイズと重なりから排他リージョンを導出
    sub2 = act_assets._venn_subsets([{"label": "A", "size": 20}, {"label": "B", "size": 10}], {"AB": 4})
    assert sub2 == (16.0, 6.0, 4.0)
    # 3集合: 包除原理で導出され、負にならない
    sub3 = act_assets._venn_subsets(
        [{"label": "A", "size": 20}, {"label": "B", "size": 15}, {"label": "C", "size": 10}],
        {"AB": 5, "AC": 4, "BC": 3, "ABC": 2})
    assert sub3 == (13.0, 9.0, 3.0, 5.0, 2.0, 1.0, 2.0)
    assert all(v >= 0 for v in sub3)
    # サイズ違いは異なるリージョン値になる(固定既定値で潰れない)
    other = act_assets._venn_subsets([{"label": "A", "size": 40}, {"label": "B", "size": 10}], {"AB": 4})
    assert other != sub2
